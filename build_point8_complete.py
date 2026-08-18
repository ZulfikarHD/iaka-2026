import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sys
import os

# --- PALETTE CONSTANTS ---
NAVY = RGBColor(32, 40, 143)         # #20288F - Primary Brand
PURPLE = RGBColor(106, 44, 201)      # #6A2CC9 - Accent
DARK_NAVY = RGBColor(26, 35, 126)    # #1A237E
GREEN = RGBColor(46, 125, 50)        # #2E7D32 - Success
DARK_GREEN = RGBColor(27, 94, 32)    # #1B5E20
RED = RGBColor(198, 40, 40)          # #C62828 - Alert
ORANGE = RGBColor(230, 81, 0)        # #E65100 - Warning
DARK_TEXT = RGBColor(42, 42, 61)     # #2A2A3D - Body Text
MUTED_TEXT = RGBColor(110, 110, 130) # #6E6E82 - Muted Text
WHITE = RGBColor(255, 255, 255)

# Container Fills
FILL_LIGHT_NAVY = RGBColor(244, 247, 254)   # #F4F7FE
FILL_LIGHT_PURPLE = RGBColor(248, 246, 254) # #F8F6FE
FILL_LIGHT_GREEN = RGBColor(232, 245, 233)  # #E8F5E9
FILL_LIGHT_AMBER = RGBColor(255, 249, 230)  # #FFF9E6
FILL_LIGHT_RED = RGBColor(255, 245, 245)    # #FFF5F5

# Borders
BORDER_NAVY = RGBColor(210, 220, 245)       # #D2DCF5
BORDER_PURPLE = RGBColor(217, 206, 247)     # #D9CEF7
BORDER_GREEN = RGBColor(165, 214, 167)      # #A5D6A7
BORDER_AMBER = RGBColor(255, 224, 130)      # #FFE082
BORDER_RED = RGBColor(255, 205, 210)        # #FFCDD2

def add_header(slide, badge_num, title_text, purpose_text):
    # Wave ornament
    if os.path.exists('wave_ornament.png'):
        slide.shapes.add_picture('wave_ornament.png', Inches(-0.118), Inches(-0.196), Inches(1.584), Inches(0.891))
    # Peruri Logo
    if os.path.exists('peruri_logo.png'):
        slide.shapes.add_picture('peruri_logo.png', Inches(12.09), Inches(0.013), Inches(1.247), Inches(0.694))

    # Badge Box
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(0.65), Inches(0.80), Inches(0.72))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PURPLE
    badge.line.color.rgb = PURPLE
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = badge_num
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    tb_title = slide.shapes.add_textbox(Inches(1.55), Inches(0.60), Inches(10.20), Inches(0.80))
    tf = tb_title.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(23)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Purpose Subtitle
    tb_purpose = slide.shapes.add_textbox(Inches(0.60), Inches(1.48), Inches(11.80), Inches(0.38))
    tf = tb_purpose.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Tujuan bagian ini:  "
    r1.font.name = 'Arial'
    r1.font.size = Pt(11.5)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = purpose_text
    r2.font.name = 'Arial'
    r2.font.size = Pt(11.5)
    r2.font.color.rgb = DARK_TEXT

    # Footer
    tb_foot = slide.shapes.add_textbox(Inches(0.60), Inches(7.02), Inches(12.13), Inches(0.24))
    tf = tb_foot.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = "IAKA 2026 — Kerangka Presentasi Peserta  ·  DSS SIRINE 4.0 Unit Cetak Pita Cukai"
    p.font.name = 'Arial'
    p.font.size = Pt(8.5)
    p.font.color.rgb = MUTED_TEXT

def add_kpi_card(slide, left, top, width, height, title, value, subtext, color, fill_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.04)
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)

    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = 'Arial'
    p0.font.size = Pt(8.5)
    p0.font.bold = True
    p0.font.color.rgb = color
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = value
    p1.font.name = 'Arial'
    p1.font.size = Pt(15.5)
    p1.font.bold = True
    p1.font.color.rgb = color
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = subtext
    p2.font.name = 'Arial'
    p2.font.size = Pt(7.2)
    p2.font.color.rgb = DARK_TEXT
    p2.alignment = PP_ALIGN.CENTER

# -------------------------------------------------------------
# BUILD SLIDE 08.1
# -------------------------------------------------------------
def build_slide_8_1(slide):
    # Clear shapes on slide
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    add_header(
        slide,
        badge_num="08.1",
        title_text="Lesson Learned: Tantangan Utama Lapangan, Mitigasi Masalah & Pembelajaran Proyek",
        purpose_text="Bagikan pembelajaran agar inovasi/kaizen dapat dipetik intisarinya dan direplikasi oleh unit produksi lain di Perum Peruri."
    )

    # Top 4 Hero KPI Cards
    add_kpi_card(slide, 0.60, 1.90, 2.86, 1.08, "ADAPTASI OPERATOR", "100% Adopsi Digital", "42 Operator 3 Gilir Beralih Penuh dari Folio", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 3.69, 1.90, 2.86, 1.08, "EFISIENSI INPUT FORM", "< 30 Detik / PO", "Autofill SAP Pangkas Resistensi Administrasi", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 6.78, 1.90, 2.86, 1.08, "RESOLUSI MASALAH", "100% Kendala Tuntas", "Silo Data, Shift Malam & Koneksi Teratasi", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 9.87, 1.90, 2.86, 1.08, "KULTUR DATA BARU", "Data-Driven Culture", "Transparansi Granular Lenyapkan Saling Tuduh", DARK_GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)

    # Left Panel: 1. TANTANGAN UTAMA LAPANGAN & STRATEGI MITIGASI RIIL
    header_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.60), Inches(3.08), Inches(5.95), Inches(0.34))
    header_left.fill.solid()
    header_left.fill.fore_color.rgb = NAVY
    header_left.line.fill.background()
    tf = header_left.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "1. TANTANGAN UTAMA LAPANGAN & STRATEGI MITIGASI RIIL"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(3.42), Inches(5.95), Inches(3.02))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = FILL_LIGHT_NAVY
    box_left.line.color.rgb = BORDER_NAVY
    box_left.line.width = Pt(1.0)

    # 3 Cards on Left Panel
    left_cards = [
        (
            Inches(3.48), Inches(0.88), RED, BORDER_RED,
            "1. TANTANGAN KULTURAL & BEBAN INPUT (HUMAN RESISTANCE)",
            "• Kendala Lapangan: ", "Operator khawatir form digital menambah beban administratif saat awasi mesin cetak.",
            "• Mitigasi Lean UX: ", "Desain Autofill SAP cerdas (<30 dtk) & shortcut Ctrl+S melenyapkan resistensi (100% patuh)."
        ),
        (
            Inches(4.42), Inches(0.92), ORANGE, BORDER_AMBER,
            "2. TANTANGAN OPERASIONAL 3 GILIR (CIRCADIAN FATIGUE)",
            "• Kendala Lapangan: ", "Kelelahan visual gilir malam (23.00–07.00 WIB) memicu risiko keterlambatan entri PO.",
            "• Mitigasi Handover: ", "Checklist verifikasi input PO diwajibkan jadi syarat mutlak tanda tangan serah terima gilir."
        ),
        (
            Inches(5.40), Inches(0.86), NAVY, BORDER_NAVY,
            "3. TANTANGAN TEKNIS DATA SILO & KOMITMEN ZERO CAPEX",
            "• Kendala Lapangan: ", "SAP ERP ZPPRSIPPC0012, verifikasi cacat, dan buku folio terisolasi tanpa anggaran software.",
            "• Mitigasi In-House: ", "Pengembangan web intranet mandiri (Laravel Service, Inertia, Vue) dengan CAPEX & OPEX Rp 0."
        )
    ]

    for top_pos, h_size, title_col, border_col, title_txt, l1_lbl, l1_val, l2_lbl, l2_val in left_cards:
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), top_pos, Inches(5.71), h_size)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = border_col
        c.line.width = Pt(1.0)
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.02)
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = title_txt
        p.font.name = 'Arial'
        p.font.size = Pt(7.3)
        p.font.bold = True
        p.font.color.rgb = title_col

        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = l1_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p1.add_run()
        r2.text = l1_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = l2_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p2.add_run()
        r2.text = l2_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

    # Label under left panel
    tb_src_l1 = slide.shapes.add_textbox(Inches(0.72), Inches(6.28), Inches(5.71), Inches(0.14))
    tf_src_l1 = tb_src_l1.text_frame
    tf_src_l1.word_wrap = False
    tf_src_l1.margin_top = Inches(0)
    tf_src_l1.margin_left = Inches(0)
    p = tf_src_l1.paragraphs[0]
    p.text = "*(Sumber: Laporan Monitoring & Evaluasi Kaizen Unit Cetak Pita Cukai Semester 1 2026)"
    p.font.name = 'Arial'
    p.font.size = Pt(6.3)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY

    # Right Panel: 2. PEMBELAJARAN KUNCI (KEY LESSONS) & PRINSIP SUKSES KAIZEN
    header_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.70), Inches(3.08), Inches(6.03), Inches(0.34))
    header_right.fill.solid()
    header_right.fill.fore_color.rgb = PURPLE
    header_right.line.fill.background()
    tf = header_right.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "2. PEMBELAJARAN KUNCI (KEY LESSONS) & PRINSIP SUKSES KAIZEN"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.70), Inches(3.42), Inches(6.03), Inches(3.02))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = FILL_LIGHT_PURPLE
    box_right.line.color.rgb = BORDER_PURPLE
    box_right.line.width = Pt(1.0)

    # 3 Cards on Right Panel
    right_cards = [
        (
            Inches(3.48), Inches(0.88), PURPLE, BORDER_PURPLE,
            "1. KEMUDAHAN PENGGUNA (LEAN UX) PENENTU UTAMA ADOPSI LAPANGAN",
            "• Intisari Kaizen: ", "Sistem secanggih apa pun akan ditolak di lini cetak jika menyita waktu operator mesin.",
            "• Prinsip Sukses: ", "Desain cepat (<30 dtk) & otomatis adalah syarat mutlak terciptanya kepatuhan digital 100%."
        ),
        (
            Inches(4.42), Inches(0.92), DARK_GREEN, BORDER_GREEN,
            "2. TRANSPARANSI DATA GRANULAR MENGUBAH KULTUR MENJADI KOLABORATIF",
            "• Intisari Kaizen: ", "Ketiadaan data granular memicu kebiasaan saling menduga kesalahan saat inschiet naik.",
            "• Prinsip Sukses: ", "Data mesin & shift objektif membuktikan anomali mekanis, mengarahkan tim ke coaching terarah."
        ),
        (
            Inches(5.40), Inches(0.86), NAVY, BORDER_NAVY,
            "3. PELEMBAGAAN LEGAL-FORMAL MENGUNCI KEBERLANJUTAN SISTEM",
            "• Intisari Kaizen: ", "Inovasi tanpa regulasi baku berisiko kembali ke cara manual lama saat rotasi personel.",
            "• Prinsip Sukses: ", "IK-PPC-2026-001, SOP-PPC-2026-004 & BA-PPC-2026-002 kunci kepatuhan permanen ISO 9001."
        )
    ]

    for top_pos, h_size, title_col, border_col, title_txt, r1_lbl, r1_val, r2_lbl, r2_val in right_cards:
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.82), top_pos, Inches(5.79), h_size)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = border_col
        c.line.width = Pt(1.0)
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.02)
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = title_txt
        p.font.name = 'Arial'
        p.font.size = Pt(7.3)
        p.font.bold = True
        p.font.color.rgb = title_col

        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = r1_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p1.add_run()
        r2.text = r1_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = r2_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p2.add_run()
        r2.text = r2_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

    # Label under right panel
    tb_src_r1 = slide.shapes.add_textbox(Inches(6.82), Inches(6.28), Inches(5.79), Inches(0.14))
    tf_src_r1 = tb_src_r1.text_frame
    tf_src_r1.word_wrap = False
    tf_src_r1.margin_top = Inches(0)
    tf_src_r1.margin_left = Inches(0)
    p = tf_src_r1.paragraphs[0]
    p.text = "*(Sumber: Sintesis Evaluasi Pelaksanaan Proyek Kaizen & Budaya Mutu Perum Peruri 2026)"
    p.font.name = 'Arial'
    p.font.size = Pt(6.3)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Bottom Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(6.52), Inches(12.13), Inches(0.44))
    banner.fill.solid()
    banner.fill.fore_color.rgb = FILL_LIGHT_PURPLE
    banner.line.color.rgb = BORDER_PURPLE
    banner.line.width = Pt(1.0)
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "• Intisari Transformasi Lapangan: "
    r1.font.name = 'Arial'
    r1.font.size = Pt(7.5)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = "Keberhasilan DSS SIRINE 4.0 membuktikan bahwa digitalisasi di lini cetak sekuriti harus memadukan kesederhanaan antarmuka (Lean UX), pembinaan budaya kerja kolaboratif 3 gilir 24/7, dan penguncian regulasi legal-formal agar menghasilkan perbaikan operasional berkelanjutan."
    r2.font.name = 'Arial'
    r2.font.size = Pt(7.4)
    r2.font.color.rgb = DARK_TEXT

# -------------------------------------------------------------
# BUILD SLIDE 08.2
# -------------------------------------------------------------
def build_slide_8_2(slide):
    # Clear shapes on slide
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    add_header(
        slide,
        badge_num="08.2",
        title_text="Lesson Learned: Rekomendasi Pengembangan Lanjutan & Kesimpulan Akhir Proyek",
        purpose_text="Sajikan peta jalan pengembangan inovasi lintas unit dan tegaskan kesimpulan ringkas atas pencapaian Kaizen DSS SIRINE 4.0."
    )

    # Top 4 Hero KPI Cards
    add_kpi_card(slide, 0.60, 1.90, 2.86, 1.08, "PENURUNAN INSCHIET", "-1,28 pp (-27,8%)", "Turun dari 4,61% ke 3,33% di Q2 2026", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 3.69, 1.90, 2.86, 1.08, "COST AVOIDANCE RIIL", "Rp 2,23 Miliar", "Realisasi Nyata S1 2026 (Proyeksi Rp 6,82 M/Thn)", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 6.78, 1.90, 2.86, 1.08, "NET VALUE & ROI", "Rp 6,82 M / Instant", "CAPEX Rp 0 & OPEX Rp 0 (Payback 0 Bulan)", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 9.87, 1.90, 2.86, 1.08, "POTENSI REPLIKASI", "4 Lini Produksi", "Meterai, Paspor RI, Khazanah & Uang Kertas", DARK_GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)

    # Left Panel: 1. REKOMENDASI PENGEMBANGAN LANJUTAN & RENCANA STRATEGIS
    header_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.60), Inches(3.08), Inches(5.95), Inches(0.34))
    header_left.fill.solid()
    header_left.fill.fore_color.rgb = NAVY
    header_left.line.fill.background()
    tf = header_left.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "1. REKOMENDASI PENGEMBANGAN LANJUTAN & RENCANA STRATEGIS"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(3.42), Inches(5.95), Inches(3.02))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = FILL_LIGHT_NAVY
    box_left.line.color.rgb = BORDER_NAVY
    box_left.line.width = Pt(1.0)

    # 3 Cards on Left Panel
    left_cards2 = [
        (
            Inches(3.48), Inches(0.88), PURPLE, BORDER_PURPLE,
            "1. PENGEMBANGAN INTERNAL UNIT CETAK (Q3–Q4 2026)",
            "• Rapor Scoring Kinerja Tim: ", "Otomasi evaluasi produktivitas komposit (Grade A–E) untuk apresiasi regu kerja cetak.",
            "• Maintenance Log Book: ", "Integrasi langsung Pareto cacat ke catatan servis teknisi Divisi Pemeliharaan."
        ),
        (
            Inches(4.42), Inches(0.92), NAVY, BORDER_NAVY,
            "2. REPLIKASI KORPORAT KE 4 LINI PRODUKSI STRATEGIS (2027)",
            "• 4 Unit Sekuriti Utama: ", "Replikasi core engine ke lini Meterai, Paspor RI, Khazanah/Finishing, dan Uang Kertas.",
            "• Efisiensi Korporat: ", "Memanfaatkan template arsitektur teruji tanpa perlu biaya riset software dari awal."
        ),
        (
            Inches(5.40), Inches(0.86), DARK_GREEN, BORDER_GREEN,
            "3. AKSELERASI INDUSTRI 4.0 & TRANSFORMASI DIGITAL PERURI",
            "• Integrasi PPIC Korporat: ", "Sinkronisasi histori stabilitas mesin ke perencanaan alokasi pesanan kerja dinamis.",
            "• Asesmen INDI 4.0: ", "Mendukung pencapaian skor INDI 4.0 Kemenperin RI pilar Smart Factory & Data-Driven."
        )
    ]

    for top_pos, h_size, title_col, border_col, title_txt, l1_lbl, l1_val, l2_lbl, l2_val in left_cards2:
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), top_pos, Inches(5.71), h_size)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = border_col
        c.line.width = Pt(1.0)
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.02)
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = title_txt
        p.font.name = 'Arial'
        p.font.size = Pt(7.3)
        p.font.bold = True
        p.font.color.rgb = title_col

        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = l1_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p1.add_run()
        r2.text = l1_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = l2_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p2.add_run()
        r2.text = l2_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

    # Label under left panel
    tb_src_l2 = slide.shapes.add_textbox(Inches(0.72), Inches(6.28), Inches(5.71), Inches(0.14))
    tf_src_l2 = tb_src_l2.text_frame
    tf_src_l2.word_wrap = False
    tf_src_l2.margin_top = Inches(0)
    tf_src_l2.margin_left = Inches(0)
    p = tf_src_l2.paragraphs[0]
    p.text = "*(Sumber: Rencana Pengembangan Kaizen Unit Cetak Pita Cukai & RJPP Transformasi Digital Peruri 2026)"
    p.font.name = 'Arial'
    p.font.size = Pt(6.3)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY

    # Right Panel: 2. KESIMPULAN RINGKAS ATAS KESELURUHAN PROYEK DSS SIRINE 4.0
    header_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.70), Inches(3.08), Inches(6.03), Inches(0.34))
    header_right.fill.solid()
    header_right.fill.fore_color.rgb = GREEN
    header_right.line.fill.background()
    tf = header_right.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "2. KESIMPULAN RINGKAS ATAS KESELURUHAN PROYEK DSS SIRINE 4.0"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.70), Inches(3.42), Inches(6.03), Inches(3.02))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = FILL_LIGHT_GREEN
    box_right.line.color.rgb = BORDER_GREEN
    box_right.line.width = Pt(1.0)

    # 3 Cards on Right Panel
    right_cards2 = [
        (
            Inches(3.48), Inches(0.88), NAVY, BORDER_NAVY,
            "1. ELIMINASI DATA SILO & PENYELESAIAN AKAR MASALAH SECARA PRESISI",
            "• Silsilah Terlacak 100%: ", "Satukan 3 pulau data: PO ➔ Mesin ➔ Shift ➔ Tim ➔ Cacat (Klausul ISO 9001:2015).",
            "• Aksi Terarah: ", "Membedakan anomali mekanis vs variasi metode kerja, memangkas downtime >50%–75%."
        ),
        (
            Inches(4.42), Inches(0.92), GREEN, BORDER_GREEN,
            "2. PENCIPTAAN NILAI FINANSIAL SIGNIFIKAN DENGAN INSTANT ROI",
            "• Efisiensi Riil Terbukti: ", "Penghematan Rp 2,23 Miliar di S1 2026 & proyeksi tahunan Rp 6,82 Miliar (-1,28 pp).",
            "• Mandiri Tanpa Lisensi: ", "100% In-house development (CAPEX & OPEX Rp 0) hasilkan Payback Period 0 Bulan."
        ),
        (
            Inches(5.40), Inches(0.86), PURPLE, BORDER_PURPLE,
            "3. KEDAULATAN SISTEM, DISIPLIN TIM & JAMINAN MUTU NEGARA",
            "• Keberlanjutan Mandiri: ", "42 operator terlatih (nilai 94,8/100) dikunci legal lewat IK-PPC & SOP-PPC resmi.",
            "• Kepuasan Kemenkeu RI: ", "Mengamankan SLA pengiriman pita cukai 100% On-Time tanpa komplain mutu resmi."
        )
    ]

    for top_pos, h_size, title_col, border_col, title_txt, r1_lbl, r1_val, r2_lbl, r2_val in right_cards2:
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.82), top_pos, Inches(5.79), h_size)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = border_col
        c.line.width = Pt(1.0)
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.02)
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = title_txt
        p.font.name = 'Arial'
        p.font.size = Pt(7.3)
        p.font.bold = True
        p.font.color.rgb = title_col

        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = r1_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p1.add_run()
        r2.text = r1_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = r2_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p2.add_run()
        r2.text = r2_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

    # Label under right panel
    tb_src_r2 = slide.shapes.add_textbox(Inches(6.82), Inches(6.28), Inches(5.79), Inches(0.14))
    tf_src_r2 = tb_src_r2.text_frame
    tf_src_r2.word_wrap = False
    tf_src_r2.margin_top = Inches(0)
    tf_src_r2.margin_left = Inches(0)
    p = tf_src_r2.paragraphs[0]
    p.text = "*(Sumber: Resume Evaluasi Komprehensif Tim Inovasi Kaizen & Manajemen SBU HSS Perum Peruri 2026)"
    p.font.name = 'Arial'
    p.font.size = Pt(6.3)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN

    # Bottom Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(6.52), Inches(12.13), Inches(0.44))
    banner.fill.solid()
    banner.fill.fore_color.rgb = FILL_LIGHT_GREEN
    banner.line.color.rgb = BORDER_GREEN
    banner.line.width = Pt(1.0)
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "• Pernyataan Penutup: "
    r1.font.name = 'Arial'
    r1.font.size = Pt(7.5)
    r1.font.bold = True
    r1.font.color.rgb = DARK_GREEN
    r2 = p.add_run()
    r2.text = "Kaizen DSS SIRINE 4.0 membuktikan bahwa transformasi digital berbasis data granular tidak hanya menghentikan pemborosan finansial miliaran rupiah secara instan (0 Bulan Payback), melainkan merevolusi tata kelola lapangan Perum Peruri menjadi lebih presisi, akuntabel, dan berdaya saing tinggi."
    r2.font.name = 'Arial'
    r2.font.size = Pt(7.4)
    r2.font.color.rgb = DARK_TEXT

# -------------------------------------------------------------
# MAIN EXECUTION
# -------------------------------------------------------------
def main():
    prs = pptx.Presentation('Kerangka_Presentasi_Peserta_IAKA_2026 (1).pptx')
    print(f"Total slides initially: {len(prs.slides)}")

    # In the current deck:
    # idx 17: Slide 18 is Point 08 template placeholder
    # idx 18: Slide 19 is Closing ("TERIMA KASIH")

    # Slide 08.1 will be built on slide index 17
    slide_8_1 = prs.slides[17]
    build_slide_8_1(slide_8_1)
    print("Slide 08.1 built successfully on index 17.")

    # Check if slide 08.2 already exists or needs to be inserted
    if len(prs.slides) == 19:
        new_slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Move new_slide to index 18 (right after index 17, before closing)
        sldIdLst = prs.slides._sldIdLst
        slide_elem = sldIdLst[-1]
        sldIdLst.remove(slide_elem)
        sldIdLst.insert(18, slide_elem)
        slide_8_2 = prs.slides[18]
    else:
        slide_8_2 = prs.slides[18]

    build_slide_8_2(slide_8_2)
    print("Slide 08.2 built successfully on index 18.")

    prs.save('Kerangka_Presentasi_Peserta_IAKA_2026 (1).pptx')
    print(f"Saved! Kerangka_Presentasi_Peserta_IAKA_2026 (1).pptx updated. Total slides: {len(prs.slides)}")

if __name__ == '__main__':
    main()
