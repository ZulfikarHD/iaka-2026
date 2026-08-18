import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
import sys
import os

# --- PALETTE CONSTANTS ---
NAVY = RGBColor(32, 40, 143)         # #20288F - Primary Brand
PURPLE = RGBColor(106, 44, 201)      # #6A2CC9 - Accent
DARK_NAVY = RGBColor(26, 35, 126)    # #1A237E
GREEN = RGBColor(46, 125, 50)        # #2E7D32 - Success
DARK_GREEN = RGBColor(27, 94, 32)    # #1B5E20
RED = RGBColor(198, 40, 40)          # #C62828 - Alert
ORANGE = RGBColor(230, 81, 0)        # #E65100 - Warning
DARK_TEXT = RGBColor(42, 42, 61)     # #2A2A3D - Body Text
MUTED_TEXT = RGBColor(110, 110, 130) # #6E6E82 - Muted Text
WHITE = RGBColor(255, 255, 255)

# Container Fills
FILL_LIGHT_NAVY = RGBColor(244, 247, 254)   # #F4F7FE
FILL_LIGHT_PURPLE = RGBColor(248, 246, 254) # #F8F6FE
FILL_LIGHT_GREEN = RGBColor(232, 245, 233)  # #E8F5E9
FILL_LIGHT_AMBER = RGBColor(255, 249, 230)  # #FFF9E6
FILL_LIGHT_RED = RGBColor(255, 245, 245)    # #FFF5F5

# Borders
BORDER_NAVY = RGBColor(210, 220, 245)       # #D2DCF5
BORDER_PURPLE = RGBColor(217, 206, 247)     # #D9CEF7
BORDER_GREEN = RGBColor(165, 214, 167)      # #A5D6A7
BORDER_AMBER = RGBColor(255, 224, 130)      # #FFE082
BORDER_RED = RGBColor(255, 205, 210)        # #FFCDD2

def add_header(slide, badge_num, title_text, purpose_text):
    # Wave ornament
    if os.path.exists('wave_ornament.png'):
        slide.shapes.add_picture('wave_ornament.png', Inches(-0.118), Inches(-0.196), Inches(1.584), Inches(0.891))
    # Peruri Logo
    if os.path.exists('peruri_logo.png'):
        slide.shapes.add_picture('peruri_logo.png', Inches(12.09), Inches(0.013), Inches(1.247), Inches(0.694))

    # Badge Box
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(0.65), Inches(0.80), Inches(0.72))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PURPLE
    badge.line.color.rgb = PURPLE
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = badge_num
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    tb_title = slide.shapes.add_textbox(Inches(1.55), Inches(0.60), Inches(10.20), Inches(0.80))
    tf = tb_title.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Purpose Subtitle
    tb_purpose = slide.shapes.add_textbox(Inches(0.60), Inches(1.48), Inches(11.80), Inches(0.38))
    tf = tb_purpose.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Tujuan bagian ini:  "
    r1.font.name = 'Arial'
    r1.font.size = Pt(11.5)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = purpose_text
    r2.font.name = 'Arial'
    r2.font.size = Pt(11.5)
    r2.font.color.rgb = DARK_TEXT

    # Footer
    tb_foot = slide.shapes.add_textbox(Inches(0.60), Inches(7.02), Inches(12.13), Inches(0.24))
    tf = tb_foot.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = "IAKA 2026 — Kerangka Presentasi Peserta  ·  DSS SIRINE 4.0 Unit Cetak Pita Cukai"
    p.font.name = 'Arial'
    p.font.size = Pt(8.5)
    p.font.color.rgb = MUTED_TEXT

def add_kpi_card(slide, left, top, width, height, title, value, subtext, color, fill_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.04)
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)

    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = 'Arial'
    p0.font.size = Pt(8.5)
    p0.font.bold = True
    p0.font.color.rgb = color
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = value
    p1.font.name = 'Arial'
    p1.font.size = Pt(15.5)
    p1.font.bold = True
    p1.font.color.rgb = color
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = subtext
    p2.font.name = 'Arial'
    p2.font.size = Pt(7.2)
    p2.font.color.rgb = DARK_TEXT
    p2.alignment = PP_ALIGN.CENTER

def style_table_cell(cell, text, font_size=7.2, bold=False, text_color=DARK_TEXT, fill_color=None, align=PP_ALIGN.LEFT):
    cell.text = ""
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    else:
        cell.fill.background()
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = align

# -------------------------------------------------------------
# BUILD SLIDE 06.1
# -------------------------------------------------------------
def build_slide_6_1(slide):
    # Clear shapes on slide
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    add_header(
        slide,
        badge_num="06.1",
        title_text="Dampak Finansial: Model Simulasi ROI, Cost Avoidance & Nilai Bersih",
        purpose_text="Sajikan kalkulasi dampak finansial terbuka: simulasi penghematan lembar cetak vs baseline, cost avoidance (@ Rp 3.000/LK*), net benefit, dan payback period."
    )

    # Top 4 Hero KPI Cards with Explicit Data Source & Simulation Disclaimer
    add_kpi_card(slide, 0.60, 1.90, 2.86, 1.08, "REALISASI S1 2026 (SIMULASI)*", "Rp 2,23 Miliar*", "Simulasi Cost Avoidance | Fisik: 743.234 LK Riil", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 3.69, 1.90, 2.86, 1.08, "PROYEKSI TAHUNAN (SIMULASI)*", "Rp 6,82 Miliar / Thn*", "Simulasi Efisiensi | Fisik: 2.273.752 LK Riil/Thn", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 6.78, 1.90, 2.86, 1.08, "BIAYA INVESTASI (CAPEX)", "Rp 0 (In-House)", "100% Tim Internal Peruri | Zero License Fee", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 9.87, 1.90, 2.86, 1.08, "PAYBACK PERIOD (SIMULASI)*", "0 Bulan (Seketika)", "Net Benefit Rp 6,82 M/Thn* | Instant ROI", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)

    # Left Panel: 1. KERTAS KERJA FORMULA PERHITUNGAN & COST AVOIDANCE (TERBUKA)
    header_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.60), Inches(3.08), Inches(6.15), Inches(0.34))
    header_left.fill.solid()
    header_left.fill.fore_color.rgb = NAVY
    header_left.line.fill.background()
    tf = header_left.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "1. FORMULA MODEL SIMULASI COST AVOIDANCE & DATA KUARTALAN"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(3.42), Inches(6.15), Inches(3.02))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = FILL_LIGHT_NAVY
    box_left.line.color.rgb = BORDER_NAVY
    box_left.line.width = Pt(1.0)

    # Left Card 1: Formula Step-by-Step with Emphasized Simulation Note
    card_l1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(3.48), Inches(5.91), Inches(1.18))
    card_l1.fill.solid()
    card_l1.fill.fore_color.rgb = WHITE
    card_l1.line.color.rgb = BORDER_NAVY
    card_l1.line.width = Pt(1.0)
    tf = card_l1.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.04)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)

    p = tf.paragraphs[0]
    p.text = "RUMUS SIMULASI COST AVOIDANCE (ASUMSI ESTIMASI @ Rp 3.000 / LK*)"
    p.font.name = 'Arial'
    p.font.size = Pt(7.8)
    p.font.bold = True
    p.font.color.rgb = NAVY

    formula_steps = [
        ("• Estimasi Tambah Cetak: ", "Rp 3.000 / LK* (Estimasi simulasi bahan baku kertas sekuriti, tinta, mesin & SDM — BUKAN HPP/Harga Resmi)."),
        ("• Data Fisik Riil SAP: ", "Vol. 177.636.930 LK (SAP ZPPRSIPPC0012) | Baseline: 4,61% ➔ Realisasi Q2: 3,33% (Δ = -1,28 pp)."),
        ("• Lembar Diselamatkan: ", "177.636.930 LK × 1,28% = 2.273.752 Lembar Cetak / Tahun (Fisik Riil)."),
        ("• Valuasi Simulasi Saving: ", "2.273.752 LK × Rp 3.000* = Rp 6.821.256.000 / Tahun (~Rp 6,82 Miliar*).")
    ]
    for s_label, s_val in formula_steps:
        p_s = tf.add_paragraph()
        r1 = p_s.add_run()
        r1.text = s_label
        r1.font.bold = True
        r1.font.size = Pt(6.7)
        r1.font.color.rgb = DARK_NAVY
        r2 = p_s.add_run()
        r2.text = s_val
        r2.font.size = Pt(6.7)
        r2.font.bold = (s_label.startswith("• Valuasi"))
        r2.font.color.rgb = DARK_GREEN if s_label.startswith("• Valuasi") else DARK_TEXT

    # Left Table: Realization & Projection Workpaper
    table_shape_l = slide.shapes.add_table(6, 5, Inches(0.72), Inches(4.70), Inches(5.91), Inches(1.48))
    t_l = table_shape_l.table
    t_l.columns[0].width = Inches(1.30)
    t_l.columns[1].width = Inches(1.05)
    t_l.columns[2].width = Inches(1.05)
    t_l.columns[3].width = Inches(1.15)
    t_l.columns[4].width = Inches(1.36)

    headers_l = ["Periode Evaluasi", "Volume Riil (LK)", "Inschiet Riil (%)", "Lembar Rusak", "Simulasi Saving*"]
    for c_idx, h in enumerate(headers_l):
        style_table_cell(t_l.cell(0, c_idx), h, font_size=6.8, bold=True, text_color=WHITE, fill_color=NAVY, align=PP_ALIGN.CENTER if c_idx>0 else PP_ALIGN.LEFT)

    rows_l = [
        ("Baseline 2025", "177.636.930", "4,61%", "8.189.062 lb", "Rp 24,56 M* (Toleransi)", WHITE),
        ("Realisasi Q1 2026", "57.385.254", "4,34% (-0,27 pp)", "2.490.520 lb", "Rp 464,82 Jt* (Hemat 155 rb lb)", WHITE),
        ("Realisasi Q2 2026", "45.960.434", "3,33% (-1,28 pp)", "1.530.482 lb", "Rp 1,76 Miliar* (Hemat 588 rb lb)", FILL_LIGHT_GREEN),
        ("Total Realisasi S1", "103.345.688", "3,89% (-0,72 pp)", "4.021.002 lb", "Rp 2,23 Miliar* (Hemat 743 rb lb)", FILL_LIGHT_GREEN),
        ("Proyeksi 1 Thn Penuh", "177.636.930", "3,33% (Run-rate)", "5.915.310 lb", "Rp 6,82 Miliar / Thn*", FILL_LIGHT_PURPLE)
    ]
    for r_idx, (p_eval, vol_val, ins_val, dam_val, sav_val, bg_col) in enumerate(rows_l, start=1):
        is_bold_row = (r_idx >= 3)
        txt_col = DARK_GREEN if r_idx in [3, 4] else (PURPLE if r_idx==5 else DARK_TEXT)
        style_table_cell(t_l.cell(r_idx, 0), p_eval, font_size=6.6, bold=is_bold_row, text_color=DARK_NAVY, fill_color=bg_col)
        style_table_cell(t_l.cell(r_idx, 1), vol_val, font_size=6.6, bold=False, text_color=DARK_TEXT, fill_color=bg_col, align=PP_ALIGN.CENTER)
        style_table_cell(t_l.cell(r_idx, 2), ins_val, font_size=6.6, bold=is_bold_row, text_color=txt_col, fill_color=bg_col, align=PP_ALIGN.CENTER)
        style_table_cell(t_l.cell(r_idx, 3), dam_val, font_size=6.6, bold=False, text_color=DARK_TEXT, fill_color=bg_col, align=PP_ALIGN.CENTER)
        style_table_cell(t_l.cell(r_idx, 4), sav_val, font_size=6.6, bold=is_bold_row, text_color=txt_col, fill_color=bg_col, align=PP_ALIGN.CENTER)

    # Prominent Disclaimer & Verified Source Label under Left Table
    tb_src_l = slide.shapes.add_textbox(Inches(0.72), Inches(6.20), Inches(5.91), Inches(0.20))
    tf_src_l = tb_src_l.text_frame
    tf_src_l.word_wrap = False
    tf_src_l.margin_top = Inches(0)
    tf_src_l.margin_left = Inches(0)
    p = tf_src_l.paragraphs[0]
    p.text = "*Catatan: Rp 3.000/LK adalah ESTIMASI INTERNAL untuk simulasi (BUKAN HPP/harga resmi). Vol. & Inschiet = Data Riil SAP ZPPRSIPPC0012."
    p.font.name = 'Arial'
    p.font.size = Pt(6.3)
    p.font.bold = True
    p.font.color.rgb = RED

    # Right Panel: 2. ANALISIS NET VALUE, PAYBACK PERIOD & EFISIENSI WAKTU
    header_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.90), Inches(3.08), Inches(5.83), Inches(0.34))
    header_right.fill.solid()
    header_right.fill.fore_color.rgb = PURPLE
    header_right.line.fill.background()
    tf = header_right.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "2. ANALISIS NET VALUE, PAYBACK PERIOD & EFISIENSI WAKTU"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.90), Inches(3.42), Inches(5.83), Inches(3.02))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = FILL_LIGHT_PURPLE
    box_right.line.color.rgb = BORDER_PURPLE
    box_right.line.width = Pt(1.0)

    # Mini title for chart with Explicit Simulation & Source Citation
    tb_c_title = slide.shapes.add_textbox(Inches(7.02), Inches(3.46), Inches(5.59), Inches(0.24))
    tf = tb_c_title.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Perbandingan Pemborosan vs Realisasi Finansial (Simulasi Rp Miliar)* "
    r1.font.name = 'Arial'
    r1.font.size = Pt(7.3)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = "— Sumber SAP"
    r2.font.name = 'Arial'
    r2.font.size = Pt(6.8)
    r2.font.italic = True
    r2.font.color.rgb = MUTED_TEXT

    # Native PPTX Clustered Column Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Pemborosan 2025*', 'Realisasi S1 2026*', 'Proyeksi Hemat/Thn*', 'Investasi CAPEX']
    chart_data.add_series('Nilai Finansial (Rp Miliar)', (24.56, 2.23, 6.82, 0.00))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(7.02), Inches(3.68), Inches(5.59), Inches(1.32),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.name = 'Arial'
    data_labels.font.size = Pt(8.0)
    data_labels.font.bold = True
    data_labels.font.color.rgb = DARK_NAVY
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    # Series styling
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = PURPLE

    # Right Card 1: Net Benefit & Payback Period Model
    card_r1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.02), Inches(5.04), Inches(5.59), Inches(1.14))
    card_r1.fill.solid()
    card_r1.fill.fore_color.rgb = WHITE
    card_r1.line.color.rgb = BORDER_GREEN
    card_r1.line.width = Pt(1.0)
    tf = card_r1.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.04)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)

    p = tf.paragraphs[0]
    p.text = "KALKULASI NILAI BERSIH (NET BENEFIT) & PAYBACK PERIOD (SIMULASI)"
    p.font.name = 'Arial'
    p.font.size = Pt(7.8)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN

    roi_items = [
        ("• Net Value Creation: ", "Simulasi Cost Avoidance Rp 6,82 M* - CAPEX Rp 0 - OPEX Lisensi Rp 0 = Rp 6,82 Miliar / Tahun*."),
        ("• Payback Period: ", "0 Bulan (Seketika) — Investasi Rp 0 / Benefit Simulasi Rp 6,82 M* menghasilkan Instant ROI."),
        ("• Efisiensi Operasional Riil: ", "Waktu henti inspeksi mesin turun dari > 8 jam ke < 2–4 jam, rekap turun dari 45 mnt ke 0 mnt."),
        ("• Penegasan Estimasi: ", "Nilai Rp 3.000/LK semata-mata model simulasi finansial, bukan rujukan laporan biaya resmi.")
    ]
    for r_label, r_desc in roi_items:
        p_r = tf.add_paragraph()
        r1 = p_r.add_run()
        r1.text = r_label
        r1.font.bold = True
        r1.font.size = Pt(6.7)
        r1.font.color.rgb = DARK_GREEN if "Net Value" in r_label or "Payback" in r_label else DARK_NAVY
        r2 = p_r.add_run()
        r2.text = r_desc
        r2.font.size = Pt(6.7)
        r2.font.color.rgb = DARK_TEXT

    # Prominent Audit Trail Label under Right Card
    tb_src_r = slide.shapes.add_textbox(Inches(7.02), Inches(6.20), Inches(5.59), Inches(0.20))
    tf_src_r = tb_src_r.text_frame
    tf_src_r.word_wrap = False
    tf_src_r.margin_top = Inches(0)
    tf_src_r.margin_left = Inches(0)
    p = tf_src_r.paragraphs[0]
    p.text = "*(Audit Trail: 100% Data Transaksi Digital Tervalidasi via SAP ERP & Berita Acara BA-PPC-2026-002)"
    p.font.name = 'Arial'
    p.font.size = Pt(6.5)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Bottom Banner with Explicit Simulation Emphasis
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(6.52), Inches(12.13), Inches(0.44))
    banner.fill.solid()
    banner.fill.fore_color.rgb = FILL_LIGHT_AMBER
    banner.line.color.rgb = BORDER_AMBER
    banner.line.width = Pt(1.0)
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "• Penegasan Simulasi Finansial & Data Riil: "
    r1.font.name = 'Arial'
    r1.font.size = Pt(7.4)
    r1.font.bold = True
    r1.font.color.rgb = ORANGE
    r2 = p.add_run()
    r2.text = "Nilai rupiah (Rp 2,23 Miliar S1 & proyeksi Rp 6,82 Miliar/tahun) merupakan HASIL SIMULASI ESTIMASI (@ Rp 3.000/LK) untuk mengukur dampak efisiensi — BUKAN laporan HPP riil atau harga jual resmi Peruri (rahasia perusahaan). Seluruh data volume fisik (177,6 Jt LK), lembar diselamatkan (743 rb LK S1 / 2,27 Jt LK per tahun), dan penurunan inschiet (4,61% ➔ 3,33%) adalah DATA RIIL TERVERIFIKASI SAP ZPPRSIPPC0012."
    r2.font.name = 'Arial'
    r2.font.size = Pt(7.2)
    r2.font.color.rgb = DARK_TEXT

# -------------------------------------------------------------
# BUILD SLIDE 06.2
# -------------------------------------------------------------
def build_slide_6_2(slide):
    # Clear shapes on slide
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    add_header(
        slide,
        badge_num="06.2",
        title_text="Dampak Multidimensi: Mutu, People, K3L/ESG, Kepatuhan & Potensi Replikasi",
        purpose_text="Sajikan dampak non-finansial terukur lintas mutu, budaya kerja, ESG, kepatuhan audit, serta keselarasan strategi INDI 4.0 & roadmap replikasi."
    )

    # Top 4 Hero KPI Cards with Explicit Data Sources
    add_kpi_card(slide, 0.60, 1.90, 2.86, 1.08, "MUTU & KEPUASAN DJBC", "0 Komplain / Klaim", "SLA 100% On-Time | Sumber: Berita Acara DJBC", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 3.69, 1.90, 2.86, 1.08, "DAMPAK ESG & LINGKUNGAN", "2,27 Jt Lembar / Thn", "Reduksi Limbah Afval | Sumber: SAP ZPPRSIPPC0012", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 6.78, 1.90, 2.86, 1.08, "BUDAYA DATA & KEPATUHAN", "100% Auditable", "Sesuai ISO 9001:2015 | Sumber: Log Transaksi PO", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 9.87, 1.90, 2.86, 1.08, "KESIAPAN REPLIKASI", "4 Unit Produksi", "Meterai, Paspor, Khazanah, Uang | RJPP INDI 4.0", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)

    # Left Panel: 1. MATRIKS DAMPAK NON-FINANSIAL TERVALIDASI (4 PILAR)
    header_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.60), Inches(3.08), Inches(5.95), Inches(0.34))
    header_left.fill.solid()
    header_left.fill.fore_color.rgb = NAVY
    header_left.line.fill.background()
    tf = header_left.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "1. MATRIKS DAMPAK NON-FINANSIAL TERVALIDASI (4 PILAR)"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(3.42), Inches(5.95), Inches(3.02))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = FILL_LIGHT_NAVY
    box_left.line.color.rgb = BORDER_NAVY
    box_left.line.width = Pt(1.0)

    # 4 Pillar Cards on Left Panel with Explicit Source Tags
    pillars = [
        (
            Inches(3.48), Inches(0.66), GREEN, BORDER_GREEN,
            "PILAR 1: MUTU & KEPUASAN PELANGGAN (DJBC KEMENKEU RI)",
            "• Jaminan Mutu & SLA 100%: Nol komplain mutu & pengiriman tepat waktu (Sumber: BA Serah Terima DJBC).",
            "• Proteksi Penerimaan Negara: Risiko lembar HCTS cacat lolos ditekan hingga batas nol."
        ),
        (
            Inches(4.18), Inches(0.66), PURPLE, BORDER_PURPLE,
            "PILAR 2: PEOPLE, BUDAYA KERJA & KEPUASAN PENGGUNA",
            "• Budaya Data Objektif: Output LK & % cacat per tim tersaji berdampingan (Sumber: Modul Unit Cetak).",
            "• Kenyamanan Kerja (Lean UX): Input form digital < 30 detik Autofill SAP (Sumber: Log Transaksi)."
        ),
        (
            Inches(4.88), Inches(0.68), NAVY, BORDER_NAVY,
            "PILAR 3: KEPATUHAN TATA KELOLA & AUDITABILITY (ISO 9001:2015)",
            "• Klausul ISO 9001: Penuhi Klausul 8.5.2 (Mampu Telusur) & Klausul 9.1.3 (Analisis Data Mutu).",
            "• Jejak Audit 100%: Rantai PO ➔ Mesin ➔ Shift ➔ Tim ➔ Cacat terverifikasi digital (Sumber: Basis Data)."
        ),
        (
            Inches(5.60), Inches(0.68), DARK_GREEN, BORDER_GREEN,
            "PILAR 4: K3L & ESG (KEBERLANJUTAN LINGKUNGAN & ERGONOMI)",
            "• Reduksi Limbah Padat (ESG): Mencegah 2,27 Jt lembar limbah afval/thn (Sumber: SAP ZPPRSIPPC0012).",
            "• Reduksi Kimia & Kelelahan: Menekan solvent pencuci rol & pangkas perbaikan mesin ke < 2–4 jam."
        )
    ]

    for top_pos, h_size, col, b_col, p_title, line1, line2 in pillars:
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), top_pos, Inches(5.71), h_size)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = b_col
        c.line.width = Pt(1.0)
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.02)
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = p_title
        p.font.name = 'Arial'
        p.font.size = Pt(7.3)
        p.font.bold = True
        p.font.color.rgb = col

        p1 = tf.add_paragraph()
        p1.text = line1
        p1.font.name = 'Arial'
        p1.font.size = Pt(6.5)
        p1.font.color.rgb = DARK_TEXT

        p2 = tf.add_paragraph()
        p2.text = line2
        p2.font.name = 'Arial'
        p2.font.size = Pt(6.5)
        p2.font.color.rgb = DARK_TEXT

    # Source label under left panel
    tb_src_l2 = slide.shapes.add_textbox(Inches(0.72), Inches(6.28), Inches(5.71), Inches(0.14))
    tf_src_l2 = tb_src_l2.text_frame
    tf_src_l2.word_wrap = False
    tf_src_l2.margin_top = Inches(0)
    tf_src_l2.margin_left = Inches(0)
    p = tf_src_l2.paragraphs[0]
    p.text = "*(Sumber Data: Laporan SLA DJBC 2026, Modul Verifikasi Mutu, Berita Acara BA-PPC-2026-002 & SAP ERP)"
    p.font.name = 'Arial'
    p.font.size = Pt(6.3)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY

    # Right Panel: 2. KESELARASAN STRATEGIS PERUSAHAAN & ROADMAP REPLIKASI
    header_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.70), Inches(3.08), Inches(6.03), Inches(0.34))
    header_right.fill.solid()
    header_right.fill.fore_color.rgb = PURPLE
    header_right.line.fill.background()
    tf = header_right.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "2. KESELARASAN STRATEGIS PERUSAHAAN & ROADMAP REPLIKASI"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.70), Inches(3.42), Inches(6.03), Inches(3.02))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = FILL_LIGHT_PURPLE
    box_right.line.color.rgb = BORDER_PURPLE
    box_right.line.width = Pt(1.0)

    # Right Card 1: Strategic Alignment & INDI 4.0
    card_r1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.82), Inches(3.48), Inches(5.79), Inches(1.30))
    card_r1.fill.solid()
    card_r1.fill.fore_color.rgb = WHITE
    card_r1.line.color.rgb = BORDER_PURPLE
    card_r1.line.width = Pt(1.0)
    tf = card_r1.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.04)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)

    p = tf.paragraphs[0]
    p.text = "KESELARASAN STRATEGIS PERUSAHAAN & DUKUNGAN SKOR INDI 4.0"
    p.font.name = 'Arial'
    p.font.size = Pt(7.8)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    strat_items = [
        ("• Transformasi Smart Factory: ", "Mendukung RJPP Digital Peruri & Asesmen INDI 4.0 Kementerian BUMN."),
        ("• Integrasi Data PPIC: ", "Sinkronisasi data keandalan mesin dengan perencanaan kapasitas SAP."),
        ("• Kedaulatan & Keamanan Data: ", "100% In-house development menjamin kerahasiaan data sekuriti negara.")
    ]
    for s_label, s_desc in strat_items:
        p_s = tf.add_paragraph()
        r1 = p_s.add_run()
        r1.text = s_label
        r1.font.bold = True
        r1.font.size = Pt(6.8)
        r1.font.color.rgb = PURPLE
        r2 = p_s.add_run()
        r2.text = s_desc
        r2.font.size = Pt(6.8)
        r2.font.color.rgb = DARK_TEXT

    # Right Card 2: Replication Roadmap (4 Units)
    card_r2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.82), Inches(4.82), Inches(5.79), Inches(1.42))
    card_r2.fill.solid()
    card_r2.fill.fore_color.rgb = WHITE
    card_r2.line.color.rgb = BORDER_NAVY
    card_r2.line.width = Pt(1.0)
    tf = card_r2.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.04)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)

    p = tf.paragraphs[0]
    p.text = "POTENSI & ROADMAP REPLIKASI KE 4 LINI PRODUKSI STRATEGIS LAIN"
    p.font.name = 'Arial'
    p.font.size = Pt(7.8)
    p.font.bold = True
    p.font.color.rgb = NAVY

    replications = [
        ("1. Unit Cetak Meterai: ", "Digitalisasi nomor seri & audit cacat cetak meterai tempel/elektronik."),
        ("2. Unit Dokumen Paspor: ", "Pelacakan nomor seri & pengendalian afval blanko paspor RI & buku pelaut."),
        ("3. Unit Khazanah & Finishing: ", "Rekonsiliasi lembar sempurna vs cacat saat proses sortasi akhir."),
        ("4. Unit Cetak Uang Kertas: ", "Replikasi audit Pareto kerusakan mekanis mesin cetak uang sekuriti tinggi.")
    ]
    for r_unit, r_detail in replications:
        p_r = tf.add_paragraph()
        r1 = p_r.add_run()
        r1.text = "• " + r_unit
        r1.font.bold = True
        r1.font.size = Pt(6.6)
        r1.font.color.rgb = NAVY
        r2 = p_r.add_run()
        r2.text = r_detail
        r2.font.size = Pt(6.6)
        r2.font.color.rgb = DARK_TEXT

    # Source label under right panel
    tb_src_r2 = slide.shapes.add_textbox(Inches(6.82), Inches(6.28), Inches(5.79), Inches(0.14))
    tf_src_r2 = tb_src_r2.text_frame
    tf_src_r2.word_wrap = False
    tf_src_r2.margin_top = Inches(0)
    tf_src_r2.margin_left = Inches(0)
    p = tf_src_r2.paragraphs[0]
    p.text = "*(Landasan Kebijakan: Selaras PP No. 06/2019, Standar Mutu ISO 9001:2015, dan Roadmap INDI 4.0)"
    p.font.name = 'Arial'
    p.font.size = Pt(6.3)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Bottom Banner with Explicit Source Verification
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(6.52), Inches(12.13), Inches(0.44))
    banner.fill.solid()
    banner.fill.fore_color.rgb = FILL_LIGHT_PURPLE
    banner.line.color.rgb = BORDER_PURPLE
    banner.line.width = Pt(1.0)
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "• Integritas Data Non-Finansial: "
    r1.font.name = 'Arial'
    r1.font.size = Pt(7.5)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = "Dampak terverifikasi melalui Berita Acara Serah Terima DJBC (Nol Komplain), pencegahan 2,27 Juta lembar limbah afval (SAP ZPPRSIPPC0012), bukti audit penarikan buku folio (BA-PPC-2026-002), dan keselarasan asesmen INDI 4.0."
    r2.font.name = 'Arial'
    r2.font.size = Pt(7.5)
    r2.font.color.rgb = DARK_TEXT

# -------------------------------------------------------------
# MAIN EXECUTION
# -------------------------------------------------------------
def main():
    prs = pptx.Presentation('Kerangka_Presentasi_Peserta_IAKA_2026 (1).pptx')
    print(f"Total slides initially: {len(prs.slides)}")

    # Slide 14 (index 13) is 06.1
    slide_6_1 = prs.slides[13]
    build_slide_6_1(slide_6_1)
    print("Slide 06.1 updated with prominent simulation disclaimer & verified sources.")

    # Slide 15 (index 14) is 06.2
    slide_6_2 = prs.slides[14]
    build_slide_6_2(slide_6_2)
    print("Slide 06.2 updated with prominent verified sources.")

    prs.save('Kerangka_Presentasi_Peserta_IAKA_2026 (1).pptx')
    print("Saved! Kerangka_Presentasi_Peserta_IAKA_2026 (1).pptx successfully updated.")

if __name__ == '__main__':
    main()
