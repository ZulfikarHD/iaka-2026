import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
import sys
import os

# ==========================================
# PALETTE & DESIGN CONSTANTS (Peruri Theme)
# ==========================================
NAVY = RGBColor(32, 40, 143)         # #20288F - Primary Brand Navy
PURPLE = RGBColor(106, 44, 201)      # #6A2CC9 - Accent Purple
DARK_NAVY = RGBColor(26, 35, 126)    # #1A237E - Deep Navy
GREEN = RGBColor(46, 125, 50)        # #2E7D32 - Success Green
DARK_GREEN = RGBColor(27, 94, 32)    # #1B5E20
RED = RGBColor(198, 40, 40)          # #C62828 - Alert Red
ORANGE = RGBColor(230, 81, 0)        # #E65100 - Warning Amber
DARK_TEXT = RGBColor(42, 42, 61)     # #2A2A3D - Body Text Charcoal
MUTED_TEXT = RGBColor(110, 110, 130) # #6E6E82 - Muted Text Gray
WHITE = RGBColor(255, 255, 255)

# Container Fills
FILL_LIGHT_NAVY = RGBColor(244, 247, 254)   # #F4F7FE
FILL_LIGHT_PURPLE = RGBColor(248, 246, 254) # #F8F6FE
FILL_LIGHT_GREEN = RGBColor(232, 245, 233)  # #E8F5E9
FILL_LIGHT_AMBER = RGBColor(255, 249, 230)  # #FFF9E6
FILL_LIGHT_RED = RGBColor(255, 245, 245)    # #FFF5F5

# Container Borders
BORDER_NAVY = RGBColor(210, 220, 245)       # #D2DCF5
BORDER_PURPLE = RGBColor(217, 206, 247)     # #D9CEF7
BORDER_GREEN = RGBColor(165, 214, 167)      # #A5D6A7
BORDER_AMBER = RGBColor(255, 224, 130)      # #FFE082
BORDER_RED = RGBColor(255, 205, 210)        # #FFCDD2

def create_presentation():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide layout
    return prs, blank_layout

def add_header(slide, badge_num, title_text, purpose_text):
    # Wave ornament top left
    if os.path.exists('wave_ornament.png'):
        slide.shapes.add_picture('wave_ornament.png', Inches(-0.118), Inches(-0.196), Inches(1.584), Inches(0.891))
    
    # Peruri Logo top right
    if os.path.exists('peruri_logo.png'):
        slide.shapes.add_picture('peruri_logo.png', Inches(12.09), Inches(0.013), Inches(1.247), Inches(0.694))

    # Badge Box
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(0.65), Inches(0.80), Inches(0.72))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PURPLE
    badge.line.color.rgb = PURPLE
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = badge_num
    p.font.name = 'Arial'
    p.font.size = Pt(20 if len(badge_num) <= 3 else 16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    tb_title = slide.shapes.add_textbox(Inches(1.55), Inches(0.60), Inches(10.20), Inches(0.80))
    tf = tb_title.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Purpose Subtitle
    tb_purpose = slide.shapes.add_textbox(Inches(0.60), Inches(1.48), Inches(12.13), Inches(0.38))
    tf = tb_purpose.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Tujuan bagian ini:  "
    r1.font.name = 'Arial'
    r1.font.size = Pt(11.5)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = purpose_text
    r2.font.name = 'Arial'
    r2.font.size = Pt(11.5)
    r2.font.color.rgb = DARK_TEXT

    # Footer
    tb_foot = slide.shapes.add_textbox(Inches(0.60), Inches(7.04), Inches(12.13), Inches(0.24))
    tf = tb_foot.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = "IAKA 2026 — Kerangka Presentasi Peserta  ·  DSS SIRINE 4.0 Unit Cetak Pita Cukai"
    p.font.name = 'Arial'
    p.font.size = Pt(8.5)
    p.font.color.rgb = MUTED_TEXT

def add_kpi_card(slide, left, top, width, height, title, value, subtext, color, fill_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(left + 0.12), Inches(top + 0.10), Inches(width - 0.24), Inches(height - 0.16))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)

    p1 = tf.paragraphs[0]
    p1.text = title.upper()
    p1.font.name = 'Arial'
    p1.font.size = Pt(8.5)
    p1.font.bold = True
    p1.font.color.rgb = color

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.name = 'Arial'
    p2.font.size = Pt(17)
    p2.font.bold = True
    p2.font.color.rgb = color

    p3 = tf.add_paragraph()
    p3.text = subtext
    p3.font.name = 'Arial'
    p3.font.size = Pt(8.5)
    p3.font.color.rgb = DARK_TEXT

def add_card(slide, left, top, width, height, fill_color, border_color, border_width=1.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(border_width)
    return card

def add_section_header(slide, left, top, width, height, title, bg_color):
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = bg_color
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    return hdr

# ==========================================
# SLIDE 1: COVER SLIDE
# ==========================================
def build_slide_1(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)

    # Top wave & Peruri logo
    if os.path.exists('wave_ornament.png'):
        slide.shapes.add_picture('wave_ornament.png', Inches(-0.118), Inches(-0.196), Inches(2.2), Inches(1.238))
    if os.path.exists('peruri_logo.png'):
        slide.shapes.add_picture('peruri_logo.png', Inches(11.4), Inches(0.35), Inches(1.4), Inches(0.78))

    # Hero Main Card
    main_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(5.3))
    main_card.fill.solid()
    main_card.fill.fore_color.rgb = FILL_LIGHT_NAVY
    main_card.line.color.rgb = BORDER_NAVY
    main_card.line.width = Pt(2.0)

    # Decorative top bar on card
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = NAVY
    top_bar.line.fill.background()

    # Event Tag Box
    tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.8), Inches(4.8), Inches(0.42))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = PURPLE
    tag_box.line.fill.background()
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = False
    tf_tag.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "INNOVATION & KAIZEN AWARD (IAKA) 2026"
    p_tag.font.name = 'Arial'
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = WHITE
    p_tag.alignment = PP_ALIGN.CENTER

    # Main Title Box
    tb_title = slide.shapes.add_textbox(Inches(1.2), Inches(2.4), Inches(10.9), Inches(1.4))
    tf_t = tb_title.text_frame
    tf_t.word_wrap = True
    tf_t.margin_top = Inches(0)
    tf_t.margin_left = Inches(0)
    p_t = tf_t.paragraphs[0]
    p_t.text = "RISALAH LATAR BELAKANG & IDENTIFIKASI MASALAH"
    p_t.font.name = 'Arial'
    p_t.font.size = Pt(28)
    p_t.font.bold = True
    p_t.font.color.rgb = NAVY

    # Subtitle Box
    tb_sub = slide.shapes.add_textbox(Inches(1.2), Inches(3.75), Inches(10.9), Inches(1.1))
    tf_s = tb_sub.text_frame
    tf_s.word_wrap = True
    tf_s.margin_top = Inches(0)
    tf_s.margin_left = Inches(0)
    p_s = tf_s.paragraphs[0]
    p_s.text = "Transformasi Tata Kelola Pengendalian Mutu Unit Cetak Pita Cukai: Menjawab Tuntutan Pasar Global, Komitmen Tender Nasional, dan Sasaran Strategis Peruri Melalui Integrasi Data Meja Mesin"
    p_s.font.name = 'Arial'
    p_s.font.size = Pt(13)
    p_s.font.italic = True
    p_s.font.color.rgb = DARK_TEXT

    # 3 Metadata Pillars Bottom
    meta_cards = [
        ("UNIT KERJA & LINI PRODUKSI", "Unit Cetak Pita Cukai", "Dept. Khazanah & Verifikasi (9 Mesin Cetak)", NAVY, FILL_LIGHT_PURPLE, BORDER_PURPLE),
        ("SISTEM INOVASI", "DSS SIRINE 4.0", "Decision Support System Terintegrasi", PURPLE, FILL_LIGHT_NAVY, BORDER_NAVY),
        ("SKALA & SASARAN STRATEGIS", "177.636.930 LK", "Integritas Fiskal DJBC Kemenkeu RI", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    ]
    for idx, (m_title, m_val, m_sub, m_col, m_fill, m_bord) in enumerate(meta_cards):
        cx = 1.2 + idx * 3.7
        add_kpi_card(slide, cx, 5.0, 3.5, 1.4, m_title, m_val, m_sub, m_col, m_fill, m_bord)

    # Footer
    tb_foot = slide.shapes.add_textbox(Inches(0.8), Inches(7.04), Inches(11.733), Inches(0.24))
    tf_f = tb_foot.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "PERUM PERCETAKAN UANG REPUBLIK INDONESIA  ·  IAKA 2026"
    p_f.font.name = 'Arial'
    p_f.font.size = Pt(8.5)
    p_f.font.color.rgb = MUTED_TEXT
    p_f.alignment = PP_ALIGN.CENTER

# ==========================================
# SLIDE 2: RINGKASAN EKSEKUTIF & HERO METRICS
# ==========================================
def build_slide_2(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "EXEC", "Ringkasan Eksekutif: Urgensi Strategis & Capaian Kunci DSS SIRINE 4.0",
               "Sajikan ringkasan komprehensif latar belakang, masalah operasional, intervensi sistem, dan hasil capaian nyata.")

    # 4 Top Hero Cards
    add_kpi_card(slide, 0.60, 1.95, 2.90, 1.15, "BASELINE INSCHIET 2025", "4,61%", "Puncak Q4: 5,11% | Vol: 177,6M LK", RED, FILL_LIGHT_RED, BORDER_RED)
    add_kpi_card(slide, 3.65, 1.95, 2.90, 1.15, "REALISASI S1 2026", "3,89%", "Q2: 3,33% (-27,8%) | Target: <4,00%", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 6.70, 1.95, 2.90, 1.15, "PENYELAMATAN FINANSIAL", "Rp 2,23 Miliar", "743.234 Lembar Cetak Diselamatkan", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 9.75, 1.95, 2.98, 1.15, "REDUKSI WAKTU HENTI", "< 2–4 Jam", "Turun 50%–75% (vs > 8 Jam Perbaikan)", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)

    # 3 Summary Pillar Panels
    panels = [
        ("1. KONTEKS STRATEGIS & MANDAT BISNIS", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY, [
            ("Mandat Negara & DJBC Kemenkeu:", " Mencetak Pita Cukai (PCHT & MMEA) sebagai instrumen fiskal resmi bernilai ratusan triliun rupiah."),
            ("Skala Produksi Masif:", " Pesanan aktual 2025 mencapai 177.636.930 lembar cetak (rata-rata tahunan 160.000.000 lembar)."),
            ("Standar Industri Sekuriti Tinggi:", " Nol deviasi mutu (*zero-defect*), kepatuhan SLA distribusi, dan akuntabilitas bahan baku (*zero leakage*).")
        ]),
        ("2. KENDALA OPERASIONAL & TITIK BUTA", RED, FILL_LIGHT_RED, BORDER_RED, [
            ("Pemisahan Aliran Data (Data Silo):", " Meja mesin mencatat di buku folio fisik; laporan SAP di kantor hanya menampilkan ringkasan global tanpa nomor mesin dan shift."),
            ("Pemeriksaan Spekulatif (> 8 Jam):", " Teknisi memeriksa 9 mesin secara coba-coba saat terjadi cacat mutu karena tidak ada atribusi mesin."),
            ("Skala Dampak Finansial:", " Potensi pemborosan biaya bahan baku & cetak mencapai Rp 22,13 M s.d. Rp 24,56 Miliar/tahun.")
        ]),
        ("3. SOLUSI DSS SIRINE 4.0 & HASIL VALIDASI", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN, [
            ("Integrasi 3 Titik Aliran Data:", " Mengintegrasikan data transaksi meja mesin (< 30 detik), SAP ZPPRSIPPC0012, dan verifikasi mutu HCTS."),
            ("Atribusi Granular & Tindakan Cepat:", " Runtutan runtut: PO -> Mesin (9 Mesin) -> Shift -> Operator -> Kategori Cacat."),
            ("Capaian Semester 1 2026:", " Inschiet turun ke 3,89% (Q2: 3,33%), saving Rp 2,23 Miliar, dan downtime perbaikan turun 50%–75%.")
        ])
    ]

    col_w = 3.95
    for idx, (p_title, p_col, p_fill, p_bord, bullets) in enumerate(panels):
        left_x = 0.60 + idx * 4.10
        add_card(slide, left_x, 3.25, col_w, 3.65, p_fill, p_bord)
        add_section_header(slide, left_x, 3.25, col_w, 0.38, p_title, p_col)

        tb = slide.shapes.add_textbox(Inches(left_x + 0.15), Inches(3.72), Inches(col_w - 0.30), Inches(3.10))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)

        for b_idx, (b_bold, b_desc) in enumerate(bullets):
            p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            r0 = p.add_run()
            r0.text = "• "
            r0.font.name = 'Arial'
            r0.font.size = Pt(10)
            r0.font.bold = True
            r0.font.color.rgb = p_col
            r1 = p.add_run()
            r1.text = b_bold
            r1.font.name = 'Arial'
            r1.font.size = Pt(10)
            r1.font.bold = True
            r1.font.color.rgb = DARK_TEXT
            r2 = p.add_run()
            r2.text = b_desc
            r2.font.name = 'Arial'
            r2.font.size = Pt(9.5)
            r2.font.color.rgb = DARK_TEXT

# ==========================================
# SLIDE 3: DINAMIKA PASAR GLOBAL
# ==========================================
def build_slide_3(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "01.1", "Dinamika Pasar Global: Standar Presisi & Keunggulan Percetakan Sekuriti",
               "Bangun konteks standar industri sekuriti tinggi internasional (Intergraf & WCO) dan 4 kriteria utama pemesan.")

    # Left Container: Standar Internasional & Multi-Layer Security
    add_card(slide, 0.60, 1.95, 4.40, 4.95, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_section_header(slide, 0.60, 1.95, 4.40, 0.40, "STANDAR INTERNASIONAL & FITUR PENSIAPAN", NAVY)

    tb_left = slide.shapes.add_textbox(Inches(0.75), Inches(2.45), Inches(4.10), Inches(4.35))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_top = Inches(0)
    tf_l.margin_left = Inches(0)

    p = tf_l.paragraphs[0]
    p.text = "Kepatuhan Terhadap Standar Global:"
    p.font.name = 'Arial'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(4)

    p_body = tf_l.add_paragraph()
    p_body.text = "Industri *high-security printing* diatur oleh pedoman ketat seperti standar Intergraf dan World Customs Organization (WCO) untuk pencegahan pemalsuan dokumen negara (*anti-counterfeiting*)."
    p_body.font.name = 'Arial'
    p_body.font.size = Pt(9.5)
    p_body.font.color.rgb = DARK_TEXT
    p_body.space_after = Pt(10)

    p_feat = tf_l.add_paragraph()
    p_feat.text = "Fitur Pengamanan Berlapis (Multi-Layer):"
    p_feat.font.name = 'Arial'
    p_feat.font.size = Pt(10.5)
    p_feat.font.bold = True
    p_feat.font.color.rgb = PURPLE
    p_feat.space_after = Pt(4)

    features = [
        ("Kertas Serat Pengaman:", " Kertas sekuriti higroskopis dengan serat pengaman tak kasat mata."),
        ("Tinta Sekuriti UV:", " Tinta sekuriti berpendar ultra-violet berpresisi tinggi."),
        ("Ornamen Guilloche & Microtext:", " Garis halus anti-fotokopi dan teks mikro rapat."),
        ("Pita Hologram:", " Aplikasi benang pengaman & hologram berakurasi mikron.")
    ]
    for fb, fd in features:
        pf = tf_l.add_paragraph()
        pf.space_after = Pt(5)
        r0 = pf.add_run()
        r0.text = "✔ "
        r0.font.name = 'Arial'
        r0.font.size = Pt(9.5)
        r0.font.bold = True
        r0.font.color.rgb = GREEN
        r1 = pf.add_run()
        r1.text = fb
        r1.font.name = 'Arial'
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = pf.add_run()
        r2.text = fd
        r2.font.name = 'Arial'
        r2.font.size = Pt(9)
        r2.font.color.rgb = DARK_TEXT

    # Right Container: 4 Kriteria Utama Pasar Sekuriti
    add_card(slide, 5.20, 1.95, 7.53, 4.95, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 5.20, 1.95, 7.53, 0.40, "4 KRITERIA UTAMA PASAR PERCETAKAN SEKURITI NEGARA", PURPLE)

    crit_cards = [
        ("1. Jaminan Autentikasi Tanpa Cacat (Zero-Defect)", RED, FILL_LIGHT_RED, BORDER_RED,
         "Cacat fisik seperti blobor (ink bleeding), noda bintik (hickies), dan pergeseran register (misregister) adalah deviasi kritis. Cacat ini dapat menggagalkan verifikasi keaslian aparat di lapangan serta memicu sengketa hukum."),
        ("2. Daya Saing Harga Pengadaan (Cost Competitiveness)", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY,
         "Evaluasi tender menuntut efisiensi anggaran (value for money). Tingginya tingkat kerusakan (inschiet) pada bahan baku mahal akan mendongkrak biaya pokok produksi dan mengurangi daya saing harga perusahaan."),
        ("3. Ketepatan Waktu Pengiriman (Strict Delivery SLA)", ORANGE, FILL_LIGHT_AMBER, BORDER_AMBER,
         "Kepastian jadwal distribusi ke seluruh wilayah sangat ketat. Keterlambatan akibat siklus cetak ulang (tambah cetak) yang panjang akan mengganggu pasokan industri dan menunda penerimaan kas APBN."),
        ("4. Akuntabilitas Bahan Baku (Chain of Custody & Zero Leakage)", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN,
         "Pengawasan total atas seluruh kertas dan tinta sekuriti. Setiap lembar rusak wajib disortir dan dimusnahkan secara resmi dengan berita acara guna mencegah kebocoran dokumen sekuriti ke pasar gelap.")
    ]

    for idx, (c_title, c_col, c_fill, c_bord, c_desc) in enumerate(crit_cards):
        row = idx // 2
        col = idx % 2
        cx = 5.40 + col * 3.65
        cy = 2.45 + row * 2.15

        sub_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(3.45), Inches(2.05))
        sub_card.fill.solid()
        sub_card.fill.fore_color.rgb = c_fill
        sub_card.line.color.rgb = c_bord
        sub_card.line.width = Pt(1.5)

        tb_c = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(cy + 0.08), Inches(3.21), Inches(1.89))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        tf_c.margin_top = Inches(0)
        tf_c.margin_left = Inches(0)

        p1 = tf_c.paragraphs[0]
        p1.text = c_title
        p1.font.name = 'Arial'
        p1.font.size = Pt(9.5)
        p1.font.bold = True
        p1.font.color.rgb = c_col
        p1.space_after = Pt(4)

        p2 = tf_c.add_paragraph()
        p2.text = c_desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = DARK_TEXT

# ==========================================
# SLIDE 4: SKALA TENDER NASIONAL DJBC KEMENKEU
# ==========================================
def build_slide_4(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "01.2", "Skala Tender Nasional: Komitmen Fiskal & Klausul Kontrak DJBC Kemenkeu RI",
               "Jelaskan skala volume ratusan juta lembar, peran fiskal APBN, dan klausul kepatuhan tender Pita Cukai (PCHT & MMEA).")

    # 3 Top KPI Cards
    add_kpi_card(slide, 0.60, 1.95, 3.90, 1.15, "TARGET VOLUME RATA-RATA", "160.000.000 LK", "Perencanaan Kapasitas PPIC Tahunan", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 4.70, 1.95, 3.90, 1.15, "PESANAN AKTUAL 2025", "177.636.930 LK", "Realisasi Modul SAP ZPPRSIPPC0012", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 8.80, 1.95, 3.93, 1.15, "PERAN FISKAL NEGARA", "Ratusan Triliun", "Penerimaan Cukai APBN (PCHT & MMEA)", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)

    # Main Contract Clauses Panel
    add_card(slide, 0.60, 3.25, 12.13, 3.65, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 0.60, 3.25, 12.13, 0.38, "KLAUSUL KUALITAS & KEPATUHAN TENDER PENGADAAN PITA CUKAI NASIONAL", PURPLE)

    clauses = [
        ("1. SPESIFIKASI MUTU MUTLAK", RED, FILL_LIGHT_RED, BORDER_RED, [
            ("Ketepatan Fitur Pengaman:", " Seluruh fitur pengamanan fisik (kertas sekuriti, tinta UV, ornamen guilloche, hologram) wajib tercetak presisi tanpa deviasi warna, register, maupun kepekatan."),
            ("Risiko Pembuktian Keaslian:", " Cacat cetak berisiko memicu kesalahan identifikasi keaslian pita cukai oleh aparat DJBC di lapangan."),
            ("Toleransi Cacat Nol:", " Standar toleransi deviasi ditekan hingga mendekati nol pada setiap lembar dokumen sekuriti negara.")
        ]),
        ("2. REKONSILIASI KETAT LEMBAR RUSAK (HCTS)", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY, [
            ("Akuntabilitas HCTS:", " Setiap lembar cetak rusak dikategorikan sebagai Hasil Cetak Tidak Sempurna (HCTS) dan wajib dipertanggungjawabkan."),
            ("Berita Acara Pemusnahan:", " Wajib melalui rekonsiliasi dan pemusnahan resmi bersama pengawas guna mencegah kebocoran kertas sekuriti."),
            ("Beban Audit Fisik:", " Tingginya tingkat inschiet menambah beban kerja sortir manual dan memperpanjang proses audit fisik.")
        ]),
        ("3. JAMINAN SERVICE LEVEL AGREEMENT (SLA)", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN, [
            ("Ketepatan Distribusi Nasional:", " Pengiriman pesanan bernilai ratusan juta lembar wajib tepat jadwal ke seluruh kantor wilayah Bea Cukai."),
            ("Kelancaran Rantai Pasok:", " Menjamin kepastian pasokan bagi industri hasil tembakau / MMEA dan kepastian arus kas penerimaan APBN."),
            ("Mitigasi Sanksi Penalti:", " Proses tambah cetak yang lambat akibat tingginya lembar rusak berisiko terkena penalti keterlambatan kontrak.")
        ])
    ]

    col_w = 3.82
    for idx, (c_title, c_col, c_fill, c_bord, bullets) in enumerate(clauses):
        cx = 0.80 + idx * 3.95
        sub_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(3.80), Inches(col_w), Inches(2.95))
        sub_card.fill.solid()
        sub_card.fill.fore_color.rgb = c_fill
        sub_card.line.color.rgb = c_bord
        sub_card.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(3.90), Inches(col_w - 0.24), Inches(2.75))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)

        p1 = tf.paragraphs[0]
        p1.text = c_title
        p1.font.name = 'Arial'
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = c_col
        p1.space_after = Pt(6)

        for b_idx, (b_head, b_desc) in enumerate(bullets):
            p = tf.add_paragraph()
            p.space_after = Pt(5)
            r0 = p.add_run()
            r0.text = "• "
            r0.font.name = 'Arial'
            r0.font.size = Pt(9)
            r0.font.bold = True
            r0.font.color.rgb = c_col
            r1 = p.add_run()
            r1.text = b_head
            r1.font.name = 'Arial'
            r1.font.size = Pt(9)
            r1.font.bold = True
            r1.font.color.rgb = DARK_TEXT
            r2 = p.add_run()
            r2.text = b_desc
            r2.font.name = 'Arial'
            r2.font.size = Pt(8.5)
            r2.font.color.rgb = DARK_TEXT

# ==========================================
# SLIDE 5: ARAH STRATEGIS & MANDAT PERURI
# ==========================================
def build_slide_5(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "01.3", "Arah Strategis Korporasi: Mandat PP No. 06/2019 & Keunggulan Operasional",
               "Selaraskan mandat tunggal pencetakan dokumen sekuriti negara dengan 3 sasaran strategis Peruri.")

    # Mandat Banner Top
    top_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(1.95), Inches(12.13), Inches(0.85))
    top_box.fill.solid()
    top_box.fill.fore_color.rgb = NAVY
    top_box.line.fill.background()
    tf_top = top_box.text_frame
    tf_top.word_wrap = True
    tf_top.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_top.margin_left = Inches(0.25)
    tf_top.margin_right = Inches(0.25)
    p_top = tf_top.paragraphs[0]
    r0 = p_top.add_run()
    r0.text = "MANDAT TUNGGAL PERATURAN PEMERINTAH NOMOR 06 TAHUN 2019:  "
    r0.font.name = 'Arial'
    r0.font.size = Pt(11)
    r0.font.bold = True
    r0.font.color.rgb = WHITE
    r1 = p_top.add_run()
    r1.text = "Perum Percetakan Uang Republik Indonesia (Peruri) mengemban amanah mencetak Uang Rupiah dan dokumen sekuriti negara bernilai tinggi dengan jaminan keaslian serta keunggulan operasional berstandar dunia."
    r1.font.name = 'Arial'
    r1.font.size = Pt(10)
    r1.font.color.rgb = WHITE

    # 3 Strategic Pillars
    pillars = [
        ("1. PENGENDALIAN BIAYA BAHAN BAKU", "Cost Leadership & Material Protection", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY, [
            ("Proteksi Bahan Baku Mahal:", " Kertas sekuriti khusus dan tinta berpengaman merupakan komponen biaya manufaktur terbesar."),
            ("Penekanan Rasio Inschiet:", " Mengurangi rasio pemborosan bahan (*inschiet*) guna menjaga efisiensi harga pokok produksi (HPP)."),
            ("Margin Usaha Sehat:", " Menjamin profitabilitas unit kerja dan efisiensi pengadaan pada tender nasional.")
        ]),
        ("2. KEUNGGULAN OPERASIONAL", "Operational Excellence & ISO 9001:2015", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE, [
            ("Standarisasi Mutu Terpadu:", " Menyelaraskan proses pencetakan dengan sistem manajemen mutu internasional ISO 9001:2015."),
            ("Stabilitas Kapasitas Produksi:", " Memastikan 9 mesin cetak beroperasi optimal memenuhi target kontrak tanpa pembengkakan afval."),
            ("Eliminasi Cacat Berulang:", " Menghilangkan potensi cacat blobor, bintik, dan misregister pada lini cetak.")
        ]),
        ("3. DIGITALISASI AREA KERJA", "Smart Factory & INDI 4.0 Readiness", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN, [
            ("Transformasi Buku Folio:", " Mengganti pencatatan manual di meja kontrol mesin menjadi aliran data digital seketika."),
            ("Pengambilan Keputusan Cepat:", " Memberikan data operasional real-time per mesin dan per shift bagi pengawas dan teknisi."),
            ("Dukungan Keputusan Cerdas:", " Fondasi integrasi sistem menuju otomasi industri percetakan sekuriti 4.0.")
        ])
    ]

    col_w = 3.90
    for idx, (p_title, p_sub, p_col, p_fill, p_bord, bullets) in enumerate(pillars):
        cx = 0.60 + idx * 4.10
        add_card(slide, cx, 2.95, col_w, 3.10, p_fill, p_bord)
        add_section_header(slide, cx, 2.95, col_w, 0.40, p_title, p_col)

        tb = slide.shapes.add_textbox(Inches(cx + 0.15), Inches(3.42), Inches(col_w - 0.30), Inches(2.55))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)

        p_subt = tf.paragraphs[0]
        p_subt.text = p_sub
        p_subt.font.name = 'Arial'
        p_subt.font.size = Pt(9)
        p_subt.font.bold = True
        p_subt.font.color.rgb = p_col
        p_subt.space_after = Pt(6)

        for b_idx, (b_head, b_desc) in enumerate(bullets):
            p = tf.add_paragraph()
            p.space_after = Pt(5)
            r0 = p.add_run()
            r0.text = "• "
            r0.font.name = 'Arial'
            r0.font.size = Pt(9)
            r0.font.bold = True
            r0.font.color.rgb = p_col
            r1 = p.add_run()
            r1.text = b_head
            r1.font.name = 'Arial'
            r1.font.size = Pt(9)
            r1.font.bold = True
            r1.font.color.rgb = DARK_TEXT
            r2 = p.add_run()
            r2.text = b_desc
            r2.font.name = 'Arial'
            r2.font.size = Pt(8.5)
            r2.font.color.rgb = DARK_TEXT

    # Bottom Callout Box
    bot_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(6.15), Inches(12.13), Inches(0.75))
    bot_box.fill.solid()
    bot_box.fill.fore_color.rgb = FILL_LIGHT_AMBER
    bot_box.line.color.rgb = BORDER_AMBER
    bot_box.line.width = Pt(1.5)
    tf_bot = bot_box.text_frame
    tf_bot.word_wrap = True
    tf_bot.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_bot.margin_left = Inches(0.20)
    tf_bot.margin_right = Inches(0.20)
    p_bot = tf_bot.paragraphs[0]
    r_b0 = p_bot.add_run()
    r_b0.text = "MUARA STRATEGIS KORPORASI:  "
    r_b0.font.name = 'Arial'
    r_b0.font.size = Pt(10)
    r_b0.font.bold = True
    r_b0.font.color.rgb = ORANGE
    r_b1 = p_bot.add_run()
    r_b1.text = "Seluruh target strategis perusahaan bermuara pada lini operasional dengan volume pekerjaan terbesar: Unit Cetak Pita Cukai. Keberhasilan pengendalian mutu di lini ini secara langsung menentukan efisiensi dan reputasi Perum Peruri."
    r_b1.font.name = 'Arial'
    r_b1.font.size = Pt(9.5)
    r_b1.font.color.rgb = DARK_TEXT

# ==========================================
# SLIDE 6: REALITAS LAPANGAN & PROFIL 9 MESIN
# ==========================================
def build_slide_6(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "02.1", "Realitas Lapangan: Skala Operasional 9 Mesin & Parameter Kapasitas 2025",
               "Paparkan profil teknis unit kerja, pola 3 shift 24/7, komposisi 9 mesin cetak offset, dan parameter baseline.")

    # 4 Top Parameter Cards
    add_kpi_card(slide, 0.60, 1.95, 2.90, 1.15, "9 MESIN CETAK OFFSET", "4 KMR · 2 RYB · 3 GTO", "Komori (1–4), Ryobi (1–2), GTO (1–3)", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 3.65, 1.95, 2.90, 1.15, "POLA KERJA 24/7", "3 Shift Bergilir", "Pagi (07-15), Sore (15-23), Malam (23-07)", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 6.70, 1.95, 2.90, 1.15, "KEKUATAN PERSONEL", "±42 Operator", "Operator Cetak & Kepala Kelompok", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 9.75, 1.95, 2.98, 1.15, "VOLUME AKTUAL 2025", "177.636.930 LK", "Inschiet Baseline: 4,61% (SAP ZPPRSIPPC0012)", RED, FILL_LIGHT_RED, BORDER_RED)

    # Full Data Table (Tabel 1.1)
    table_shape = slide.shapes.add_table(11, 4, Inches(0.60), Inches(3.25), Inches(12.13), Inches(3.65))
    table = table_shape.table
    table.columns[0].width = Inches(4.30)
    table.columns[1].width = Inches(2.60)
    table.columns[2].width = Inches(1.80)
    table.columns[3].width = Inches(3.43)

    table_data = [
        ("Parameter Operasional / Periode", "Nilai / Angka", "Satuan", "Sumber Data Terverifikasi"),
        ("Jumlah Mesin Cetak Aktif", "9 Mesin (4 Komori, 2 Ryobi, 3 GTO)", "Unit Mesin", "Inventaris Aset Dept. Khazanah & Verifikasi"),
        ("Pola Gilir Kerja (Shift)", "3 (Pagi, Sore, Malam)", "Shift / Hari", "Standar Pola Penugasan Gilir Unit Cetak"),
        ("Durasi Operasional Lini", "24", "Jam / Hari", "Standard Operating Procedure (SOP) Unit Cetak"),
        ("Total Personel Operator Cetak", "±42", "Personel", "Data Penugasan Kerja Seksi Cetak"),
        ("Target Volume Rata-Rata Tahunan", "160.000.000", "Lembar Cetak", "Perencanaan Kapasitas PPIC Peruri"),
        ("Total Volume Pesanan Aktual 2025", "177.636.930", "Lembar Cetak", "Modul SAP Production Order (ZPPRSIPPC0012)"),
        ("Inschiet Kuartal 1 (Q1 2025)", "4,72%", "%", "Rekap Verifikasi Mutu & Modul SAP"),
        ("Inschiet Kuartal 2 (Q2 2025)", "3,97%", "%", "Rekap Verifikasi Mutu & Modul SAP"),
        ("Inschiet Kuartal 4 (Q4 2025)", "5,11% (Puncak Lonjakan)", "%", "Rekap Verifikasi Mutu & Modul SAP"),
        ("RATA-RATA BASELINE INSCHIET 2025", "4,61%", "%", "Konsolidasi Tahunan Verifikasi Mutu & SAP")
    ]

    for row_idx, row in enumerate(table_data):
        for col_idx, text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.name = 'Arial'

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                p.font.size = Pt(9.5)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER if col_idx in [1, 2] else PP_ALIGN.LEFT
            elif row_idx == 10:
                cell.fill.solid()
                cell.fill.fore_color.rgb = FILL_LIGHT_RED
                p.font.size = Pt(9.5)
                p.font.bold = True
                p.font.color.rgb = RED
                p.alignment = PP_ALIGN.CENTER if col_idx in [1, 2] else PP_ALIGN.LEFT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 1 else FILL_LIGHT_NAVY
                p.font.size = Pt(8.5)
                p.font.color.rgb = DARK_TEXT
                if col_idx in [1, 2]:
                    p.alignment = PP_ALIGN.CENTER
                    p.font.bold = (col_idx == 1)
                else:
                    p.alignment = PP_ALIGN.LEFT

# ==========================================
# SLIDE 7: ANALISIS FLUKTUASI BASELINE 2025
# ==========================================
def build_slide_7(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "02.2", "Fluktuasi Baseline Inschiet 2025: Bukti Kapabilitas vs Lonjakan Q4",
               "Analisis dinamika tren inschiet kuartalan 2025 dan pembuktian kapabilitas teknis mesin vs lonjakan order desain baru.")

    # Left Container: Native Column Chart
    add_card(slide, 0.60, 1.95, 6.00, 4.95, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 0.60, 1.95, 6.00, 0.38, "TREN INSCHIET BASELINE PER KUARTAL 2025 (%)", PURPLE)

    # Add Native Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1 (Jan-Mar)', 'Q2 (Apr-Jun)', 'Q3 (Jul-Sep)', 'Q4 (Okt-Des)', 'Rata-rata 2025']
    chart_data.add_series('Inschiet (%)', (4.72, 3.97, 4.64, 5.11, 4.61))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.80), Inches(2.45), Inches(5.60), Inches(3.60),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    plots = chart.plots[0]
    plots.has_data_labels = True
    for series in plots.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = PURPLE
        dl = series.data_labels
        dl.font.name = 'Arial'
        dl.font.size = Pt(11)
        dl.font.bold = True
        dl.font.color.rgb = DARK_NAVY
        dl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    # Chart Note Bottom
    tb_cn = slide.shapes.add_textbox(Inches(0.80), Inches(6.15), Inches(5.60), Inches(0.65))
    tf_cn = tb_cn.text_frame
    tf_cn.word_wrap = True
    p_cn = tf_cn.paragraphs[0]
    p_cn.text = "Sumber Data: Modul SAP Production Order (ZPPRSIPPC0012) & Unit Verifikasi Mutu Pita Cukai."
    p_cn.font.name = 'Arial'
    p_cn.font.size = Pt(8.5)
    p_cn.font.italic = True
    p_cn.font.color.rgb = MUTED_TEXT

    # Right Container 1: Pembuktian Kapabilitas Q2
    add_card(slide, 6.80, 1.95, 5.93, 2.35, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_section_header(slide, 6.80, 1.95, 5.93, 0.35, "PEMBUKTIAN KAPABILITAS MESIN (Q2 2025: 3,97%)", GREEN)

    tb_r1 = slide.shapes.add_textbox(Inches(6.95), Inches(2.38), Inches(5.63), Inches(1.82))
    tf_r1 = tb_r1.text_frame
    tf_r1.word_wrap = True
    tf_r1.margin_top = Inches(0)
    tf_r1.margin_left = Inches(0)

    p1 = tf_r1.paragraphs[0]
    r0 = p1.add_run()
    r0.text = "• Potensi Teknis di Bawah 4,00%: "
    r0.font.name = 'Arial'
    r0.font.size = Pt(9.5)
    r0.font.bold = True
    r0.font.color.rgb = GREEN
    r1 = p1.add_run()
    r1.text = "Pada Q2 2025, inschiet berhasil ditekan ke angka 3,97%. Capaian ini membuktikan bahwa 9 mesin cetak dan operator secara teknis mampu beroperasi stabil memenuhi toleransi mutu standar."
    r1.font.name = 'Arial'
    r1.font.size = Pt(9)
    r1.font.color.rgb = DARK_TEXT
    p1.space_after = Pt(4)

    p2 = tf_r1.add_paragraph()
    r2 = p2.add_run()
    r2.text = "• Ketergantungan Kondisi Ideal: "
    r2.font.name = 'Arial'
    r2.font.size = Pt(9.5)
    r2.font.bold = True
    r2.font.color.rgb = DARK_TEXT
    r3 = p2.add_run()
    r3.text = "Stabilitas ini hanya tercapai saat pesanan berulang dan parameter mesin tidak berubah. Belum ada sistem pengendali saat variasi order meningkat."
    r3.font.name = 'Arial'
    r3.font.size = Pt(9)
    r3.font.color.rgb = DARK_TEXT

    # Right Container 2: Anomali Lonjakan Q4
    add_card(slide, 6.80, 4.50, 5.93, 2.40, FILL_LIGHT_RED, BORDER_RED)
    add_section_header(slide, 6.80, 4.50, 5.93, 0.35, "LONJAKAN DESAIN BARU & TITIK BUTA DATA (Q4 2025: 5,11%)", RED)

    tb_r2 = slide.shapes.add_textbox(Inches(6.95), Inches(4.93), Inches(5.63), Inches(1.87))
    tf_r2 = tb_r2.text_frame
    tf_r2.word_wrap = True
    tf_r2.margin_top = Inches(0)
    tf_r2.margin_left = Inches(0)

    p3 = tf_r2.paragraphs[0]
    r4 = p3.add_run()
    r4.text = "• Lonjakan Tajam Menjelang Akhir Tahun: "
    r4.font.name = 'Arial'
    r4.font.size = Pt(9.5)
    r4.font.bold = True
    r4.font.color.rgb = RED
    r5 = p3.add_run()
    r5.text = "Inschiet melonjak ke level puncaknya 5,11% (+1,14 pp vs Q2) saat unit menerima pesanan pita cukai berdesain baru dalam jumlah besar."
    r5.font.name = 'Arial'
    r5.font.size = Pt(9)
    r5.font.color.rgb = DARK_TEXT
    p3.space_after = Pt(4)

    p4 = tf_r2.add_paragraph()
    r6 = p4.add_run()
    r6.text = "• Ketiadaan Diagnostik Meja Mesin: "
    r6.font.name = 'Arial'
    r6.font.size = Pt(9.5)
    r6.font.bold = True
    r6.font.color.rgb = DARK_TEXT
    r7 = p4.add_run()
    r7.text = "Penyetelan awal (make-ready) memakan waktu lama dan deviasi mutu terlambat terdeteksi, sehingga lembar rusak membengkak seiring naiknya kecepatan cetak."
    r7.font.name = 'Arial'
    r7.font.size = Pt(9)
    r7.font.color.rgb = DARK_TEXT

# ==========================================
# SLIDE 8: KESENJANGAN OPERASIONAL (DATA SILO)
# ==========================================
def build_slide_8(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "02.3", "Kesenjangan Operasional: Pemisahan Aliran Data Meja Mesin vs SAP",
               "Tunjukkan anatomi pemisahan data (data silo) antara pencatatan meja mesin dan laporan kantor.")

    # Two Split Panels
    col_w = 5.85
    # Left: Meja Mesin
    add_card(slide, 0.60, 1.95, col_w, 3.65, FILL_LIGHT_AMBER, BORDER_AMBER)
    add_section_header(slide, 0.60, 1.95, col_w, 0.40, "PENCATATAN TRANSAKSI DI MEJA MESIN CETAK", ORANGE)

    tb_l = slide.shapes.add_textbox(Inches(0.80), Inches(2.45), Inches(col_w - 0.40), Inches(3.05))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_top = Inches(0)
    tf_l.margin_left = Inches(0)

    bullets_l = [
        ("Buku Folio Fisik Manual:", " Dicatat manual dengan pena pada buku folio fisik di meja kontrol 9 mesin cetak (KMR 1–4, RYB 1–2, GTO 1–3)."),
        ("Data Terisolasi di Meja Mesin:", " Catatan jam jalan mesin, jumlah cetak, dan nama regu hanya tersimpan di area mesin dan tidak terhubung ke jaringan."),
        ("Rekapitulasi Manual Berkala:", " Rekap baru dilakukan oleh Kepala Kelompok saat evaluasi triwulan atau akhir masa kontrak kerja pegawai."),
        ("Proses Lambat & Rawan Human Error:", " Rentan terhadap kesalahan hitung manual, keterlambatan pelaporan, dan kehilangan buku fisik.")
    ]
    for idx, (bh, bd) in enumerate(bullets_l):
        p = tf_l.paragraphs[0] if idx == 0 else tf_l.add_paragraph()
        p.space_after = Pt(6)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(9.5)
        r0.font.bold = True
        r0.font.color.rgb = ORANGE
        r1 = p.add_run()
        r1.text = bh
        r1.font.name = 'Arial'
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = bd
        r2.font.name = 'Arial'
        r2.font.size = Pt(9)
        r2.font.color.rgb = DARK_TEXT

    # Right: SAP & Verifikasi
    add_card(slide, 6.88, 1.95, col_w, 3.65, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 6.88, 1.95, col_w, 0.40, "HASIL SORTIR DI VERIFIKASI & MODUL SAP KANTOR", PURPLE)

    tb_r = slide.shapes.add_textbox(Inches(7.08), Inches(2.45), Inches(col_w - 0.40), Inches(3.05))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_top = Inches(0)
    tf_r.margin_left = Inches(0)

    bullets_r = [
        ("Waktu Tunggu Sortir 1–2 Hari:", " Lembar cetak diperiksa di Unit Verifikasi Mutu dengan jeda waktu pemeriksaan fisik 1 s.d. 2 hari pasca-cetak."),
        ("Ringkasan Kerusakan Global:", " Data diinput ke SAP ZPPRSIPPC0012 sebagai angka kerusakan total di tingkat unit (unit-wide summary)."),
        ("Data Pasif (Dormant Data):", " Tabel CSV mentah SAP tersimpan di komputer kantor dan tidak dapat diakses langsung oleh operator di lapangan."),
        ("Ketiadaan Atribusi Spesifik:", " Laporan TIDAK merekam nomor mesin cetak, nomor PO, maupun kelompok gilir kerja (shift) yang mencetaknya.")
    ]
    for idx, (bh, bd) in enumerate(bullets_r):
        p = tf_r.paragraphs[0] if idx == 0 else tf_r.add_paragraph()
        p.space_after = Pt(6)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(9.5)
        r0.font.bold = True
        r0.font.color.rgb = PURPLE
        r1 = p.add_run()
        r1.text = bh
        r1.font.name = 'Arial'
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = bd
        r2.font.name = 'Arial'
        r2.font.size = Pt(9)
        r2.font.color.rgb = DARK_TEXT

    # Bottom Full Width Impact Box
    bot_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(5.75), Inches(12.13), Inches(1.15))
    bot_card.fill.solid()
    bot_card.fill.fore_color.rgb = FILL_LIGHT_RED
    bot_card.line.color.rgb = BORDER_RED
    bot_card.line.width = Pt(1.5)

    tb_bot = slide.shapes.add_textbox(Inches(0.80), Inches(5.82), Inches(11.73), Inches(1.00))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    tf_bot.margin_top = Inches(0)
    tf_bot.margin_left = Inches(0)

    p_b1 = tf_bot.paragraphs[0]
    r_b0 = p_b1.add_run()
    r_b0.text = "AKIBAT PEMISAHAN ALIRAN DATA (DATA SILO): KEBUTAAN ATRIBUSI (ATTRIBUTION BLINDNESS)\n"
    r_b0.font.name = 'Arial'
    r_b0.font.size = Pt(10)
    r_b0.font.bold = True
    r_b0.font.color.rgb = RED

    r_b1 = p_b1.add_run()
    r_b1.text = "Modul SAP hanya menyajikan 'apa jenis kerusakannya' di level unit (blobor, noda bintik, misregister), tetapi TIDAK BISA MENJAWAB 'di mesin mana', 'pada shift berapa', dan 'siapa operatornya'. Manajemen dan teknisi kehilangan dasar data untuk melakukan tindakan korektif cepat di area mesin."
    r_b1.font.name = 'Arial'
    r_b1.font.size = Pt(9)
    r_b1.font.color.rgb = DARK_TEXT

# ==========================================
# SLIDE 9: IMPLIKASI TITIK BUTA DATA
# ==========================================
def build_slide_9(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "02.4", "Implikasi Titik Buta Data: Kendala Penelusuran Mesin vs Operator",
               "Jelaskan 3 kendala mendasar lapangan: pemeriksaan spekulatif > 8 jam, dilema mesin vs shift, dan evaluasi tertunda.")

    # 3 Column Cards
    col_w = 3.90
    cards = [
        ("1. PEMERIKSAAN SPEKULATIF (> 8 JAM)", RED, FILL_LIGHT_RED, BORDER_RED, [
            ("Pencarian Mesin Trial-and-Error:", " Saat bagian verifikasi melaporkan kenaikan cacat blobor atau bintik, teknisi tidak tahu mesin mana yang bermasalah."),
            ("Pemeriksaan Seluruh 9 Mesin:", " Teknisi harus memeriksa 9 mesin satu per satu secara coba-coba tanpa skala prioritas data."),
            ("Jam Henti Produktif Melonjak:", " Waktu henti mesin (*downtime*) melampaui > 1 shift (> 8 jam per mesin) dan mesin terus mencetak lembar rusak.")
        ]),
        ("2. DILEMA MESIN VS CARA KERJA OPERATOR", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY, [
            ("Penurunan Komponen Fisik Mesin:", " Penurunan performa komponen seiring umur mesin: rol karet mengeras/licin (*glazing*), selimut karet (*blanket*) kempes/turun elastisitas, penjepit silinder melemah."),
            ("vs Variasi Penyetelan & Shift Malam:", " Variasi make-ready antar operator dan kelelahan ritme sirkadian pada Shift Malam (23.00–07.00 WIB)."),
            ("Akar Masalah Tidak Terurai:", " Tanpa data granular, manajemen tidak dapat memastikan apakah cacat dipicu mesin atau operator.")
        ]),
        ("3. EVALUASI KINERJA OPERATOR TERTUNDA", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE, [
            ("Umpan Balik Harian Terputus:", " Kepala Unit dan Kepala Kelompok tidak dapat memberikan bimbingan teknis harian kepada ±42 operator."),
            ("Rekapitulasi Manual Terlambat:", " Rekam jejak kerja baru diketahui berbulan-bulan kemudian saat buku folio direkap manual."),
            ("Friksi Kerja Antar-Gilir:", " Menimbulkan prasangka dan saling menyalahkan antar-shift saat terjadi lonjakan lembar rusak.")
        ])
    ]

    for idx, (c_title, c_col, c_fill, c_bord, bullets) in enumerate(cards):
        cx = 0.60 + idx * 4.10
        add_card(slide, cx, 1.95, col_w, 4.95, c_fill, c_bord)
        add_section_header(slide, cx, 1.95, col_w, 0.40, c_title, c_col)

        tb = slide.shapes.add_textbox(Inches(cx + 0.15), Inches(2.45), Inches(col_w - 0.30), Inches(4.35))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)

        for b_idx, (b_head, b_desc) in enumerate(bullets):
            p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            r0 = p.add_run()
            r0.text = "• "
            r0.font.name = 'Arial'
            r0.font.size = Pt(9.5)
            r0.font.bold = True
            r0.font.color.rgb = c_col
            r1 = p.add_run()
            r1.text = b_head
            r1.font.name = 'Arial'
            r1.font.size = Pt(9.5)
            r1.font.bold = True
            r1.font.color.rgb = DARK_TEXT
            r2 = p.add_run()
            r2.text = b_desc
            r2.font.name = 'Arial'
            r2.font.size = Pt(9)
            r2.font.color.rgb = DARK_TEXT

# ==========================================
# SLIDE 10: SIMULASI FINANSIAL BASELINE 2025
# ==========================================
def build_slide_10(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "03.1", "Skala Dampak Finansial: Beban Baseline 2025 & Valuasi 1% Inschiet",
               "Sajikan simulasi finansial transparan (Skenario A & B) dan nilai efisiensi per 1,00% penurunan inschiet.")

    # Top Metric Cards
    add_kpi_card(slide, 0.60, 1.95, 3.90, 1.15, "ESTIMASI BIAYA CETAK*", "Rp 3.000 / LK", "Kertas Sekuriti, Tinta UV & Biaya Mesin", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 4.70, 1.95, 3.90, 1.15, "BEBAN BASELINE TAHUNAN", "Rp 22,13 – 24,56 M", "Rp 1,84 – 2,05 Miliar / Bulan", RED, FILL_LIGHT_RED, BORDER_RED)
    add_kpi_card(slide, 8.80, 1.95, 3.93, 1.15, "VALUASI TIAP 1% PENURUNAN", "Rp 4,80 – 5,33 M", "Penyelamatan 1,60M – 1,78M Lembar/Thn", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)

    # 3 Detailed Calculation Panels
    col_w = 3.90
    # Skenario A
    add_card(slide, 0.60, 3.25, col_w, 2.95, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_section_header(slide, 0.60, 3.25, col_w, 0.38, "SKENARIO A: STANDAR VOLUME TAHUNAN", NAVY)

    tb_a = slide.shapes.add_textbox(Inches(0.75), Inches(3.72), Inches(col_w - 0.30), Inches(2.40))
    tf_a = tb_a.text_frame
    tf_a.word_wrap = True
    tf_a.margin_top = Inches(0)
    tf_a.margin_left = Inches(0)

    p_a1 = tf_a.paragraphs[0]
    p_a1.text = "• Volume Target Tahunan:"
    p_a1.font.name = 'Arial'
    p_a1.font.size = Pt(9)
    p_a1.font.bold = True
    p_a1.font.color.rgb = DARK_NAVY
    p_a1.space_after = Pt(1)

    p_a1_val = tf_a.add_paragraph()
    p_a1_val.text = "  160.000.000 Lembar Cetak (PPIC)"
    p_a1_val.font.name = 'Arial'
    p_a1_val.font.size = Pt(9)
    p_a1_val.font.color.rgb = DARK_TEXT
    p_a1_val.space_after = Pt(4)

    p_a2 = tf_a.add_paragraph()
    p_a2.text = "• Estimasi Lembar Rusak (Baseline 4,61%):"
    p_a2.font.name = 'Arial'
    p_a2.font.size = Pt(9)
    p_a2.font.bold = True
    p_a2.font.color.rgb = DARK_NAVY
    p_a2.space_after = Pt(1)

    p_a2_val = tf_a.add_paragraph()
    p_a2_val.text = "  160.000.000 × 4,61% = 7.376.000 LK / Thn"
    p_a2_val.font.name = 'Arial'
    p_a2_val.font.size = Pt(9)
    p_a2_val.font.color.rgb = DARK_TEXT
    p_a2_val.space_after = Pt(4)

    p_a3 = tf_a.add_paragraph()
    p_a3.text = "• Nilai Beban Kerugian Finansial:"
    p_a3.font.name = 'Arial'
    p_a3.font.size = Pt(9)
    p_a3.font.bold = True
    p_a3.font.color.rgb = RED
    p_a3.space_after = Pt(1)

    p_a3_val = tf_a.add_paragraph()
    p_a3_val.text = "  7.376.000 LK × Rp 3.000 = Rp 22,13 Miliar / Thn"
    p_a3_val.font.name = 'Arial'
    p_a3_val.font.size = Pt(10)
    p_a3_val.font.bold = True
    p_a3_val.font.color.rgb = RED
    p_a3_val.space_after = Pt(2)

    p_a3_sub = tf_a.add_paragraph()
    p_a3_sub.text = "  (Setara ± Rp 1,84 Miliar / Bulan)"
    p_a3_sub.font.name = 'Arial'
    p_a3_sub.font.size = Pt(8.5)
    p_a3_sub.font.bold = True
    p_a3_sub.font.color.rgb = DARK_NAVY

    # Skenario B
    add_card(slide, 4.70, 3.25, col_w, 2.95, FILL_LIGHT_RED, BORDER_RED)
    add_section_header(slide, 4.70, 3.25, col_w, 0.38, "SKENARIO B: REALISASI PESANAN 2025", RED)

    tb_b = slide.shapes.add_textbox(Inches(4.85), Inches(3.72), Inches(col_w - 0.30), Inches(2.40))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    tf_b.margin_top = Inches(0)
    tf_b.margin_left = Inches(0)

    p_b1 = tf_b.paragraphs[0]
    p_b1.text = "• Total Volume Aktual 2025:"
    p_b1.font.name = 'Arial'
    p_b1.font.size = Pt(9)
    p_b1.font.bold = True
    p_b1.font.color.rgb = RED
    p_b1.space_after = Pt(1)

    p_b1_val = tf_b.add_paragraph()
    p_b1_val.text = "  177.636.930 Lembar (SAP ZPPRSIPPC0012)"
    p_b1_val.font.name = 'Arial'
    p_b1_val.font.size = Pt(9)
    p_b1_val.font.color.rgb = DARK_TEXT
    p_b1_val.space_after = Pt(4)

    p_b2 = tf_b.add_paragraph()
    p_b2.text = "• Lembar Rusak Aktual Baseline (4,61%):"
    p_b2.font.name = 'Arial'
    p_b2.font.size = Pt(9)
    p_b2.font.bold = True
    p_b2.font.color.rgb = RED
    p_b2.space_after = Pt(1)

    p_b2_val = tf_b.add_paragraph()
    p_b2_val.text = "  177.636.930 × 4,61% = 8.189.062 LK / Thn"
    p_b2_val.font.name = 'Arial'
    p_b2_val.font.size = Pt(9)
    p_b2_val.font.color.rgb = DARK_TEXT
    p_b2_val.space_after = Pt(4)

    p_b3 = tf_b.add_paragraph()
    p_b3.text = "• Nilai Kerugian Aktual Baseline 2025:"
    p_b3.font.name = 'Arial'
    p_b3.font.size = Pt(9)
    p_b3.font.bold = True
    p_b3.font.color.rgb = RED
    p_b3.space_after = Pt(1)

    p_b3_val = tf_b.add_paragraph()
    p_b3_val.text = "  8.189.062 LK × Rp 3.000 = Rp 24,56 Miliar / Thn"
    p_b3_val.font.name = 'Arial'
    p_b3_val.font.size = Pt(10)
    p_b3_val.font.bold = True
    p_b3_val.font.color.rgb = RED
    p_b3_val.space_after = Pt(2)

    p_b3_sub = tf_b.add_paragraph()
    p_b3_sub.text = "  (Setara ± Rp 2,05 Miliar / Bulan)"
    p_b3_sub.font.name = 'Arial'
    p_b3_sub.font.size = Pt(8.5)
    p_b3_sub.font.bold = True
    p_b3_sub.font.color.rgb = DARK_NAVY

    # Valuasi 1% Inschiet
    add_card(slide, 8.80, 3.25, col_w + 0.03, 2.95, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_section_header(slide, 8.80, 3.25, col_w + 0.03, 0.38, "VALUASI PENURUNAN TIAP 1,00% INSCHIET", GREEN)

    tb_v = slide.shapes.add_textbox(Inches(8.95), Inches(3.72), Inches(col_w - 0.27), Inches(2.40))
    tf_v = tb_v.text_frame
    tf_v.word_wrap = True
    tf_v.margin_top = Inches(0)
    tf_v.margin_left = Inches(0)

    p_v0 = tf_v.paragraphs[0]
    p_v0.text = "Sensitivitas Efisiensi Tiap 100 bps (1,00%):"
    p_v0.font.name = 'Arial'
    p_v0.font.size = Pt(9)
    p_v0.font.bold = True
    p_v0.font.color.rgb = DARK_GREEN
    p_v0.space_after = Pt(4)

    p_v1 = tf_v.add_paragraph()
    p_v1.text = "• Pada Standar Volume (160 Juta LK):"
    p_v1.font.name = 'Arial'
    p_v1.font.size = Pt(8.5)
    p_v1.font.bold = True
    p_v1.font.color.rgb = DARK_TEXT
    p_v1.space_after = Pt(1)

    p_v1_val = tf_v.add_paragraph()
    p_v1_val.text = "  Penyelamatan: 1.600.000 Lembar Kertas\n  Efisiensi: Rp 4,80 Miliar / Tahun"
    p_v1_val.font.name = 'Arial'
    p_v1_val.font.size = Pt(9)
    p_v1_val.font.bold = True
    p_v1_val.font.color.rgb = GREEN
    p_v1_val.space_after = Pt(4)

    p_v2 = tf_v.add_paragraph()
    p_v2.text = "• Pada Volume Aktual 2025 (177,6 Juta LK):"
    p_v2.font.name = 'Arial'
    p_v2.font.size = Pt(8.5)
    p_v2.font.bold = True
    p_v2.font.color.rgb = DARK_TEXT
    p_v2.space_after = Pt(1)

    p_v2_val = tf_v.add_paragraph()
    p_v2_val.text = "  Penyelamatan: 1.776.369 Lembar Kertas\n  Efisiensi: Rp 5,33 Miliar / Tahun"
    p_v2_val.font.name = 'Arial'
    p_v2_val.font.size = Pt(9)
    p_v2_val.font.bold = True
    p_v2_val.font.color.rgb = GREEN

    # Bottom Confidentiality Note Box
    note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(6.30), Inches(12.13), Inches(0.60))
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = FILL_LIGHT_NAVY
    note_box.line.color.rgb = BORDER_NAVY
    note_box.line.width = Pt(1.0)
    tf_n = note_box.text_frame
    tf_n.word_wrap = True
    tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_n.margin_left = Inches(0.15)
    tf_n.margin_right = Inches(0.15)
    p_n = tf_n.paragraphs[0]
    p_n.text = "*Catatan Asumsi Finansial: Nilai estimasi biaya cetak Rp 3.000/lembar digunakan khusus untuk simulasi internal (cost avoidance) dan bukan merupakan rincian harga jual resmi atau HPP rahasia Perum Peruri."
    p_n.font.name = 'Arial'
    p_n.font.size = Pt(8.5)
    p_n.font.italic = True
    p_n.font.color.rgb = MUTED_TEXT

# ==========================================
# SLIDE 11: MATRIKS EVALUASI 5 PILAR COST OF INACTION
# ==========================================
def build_slide_11(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "03.2", "Evaluasi Risiko Pembiaran: Matriks 5 Pilar Cost of Inaction",
               "Paparkan ancaman multidimensi (Biaya, Mutu, Kepatuhan, K3L/ESG, Layanan SLA) bila data silo dibiarkan.")

    # 5 Structured Cards Grid (Tabel 1.2)
    col_w = 2.30
    gap = 0.15
    pillars = [
        ("1. BIAYA (COST)", "KRITIS", RED, FILL_LIGHT_RED, BORDER_RED,
         "Akumulasi pemborosan biaya bahan baku kertas sekuriti dan tinta khusus mencapai Rp 22,13 – Rp 24,56 Miliar/tahun.",
         "Pembengkakan biaya tambah cetak & penurunan margin laba unit."),
        ("2. MUTU (QUALITY)", "TINGGI", RED, FILL_LIGHT_RED, BORDER_RED,
         "Tingkat inschiet berfluktuasi hingga 5,11% akibat penanganan suku cadang mesin yang terlambat dan coba-coba.",
         "Tingginya persentase cacat mutu HCTS di unit kerja."),
        ("3. KEPATUHAN", "TINGGI", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE,
         "Pencatatan manual di buku folio menyulitkan penelusuran riwayat lot produksi saat audit mutu ISO 9001:2015.",
         "Potensi temuan audit & hilangnya rekam jejak digital per PO."),
        ("4. K3L / ESG", "SEDANG", ORANGE, FILL_LIGHT_AMBER, BORDER_AMBER,
         "Timbulan limbah lembar rusak mencapai 7,37 – 8,18 Juta LK/tahun (±60–65 Ton kertas) & kelelahan shift malam.",
         "Pemborosan sumber daya kertas & beban fisik operator."),
        ("5. LAYANAN (SLA)", "TINGGI", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY,
         "Siklus tambah cetak lambat memperlambat serah terima pita cukai ke DJBC dan mengganggu pasokan industri.",
         "Ancaman denda penalti SLA & risiko komplain pemesan.")
    ]

    for idx, (p_title, p_sev, p_col, p_fill, p_bord, p_risk, p_impact) in enumerate(pillars):
        cx = 0.60 + idx * (col_w + gap)
        add_card(slide, cx, 1.95, col_w, 4.95, p_fill, p_bord)

        # Header box
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(1.95), Inches(col_w), Inches(0.60))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = p_col
        hdr.line.fill.background()
        tf_h = hdr.text_frame
        tf_h.word_wrap = True
        tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_h.margin_left = Inches(0.08)
        tf_h.margin_right = Inches(0.08)
        p1 = tf_h.paragraphs[0]
        p1.text = p_title
        p1.font.name = 'Arial'
        p1.font.size = Pt(9.5)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p1.alignment = PP_ALIGN.CENTER

        # Severity Badge
        sev_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx + 0.35), Inches(2.65), Inches(col_w - 0.70), Inches(0.35))
        sev_box.fill.solid()
        sev_box.fill.fore_color.rgb = WHITE
        sev_box.line.color.rgb = p_col
        sev_box.line.width = Pt(1.5)
        tf_s = sev_box.text_frame
        tf_s.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_s = tf_s.paragraphs[0]
        p_s.text = f"TINGKAT: {p_sev}"
        p_s.font.name = 'Arial'
        p_s.font.size = Pt(8.5)
        p_s.font.bold = True
        p_s.font.color.rgb = p_col
        p_s.alignment = PP_ALIGN.CENTER

        # Content Box
        tb_c = slide.shapes.add_textbox(Inches(cx + 0.10), Inches(3.10), Inches(col_w - 0.20), Inches(3.70))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        tf_c.margin_top = Inches(0)
        tf_c.margin_left = Inches(0)

        p_r_lbl = tf_c.paragraphs[0]
        p_r_lbl.text = "Bentuk Risiko Lapangan:"
        p_r_lbl.font.name = 'Arial'
        p_r_lbl.font.size = Pt(9)
        p_r_lbl.font.bold = True
        p_r_lbl.font.color.rgb = p_col
        p_r_lbl.space_after = Pt(2)

        p_r_txt = tf_c.add_paragraph()
        p_r_txt.text = p_risk
        p_r_txt.font.name = 'Arial'
        p_r_txt.font.size = Pt(8.5)
        p_r_txt.font.color.rgb = DARK_TEXT
        p_r_txt.space_after = Pt(10)

        p_i_lbl = tf_c.add_paragraph()
        p_i_lbl.text = "Indikator Dampak Terukur:"
        p_i_lbl.font.name = 'Arial'
        p_i_lbl.font.size = Pt(9)
        p_i_lbl.font.bold = True
        p_i_lbl.font.color.rgb = DARK_NAVY
        p_i_lbl.space_after = Pt(2)

        p_i_txt = tf_c.add_paragraph()
        p_i_txt.text = p_impact
        p_i_txt.font.name = 'Arial'
        p_i_txt.font.size = Pt(8.5)
        p_i_txt.font.color.rgb = DARK_TEXT

# ==========================================
# SLIDE 12: SOLUSI INTEGRASI DATA DSS SIRINE 4.0
# ==========================================
def build_slide_12(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "04.1", "Solusi Terintegrasi: Arsitektur 3 Titik Aliran Data DSS SIRINE 4.0",
               "Tunjukkan konsep solusi sistemik yang menghubungkan transaksi meja mesin, modul SAP, dan verifikasi mutu.")

    # Top Flow Card (3 Integrated Data Points)
    add_card(slide, 0.60, 1.95, 12.13, 1.65, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 0.60, 1.95, 12.13, 0.35, "INTEGRASI 3 TITIK ALIRAN DATA OPERASIONAL SEKETIKA", PURPLE)

    # 3 Flow Blocks
    f_blocks = [
        ("DATA TRANSAKSI MEJA MESIN", "Input digital < 30 detik per PO di 9 meja mesin kontrol (menggantikan buku folio).", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE),
        ("MODUL SAP PRODUCTION ORDER", "Modul SAP ZPPRSIPPC0012 terhubung otomatis dengan spesifikasi & target lembar PO.", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY),
        ("DATA SORTIR VERIFIKASI MUTU", "Hasil sortir HCTS terpetakan seketika ke nomor mesin, shift, dan kategori cacat cetak.", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    ]
    for idx, (fb_title, fb_desc, fb_col, fb_fill, fb_bord) in enumerate(f_blocks):
        cx = 0.80 + idx * 3.95
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(2.40), Inches(3.45), Inches(1.05))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = fb_col
        box.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(Inches(cx + 0.10), Inches(2.45), Inches(3.25), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        p = tf.paragraphs[0]
        p.text = fb_title
        p.font.name = 'Arial'
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = fb_col
        p.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = fb_desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = DARK_TEXT

        # Arrow connector
        if idx < 2:
            arr = slide.shapes.add_textbox(Inches(cx + 3.45), Inches(2.65), Inches(0.50), Inches(0.50))
            tf_a = arr.text_frame
            p_a = tf_a.paragraphs[0]
            p_a.text = "⇄"
            p_a.font.name = 'Arial'
            p_a.font.size = Pt(18)
            p_a.font.bold = True
            p_a.font.color.rgb = PURPLE
            p_a.alignment = PP_ALIGN.CENTER

    # Middle Granular Chain Banner
    mid_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(3.75), Inches(12.13), Inches(0.65))
    mid_card.fill.solid()
    mid_card.fill.fore_color.rgb = DARK_NAVY
    mid_card.line.fill.background()
    tf_mid = mid_card.text_frame
    tf_mid.word_wrap = True
    tf_mid.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_mid = tf_mid.paragraphs[0]
    p_mid.text = "RUNTUTAN DATA GRANULAR:   Nomor PO  ➔  Nomor Mesin (9 Mesin)  ➔  Pola Gilir (Shift 1/2/3)  ➔  Tim Operator  ➔  Kategori Cacat"
    p_mid.font.name = 'Arial'
    p_mid.font.size = Pt(11)
    p_mid.font.bold = True
    p_mid.font.color.rgb = WHITE
    p_mid.alignment = PP_ALIGN.CENTER

    # 3 Transformational Feature Columns Bottom
    feats = [
        ("1. ELIMINASI BUKU FOLIO FISIK", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE, [
            ("Form Digital < 30 Detik:", " Entry PO cepat langsung di meja kontrol mesin dengan fitur auto-complete."),
            ("Validasi Data Real-Time:", " Mencegah salah hitung dan memastikan data jam mesin terekam seketika."),
            ("Integrasi Jadwal Operator:", " Template penugasan gilir operator otomatis terhubung ke kartu pesanan.")
        ]),
        ("2. ATRIBUSI DATA MUTU PRESISI", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY, [
            ("Pelacakan Mesin & Shift:", " Menghubungkan temuan cacat verifikasi langsung ke nomor mesin dan shift pencetak."),
            ("Pemberantasan Titik Buta:", " Mengakhiri era ringkasan global yang tidak bisa ditindaklanjuti di lapangan."),
            ("Transparansi Hasil Cetak:", " Menampilkan HCS (sempurna) vs HCTS (rusak) per kelompok kerja per PO.")
        ]),
        ("3. TINDAKAN PRESKRIPTIF CEPAT", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN, [
            ("Membedakan Akar Masalah:", " Membedakan penurunan komponen mekanis mesin vs variasi cara kerja operator."),
            ("Perbaikan Terarah (< 2–4 Jam):", " Teknisi langsung menuju mesin bermasalah tanpa memeriksa 9 mesin bergilir."),
            ("Umpan Balik Harian:", " Kepala Kelompok dapat memberikan coaching objektif kepada operator setiap hari.")
        ])
    ]

    col_w = 3.90
    for idx, (f_title, f_col, f_fill, f_bord, bullets) in enumerate(feats):
        cx = 0.60 + idx * 4.10
        add_card(slide, cx, 4.50, col_w, 2.40, f_fill, f_bord)
        add_section_header(slide, cx, 4.50, col_w, 0.35, f_title, f_col)

        tb = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(4.90), Inches(col_w - 0.24), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)

        for b_idx, (b_head, b_desc) in enumerate(bullets):
            p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
            p.space_after = Pt(4)
            r0 = p.add_run()
            r0.text = "• "
            r0.font.name = 'Arial'
            r0.font.size = Pt(9)
            r0.font.bold = True
            r0.font.color.rgb = f_col
            r1 = p.add_run()
            r1.text = b_head
            r1.font.name = 'Arial'
            r1.font.size = Pt(9)
            r1.font.bold = True
            r1.font.color.rgb = DARK_TEXT
            r2 = p.add_run()
            r2.text = b_desc
            r2.font.name = 'Arial'
            r2.font.size = Pt(8.5)
            r2.font.color.rgb = DARK_TEXT

# ==========================================
# SLIDE 13: VALIDASI & PEMBUKTIAN REALISASI S1 2026
# ==========================================
def build_slide_13(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "04.2", "Validasi & Pembuktian S1 2026: Penurunan Inschiet & Efisiensi Finansial",
               "Buktikan efektivitas DSS SIRINE 4.0 melalui data realisasi S1 2026, penurunan inschiet ke 3,33%, dan saving Rp 2,23 Miliar.")

    # 4 Top Hero Cards
    add_kpi_card(slide, 0.60, 1.95, 2.90, 1.15, "INSCHIET AKHIR Q2 2026", "3,33%", "Turun -1,28 pp (-27,8%) vs Baseline 4,61%", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 3.65, 1.95, 2.90, 1.15, "RATA-RATA SEMESTER 1", "3,89%", "Volume: 103.345.688 LK (< 4,00%)", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 6.70, 1.95, 2.90, 1.15, "LEMBAR DISELAMATKAN", "743.234 LK", "Kertas Sekuriti Bernilai Tinggi (S1)", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 9.75, 1.95, 2.98, 1.15, "EFISIENSI FINANSIAL S1", "Rp 2,23 Miliar", "Proyeksi Tahunan: Rp 6,82 Miliar", RED, FILL_LIGHT_AMBER, BORDER_AMBER)

    # Left: Native Chart Inschiet Trend
    add_card(slide, 0.60, 3.25, 4.80, 3.65, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 0.60, 3.25, 4.80, 0.35, "TREN PENURUNAN INSCHIET (%)", PURPLE)

    chart_data = CategoryChartData()
    chart_data.categories = ['Baseline 2025', 'Q1 2026 (Adaptasi)', 'Q2 2026 (Presisi)']
    chart_data.add_series('Inschiet (%)', (4.61, 4.34, 3.33))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.75), Inches(3.65), Inches(4.50), Inches(2.65),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    plots = chart.plots[0]
    plots.has_data_labels = True
    for series in plots.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = GREEN
        dl = series.data_labels
        dl.font.name = 'Arial'
        dl.font.size = Pt(11)
        dl.font.bold = True
        dl.font.color.rgb = DARK_GREEN
        dl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    tb_c_sub = slide.shapes.add_textbox(Inches(0.75), Inches(6.35), Inches(4.50), Inches(0.50))
    tf_c_sub = tb_c_sub.text_frame
    p_c_sub = tf_c_sub.paragraphs[0]
    p_c_sub.text = "Penurunan Kuartal 2: -1,28 pp (-27,77% vs Baseline 4,61%)"
    p_c_sub.font.name = 'Arial'
    p_c_sub.font.size = Pt(8.5)
    p_c_sub.font.bold = True
    p_c_sub.font.color.rgb = GREEN

    # Right: Workpaper Table (Tabel 1.3)
    add_card(slide, 5.55, 3.25, 7.18, 3.65, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_section_header(slide, 5.55, 3.25, 7.18, 0.35, "KERTAS KERJA REALISASI PENGHEMATAN BIAYA SEMESTER 1 2026", GREEN)

    table_shape = slide.shapes.add_table(4, 6, Inches(5.65), Inches(3.70), Inches(6.98), Inches(2.20))
    table = table_shape.table
    table.columns[0].width = Inches(1.48)
    table.columns[1].width = Inches(1.15)
    table.columns[2].width = Inches(0.90)
    table.columns[3].width = Inches(1.05)
    table.columns[4].width = Inches(1.15)
    table.columns[5].width = Inches(1.25)

    headers = ["Periode", "Volume (n)", "Inschiet", "Deviasi", "Diselamatkan", "Nilai Saving*"]
    rows = [
        ("Q1 2026 (Adaptasi)", "57.385.254", "4,34%", "-0,27 pp", "154.940 LK", "Rp 464,82 Juta"),
        ("Q2 2026 (Presisi)", "45.960.434", "3,33%", "-1,28 pp", "588.294 LK", "Rp 1,76 Miliar"),
        ("TOTAL S1 2026", "103.345.688", "3,89%", "-0,72 pp", "743.234 LK", "Rp 2,23 Miliar")
    ]

    for col_idx, h_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.name = 'Arial'
        p.font.size = Pt(8.5)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for row_idx, r_data in enumerate(rows):
        for col_idx, val in enumerate(r_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = FILL_LIGHT_GREEN if row_idx == 2 else (WHITE if row_idx % 2 == 0 else FILL_LIGHT_NAVY)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = 'Arial'
            p.font.size = Pt(8.5)
            p.font.color.rgb = DARK_GREEN if row_idx == 2 else DARK_TEXT
            p.font.bold = (row_idx == 2 or col_idx in [0, 4, 5])
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

    # Right Bottom Box: Downtime Reduction Highlight
    dt_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.65), Inches(6.00), Inches(6.98), Inches(0.80))
    dt_box.fill.solid()
    dt_box.fill.fore_color.rgb = WHITE
    dt_box.line.color.rgb = GREEN
    dt_box.line.width = Pt(1.5)
    tf_dt = dt_box.text_frame
    tf_dt.word_wrap = True
    tf_dt.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_dt.margin_left = Inches(0.12)
    tf_dt.margin_right = Inches(0.12)
    p_dt = tf_dt.paragraphs[0]
    r_d0 = p_dt.add_run()
    r_d0.text = "EFISIENSI WAKTU PEMERIKSAAN & PERBAIKAN MESIN (DOWNTIME):  "
    r_d0.font.name = 'Arial'
    r_d0.font.size = Pt(9.5)
    r_d0.font.bold = True
    r_d0.font.color.rgb = GREEN
    r_d1 = p_dt.add_run()
    r_d1.text = "Durasi penanganan mesin bermasalah terpangkas dari > 1 shift (> 8 jam) menjadi < 2–4 jam per mesin (reduksi jam henti sebesar 50% s.d. 75%)."
    r_d1.font.name = 'Arial'
    r_d1.font.size = Pt(9)
    r_d1.font.color.rgb = DARK_TEXT

# ==========================================
# SLIDE 14: KESIMPULAN STRATEGIS & CLOSING
# ==========================================
def build_slide_14(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "04.3", "Kesimpulan Eksekutif: Transformasi Menyeluruh & Nilai Tambah Korporasi",
               "Rangkum pencapaian menyeluruh dari latar belakang hingga hasil, serta kesiapan unit cetak menyongsong kontrak tender mendatang.")

    # 4 Strategic Impact Pillar Cards
    pillars = [
        ("1. OPERASIONAL & MUTU", "Stabilitas Mutu Terkendali", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY, [
            ("Inschiet Tembus 3,33%:", " Keberhasilan menembus angka 3,33% di Q2 2026 membuktikan efektivitas sistem intervensi presisi."),
            ("Penyelamatan 743.234 Lembar:", " Mencegah pemborosan bahan baku bernilai tinggi selama Semester 1 2026."),
            ("Pemotongan Waktu Henti 50%–75%:", " Teknisi bekerja presisi dengan waktu perbaikan < 2–4 jam (vs > 8 jam).")
        ]),
        ("2. FINANSIAL & EFISIENSI", "Cost Avoidance Miliaran Rupiah", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN, [
            ("Efisiensi Rp 2,23 Miliar (S1):", " Penghematan biaya bahan baku dan operasional cetak riil selama 6 bulan implementasi."),
            ("Proyeksi Rp 6,82 Miliar / Tahun:", " Potensi efisiensi tahunan terproyeksi secara berkelanjutan."),
            ("Investasi Mandiri (In-House):", " Biaya pengembangan sistem Rp 0,- dengan tingkat pengembalian seketika.")
        ]),
        ("3. BUDAYA KERJA & DATA", "Transparansi & Pembinaan", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE, [
            ("Otomasi Meja Mesin < 30 Detik:", " Mengeliminasi buku folio fisik dan rekap manual yang lambat."),
            ("Feedback Objektif Harian:", " Memberikan sarana coaching terarah bagi ±42 operator tanpa friksi antar-shift."),
            ("Kesiapan Smart Factory:", " Membangun fondasi digitalisasi unit kerja sesuai standar INDI 4.0.")
        ]),
        ("4. KEPATUHAN & TENDER DJBC", "Keunggulan Bersaing Berkelanjutan", ORANGE, FILL_LIGHT_AMBER, BORDER_AMBER, [
            ("Integritas Fiskal APBN:", " Menjamin suplai dokumen sekuriti negara PCHT & MMEA tanpa risiko deviasi keaslian."),
            ("Kepatuhan Mutlak SLA:", " Menghilangkan risiko penalti keterlambatan pengiriman ke Bea Cukai."),
            ("Daya Saing Tender Nasional:", " Efisiensi HPP memperkuat posisi tawar Peruri dalam tender pengadaan pita cukai.")
        ])
    ]

    col_w = 2.90
    for idx, (p_title, p_sub, p_col, p_fill, p_bord, bullets) in enumerate(pillars):
        cx = 0.60 + idx * 3.07
        add_card(slide, cx, 1.95, col_w, 4.05, p_fill, p_bord)
        add_section_header(slide, cx, 1.95, col_w, 0.38, p_title, p_col)

        tb = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(2.40), Inches(col_w - 0.24), Inches(3.55))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)

        p_s = tf.paragraphs[0]
        p_s.text = p_sub
        p_s.font.name = 'Arial'
        p_s.font.size = Pt(9)
        p_s.font.bold = True
        p_s.font.color.rgb = p_col
        p_s.space_after = Pt(6)

        for b_idx, (b_head, b_desc) in enumerate(bullets):
            p = tf.add_paragraph()
            p.space_after = Pt(5)
            r0 = p.add_run()
            r0.text = "• "
            r0.font.name = 'Arial'
            r0.font.size = Pt(8.5)
            r0.font.bold = True
            r0.font.color.rgb = p_col
            r1 = p.add_run()
            r1.text = b_head
            r1.font.name = 'Arial'
            r1.font.size = Pt(8.5)
            r1.font.bold = True
            r1.font.color.rgb = DARK_TEXT
            r2 = p.add_run()
            r2.text = b_desc
            r2.font.name = 'Arial'
            r2.font.size = Pt(8)
            r2.font.color.rgb = DARK_TEXT

    # Closing Executive Commitment Banner
    close_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(6.15), Inches(12.13), Inches(0.75))
    close_box.fill.solid()
    close_box.fill.fore_color.rgb = NAVY
    close_box.line.fill.background()
    tf_close = close_box.text_frame
    tf_close.word_wrap = True
    tf_close.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_close.margin_left = Inches(0.20)
    tf_close.margin_right = Inches(0.20)
    p_close = tf_close.paragraphs[0]
    r_c0 = p_close.add_run()
    r_c0.text = "KOMITMEN TRANSFORMASI MUTU UNIT CETAK PITA CUKAI:  "
    r_c0.font.name = 'Arial'
    r_c0.font.size = Pt(10)
    r_c0.font.bold = True
    r_c0.font.color.rgb = WHITE
    r_c1 = p_close.add_run()
    r_c1.text = "Melalui DSS SIRINE 4.0, Unit Cetak Pita Cukai membuktikan bahwa integrasi data meja mesin adalah kunci fundamental dalam mewujudkan operational excellence, melindungi bahan baku negara, dan memenangkan masa depan percetakan sekuriti Peruri."
    r_c1.font.name = 'Arial'
    r_c1.font.size = Pt(9.5)
    r_c1.font.color.rgb = WHITE

# ==========================================
# MAIN EXECUTION
# ==========================================
def main():
    prs, blank_layout = create_presentation()

    print("Building Slide 1: Cover Slide...")
    build_slide_1(prs, blank_layout)

    print("Building Slide 2: Ringkasan Eksekutif & Hero Metrics...")
    build_slide_2(prs, blank_layout)

    print("Building Slide 3: Dinamika Pasar Global...")
    build_slide_3(prs, blank_layout)

    print("Building Slide 4: Skala Tender Nasional DJBC...")
    build_slide_4(prs, blank_layout)

    print("Building Slide 5: Arah Strategis & Mandat Peruri...")
    build_slide_5(prs, blank_layout)

    print("Building Slide 6: Realitas Lapangan & Profil 9 Mesin...")
    build_slide_6(prs, blank_layout)

    print("Building Slide 7: Analisis Fluktuasi Baseline 2025...")
    build_slide_7(prs, blank_layout)

    print("Building Slide 8: Kesenjangan Operasional (Data Silo)...")
    build_slide_8(prs, blank_layout)

    print("Building Slide 9: Implikasi Titik Buta Data...")
    build_slide_9(prs, blank_layout)

    print("Building Slide 10: Simulasi Finansial Baseline 2025...")
    build_slide_10(prs, blank_layout)

    print("Building Slide 11: Matriks Evaluasi 5 Pilar Cost of Inaction...")
    build_slide_11(prs, blank_layout)

    print("Building Slide 12: Solusi Integrasi Data DSS SIRINE 4.0...")
    build_slide_12(prs, blank_layout)

    print("Building Slide 13: Validasi & Pembuktian Realisasi S1 2026...")
    build_slide_13(prs, blank_layout)

    print("Building Slide 14: Kesimpulan Strategis & Closing...")
    build_slide_14(prs, blank_layout)

    output_path = "Presentasi_Risalah_Latar_Belakang_IAKA_2026.pptx"
    prs.save(output_path)
    print(f"\n[SUCCESS] Master Deck successfully generated and saved to: {output_path}")

if __name__ == "__main__":
    main()
