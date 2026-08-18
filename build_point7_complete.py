import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sys
import os

# --- PALETTE CONSTANTS ---
NAVY = RGBColor(32, 40, 143)         # #20288F - Primary Brand
PURPLE = RGBColor(106, 44, 201)      # #6A2CC9 - Accent
DARK_NAVY = RGBColor(26, 35, 126)    # #1A237E
GREEN = RGBColor(46, 125, 50)        # #2E7D32 - Success
DARK_GREEN = RGBColor(27, 94, 32)    # #1B5E20
RED = RGBColor(198, 40, 40)          # #C62828 - Alert
ORANGE = RGBColor(230, 81, 0)        # #E65100 - Warning
DARK_TEXT = RGBColor(42, 42, 61)     # #2A2A3D - Body Text
MUTED_TEXT = RGBColor(110, 110, 130) # #6E6E82 - Muted Text
WHITE = RGBColor(255, 255, 255)

# Container Fills
FILL_LIGHT_NAVY = RGBColor(244, 247, 254)   # #F4F7FE
FILL_LIGHT_PURPLE = RGBColor(248, 246, 254) # #F8F6FE
FILL_LIGHT_GREEN = RGBColor(232, 245, 233)  # #E8F5E9
FILL_LIGHT_AMBER = RGBColor(255, 249, 230)  # #FFF9E6
FILL_LIGHT_RED = RGBColor(255, 245, 245)    # #FFF5F5

# Borders
BORDER_NAVY = RGBColor(210, 220, 245)       # #D2DCF5
BORDER_PURPLE = RGBColor(217, 206, 247)     # #D9CEF7
BORDER_GREEN = RGBColor(165, 214, 167)      # #A5D6A7
BORDER_AMBER = RGBColor(255, 224, 130)      # #FFE082
BORDER_RED = RGBColor(255, 205, 210)        # #FFCDD2

def add_header(slide, badge_num, title_text, purpose_text):
    # Wave ornament
    if os.path.exists('wave_ornament.png'):
        slide.shapes.add_picture('wave_ornament.png', Inches(-0.118), Inches(-0.196), Inches(1.584), Inches(0.891))
    # Peruri Logo
    if os.path.exists('peruri_logo.png'):
        slide.shapes.add_picture('peruri_logo.png', Inches(12.09), Inches(0.013), Inches(1.247), Inches(0.694))

    # Badge Box
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(0.65), Inches(0.80), Inches(0.72))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PURPLE
    badge.line.color.rgb = PURPLE
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = badge_num
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    tb_title = slide.shapes.add_textbox(Inches(1.55), Inches(0.60), Inches(10.20), Inches(0.80))
    tf = tb_title.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(23)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Purpose Subtitle
    tb_purpose = slide.shapes.add_textbox(Inches(0.60), Inches(1.48), Inches(11.80), Inches(0.38))
    tf = tb_purpose.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Tujuan bagian ini:  "
    r1.font.name = 'Arial'
    r1.font.size = Pt(11.5)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = purpose_text
    r2.font.name = 'Arial'
    r2.font.size = Pt(11.5)
    r2.font.color.rgb = DARK_TEXT

    # Footer
    tb_foot = slide.shapes.add_textbox(Inches(0.60), Inches(7.02), Inches(12.13), Inches(0.24))
    tf = tb_foot.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = "IAKA 2026 — Kerangka Presentasi Peserta  ·  DSS SIRINE 4.0 Unit Cetak Pita Cukai"
    p.font.name = 'Arial'
    p.font.size = Pt(8.5)
    p.font.color.rgb = MUTED_TEXT

def add_kpi_card(slide, left, top, width, height, title, value, subtext, color, fill_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.04)
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)

    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = 'Arial'
    p0.font.size = Pt(8.5)
    p0.font.bold = True
    p0.font.color.rgb = color
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = value
    p1.font.name = 'Arial'
    p1.font.size = Pt(15.5)
    p1.font.bold = True
    p1.font.color.rgb = color
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = subtext
    p2.font.name = 'Arial'
    p2.font.size = Pt(7.2)
    p2.font.color.rgb = DARK_TEXT
    p2.alignment = PP_ALIGN.CENTER

def style_table_cell(cell, text, font_size=7.2, bold=False, text_color=DARK_TEXT, fill_color=None, align=PP_ALIGN.LEFT):
    cell.text = ""
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    else:
        cell.fill.background()
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = align

# -------------------------------------------------------------
# BUILD SLIDE 07.1
# -------------------------------------------------------------
def build_slide_7_1(slide):
    # Clear shapes on slide
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    add_header(
        slide,
        badge_num="07.1",
        title_text="Sustainability: Standardisasi Tata Kelola, Registrasi SOP/IK Baru & Penutupan Sistem Lama",
        purpose_text="Kunci cara kerja baru agar mengikat secara legal-formal dalam tata kelola mutu operasional dan tidak bergantung pada ingatan individu."
    )

    # Top 4 Hero KPI Cards
    add_kpi_card(slide, 0.60, 1.90, 2.86, 1.08, "DOKUMEN RESMI DISAHKAN", "2 SOP / IK Baru", "IK-PPC-2026-001 & SOP-PPC-2026-004", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 3.69, 1.90, 2.86, 1.08, "PENUTUPAN SISTEM LAMA", "100% Non-Aktif", "Berita Acara BA-PPC-2026-002 (Buku Folio Ditarik)", RED, FILL_LIGHT_RED, BORDER_RED)
    add_kpi_card(slide, 6.78, 1.90, 2.86, 1.08, "STANDAR MUTU ISO", "ISO 9001:2015", "Klausul 8.5.2 (Traceability) & 9.1.3 (Evaluasi Data)", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 9.87, 1.90, 2.86, 1.08, "DUAL-SYSTEM ELIMINATION", "0 Form Manual", "Nol Beban Ganda | Efisiensi 100% Digital Real-Time", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)

    # Left Panel: 1. REGISTRASI STANDAR OPERASIONAL PROSEDUR & INSTRUKSI KERJA BARU
    header_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.60), Inches(3.08), Inches(5.95), Inches(0.34))
    header_left.fill.solid()
    header_left.fill.fore_color.rgb = NAVY
    header_left.line.fill.background()
    tf = header_left.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "1. REGISTRASI STANDAR OPERASIONAL PROSEDUR & INSTRUKSI KERJA BARU"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(3.42), Inches(5.95), Inches(3.02))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = FILL_LIGHT_NAVY
    box_left.line.color.rgb = BORDER_NAVY
    box_left.line.width = Pt(1.0)

    # 3 Cards on Left Panel
    left_cards = [
        (
            Inches(3.48), Inches(0.88), PURPLE, BORDER_PURPLE,
            "1. IK PENGISIAN KONFIRMASI PO CETAK DIGITAL (IK-PPC-2026-001)",
            "• Nomor Dokumen & Pengesahan: ", "No. IK-PPC-2026-001 (Rev. 00) | Disahkan oleh Kadep SBU High Security Solution.",
            "• Ketetapan Operasional: ", "Wajib entri real-time tiap PO selesai (< 30 detik via Autofill) di meja kontrol 6 mesin cetak."
        ),
        (
            Inches(4.42), Inches(0.92), NAVY, BORDER_NAVY,
            "2. SOP TINDAKAN PERBAIKAN MESIN BERBASIS PARETO SIRINE (SOP-PPC-2026-004)",
            "• Nomor Dokumen & Pengesahan: ", "No. SOP-PPC-2026-004 (Rev. 00) | Disahkan oleh Kepala Divisi Teknik & Pemeliharaan.",
            "• Ketetapan Operasional: ", "Teknisi wajib membuka modul Pareto cacat mesin spesifik sebelum servis mekanis komponen."
        ),
        (
            Inches(5.40), Inches(0.86), DARK_GREEN, BORDER_GREEN,
            "3. INTEGRASI DOKUMEN TURUNAN & CHECKLIST SERAH TERIMA GILIR",
            "• Pembaruan Formulir Serah Terima: ", "Verifikasi kelengkapan data PO jadi syarat mutlak serah terima antar-shift.",
            "• Standar Kontrol Mutu: ", "Pemeriksaan konsistensi setting mesin (rol karet, blanket, register) terdokumentasi digital."
        )
    ]

    for top_pos, h_size, title_col, border_col, title_txt, l1_lbl, l1_val, l2_lbl, l2_val in left_cards:
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), top_pos, Inches(5.71), h_size)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = border_col
        c.line.width = Pt(1.0)
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.02)
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = title_txt
        p.font.name = 'Arial'
        p.font.size = Pt(7.3)
        p.font.bold = True
        p.font.color.rgb = title_col

        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = l1_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p1.add_run()
        r2.text = l1_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = l2_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p2.add_run()
        r2.text = l2_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

    # Label under left panel
    tb_src_l1 = slide.shapes.add_textbox(Inches(0.72), Inches(6.28), Inches(5.71), Inches(0.14))
    tf_src_l1 = tb_src_l1.text_frame
    tf_src_l1.word_wrap = False
    tf_src_l1.margin_top = Inches(0)
    tf_src_l1.margin_left = Inches(0)
    p = tf_src_l1.paragraphs[0]
    p.text = "*(Dokumen Resmi Terdaftar di Portal Manajemen Mutu Perum Peruri & Berlaku Efektif Sejak Jan 2026)"
    p.font.name = 'Arial'
    p.font.size = Pt(6.3)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY

    # Right Panel: 2. PENUTUPAN SISTEM LAMA & JAMINAN KEPATUHAN ISO 9001
    header_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.70), Inches(3.08), Inches(6.03), Inches(0.34))
    header_right.fill.solid()
    header_right.fill.fore_color.rgb = PURPLE
    header_right.line.fill.background()
    tf = header_right.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "2. PENUTUPAN SISTEM LAMA & JAMINAN KEPATUHAN ISO 9001"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.70), Inches(3.42), Inches(6.03), Inches(3.02))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = FILL_LIGHT_PURPLE
    box_right.line.color.rgb = BORDER_PURPLE
    box_right.line.width = Pt(1.0)

    # 3 Cards on Right Panel
    right_cards = [
        (
            Inches(3.48), Inches(0.88), RED, BORDER_RED,
            "1. PENARIKAN RESMI BUKU FOLIO FISIK (BERITA ACARA BA-PPC-2026-002)",
            "• Eliminasi Total Sistem Lama: ", "Buku folio ditarik serentak dari 6 meja mesin cetak per 1 Januari 2026.",
            "• Mencegah Beban Ganda: ", "Menghilangkan kebiasaan mencatat ganda (kertas & PC) serta risiko hilang/rusak fisik."
        ),
        (
            Inches(4.42), Inches(0.92), PURPLE, BORDER_PURPLE,
            "2. PEMENUHAN STANDAR PENJAMINAN MUTU INTERNASIONAL (ISO 9001:2015)",
            "• Klausul 8.5.2 (Mampu Telusur): ", "Silsilah audit lengkap: Nomor PO ➔ Mesin ➔ Gilir/Shift ➔ Tim ➔ Lembar Rusak.",
            "• Klausul 9.1.3 (Analisis Data): ", "Data diolah otomatis jadi grafik tren & Pareto tanpa delay waktu rekap manual."
        ),
        (
            Inches(5.40), Inches(0.86), NAVY, BORDER_NAVY,
            "3. KEAMANAN & KEDAULATAN DATA SEKURITI NEGARA (ON-PREMISE INTRANET)",
            "• Infrastruktur Server Internal: ", "Tersimpan 100% di server intranet tertutup Peruri (zero vendor lock-in & leak risk).",
            "• Hak Akses Berjenjang: ", "Otorisasi role bertingkat: Operator Meja Mesin ➔ Kepala Kelompok ➔ Kepala Unit."
        )
    ]

    for top_pos, h_size, title_col, border_col, title_txt, r1_lbl, r1_val, r2_lbl, r2_val in right_cards:
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.82), top_pos, Inches(5.79), h_size)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = border_col
        c.line.width = Pt(1.0)
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.02)
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = title_txt
        p.font.name = 'Arial'
        p.font.size = Pt(7.3)
        p.font.bold = True
        p.font.color.rgb = title_col

        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = r1_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p1.add_run()
        r2.text = r1_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = r2_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p2.add_run()
        r2.text = r2_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

    # Label under right panel
    tb_src_r1 = slide.shapes.add_textbox(Inches(6.82), Inches(6.28), Inches(5.79), Inches(0.14))
    tf_src_r1 = tb_src_r1.text_frame
    tf_src_r1.word_wrap = False
    tf_src_r1.margin_top = Inches(0)
    tf_src_r1.margin_left = Inches(0)
    p = tf_src_r1.paragraphs[0]
    p.text = "*(Audit Kepatuhan: Berita Acara BA-PPC-2026-002 & Sistem Manajemen Terintegrasi Peruri)"
    p.font.name = 'Arial'
    p.font.size = Pt(6.3)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Bottom Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(6.52), Inches(12.13), Inches(0.44))
    banner.fill.solid()
    banner.fill.fore_color.rgb = FILL_LIGHT_PURPLE
    banner.line.color.rgb = BORDER_PURPLE
    banner.line.width = Pt(1.0)
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "• Jaminan Legalitas & Tata Kelola Baku: "
    r1.font.name = 'Arial'
    r1.font.size = Pt(7.5)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = "Metode kerja baru telah terikat secara formal melalui IK-PPC-2026-001 (Konfirmasi PO Digital) & SOP-PPC-2026-004 (Servis Mesin Pareto), didukung Berita Acara BA-PPC-2026-002 penarikan buku folio manual, menjamin keberlanjutan proses kerja digital secara permanen."
    r2.font.name = 'Arial'
    r2.font.size = Pt(7.4)
    r2.font.color.rgb = DARK_TEXT

# -------------------------------------------------------------
# BUILD SLIDE 07.2
# -------------------------------------------------------------
def build_slide_7_2(slide):
    # Clear shapes on slide
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    add_header(
        slide,
        badge_num="07.2",
        title_text="Sustainability: Transfer Knowledge, Pelatihan Mandiri & Audit Keberlanjutan",
        purpose_text="Kunci cara kerja baru agar tidak bergantung pada ingatan atau satu orang melalui transfer pengetahuan terstruktur, modul in-app, dan audit berkala."
    )

    # Top 4 Hero KPI Cards
    add_kpi_card(slide, 0.60, 1.90, 2.86, 1.08, "PARTISIPASI PELATIHAN", "100% Operator (42 Org)", "Mencakup 3 Shift Penuh (Pagi, Sore, Malam)", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 3.69, 1.90, 2.86, 1.08, "EVALUASI KOMPETENSI", "Rata-rata 94,8 / 100", "Post-Test Simulasi Meja Mesin & Troubleshooting", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 6.78, 1.90, 2.86, 1.08, "IN-APP GUIDANCE SYSTEM", "Tersedia Real-Time", "Panduan Interaktif Tersemat di Setiap Fitur Web", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 9.87, 1.90, 2.86, 1.08, "RUTINITAS AUDIT DATA", "Harian & Mingguan", "Verifikasi Checklist Shift & Rekonsiliasi Log SAP", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)

    # Left Panel: 1. MATRIKS PROGRAM PELATIHAN & BUKTI TRANSFER KNOWLEDGE
    header_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.60), Inches(3.08), Inches(6.05), Inches(0.34))
    header_left.fill.solid()
    header_left.fill.fore_color.rgb = NAVY
    header_left.line.fill.background()
    tf = header_left.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "1. MATRIKS PROGRAM PELATIHAN & BUKTI TRANSFER KNOWLEDGE"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(3.42), Inches(6.05), Inches(3.02))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = FILL_LIGHT_NAVY
    box_left.line.color.rgb = BORDER_NAVY
    box_left.line.width = Pt(1.0)

    # Table of Training Matrix
    table_shape_l = slide.shapes.add_table(5, 5, Inches(0.72), Inches(3.48), Inches(5.81), Inches(1.82))
    t_l = table_shape_l.table
    t_l.columns[0].width = Inches(1.35)
    t_l.columns[1].width = Inches(1.15)
    t_l.columns[2].width = Inches(0.95)
    t_l.columns[3].width = Inches(1.15)
    t_l.columns[4].width = Inches(1.21)

    headers_l = ["Program & Modul", "Sasaran Peserta", "Waktu / Tgl", "Metode Pelaksanaan", "Hasil Capaian"]
    for c_idx, h in enumerate(headers_l):
        style_table_cell(t_l.cell(0, c_idx), h, font_size=6.8, bold=True, text_color=WHITE, fill_color=NAVY, align=PP_ALIGN.CENTER if c_idx>0 else PP_ALIGN.LEFT)

    rows_l = [
        ("Workshop Form PO Digital", "42 Operator & PIC (3 Gilir)", "15–18 Des 2025", "Praktik Langsung Meja Mesin", "100% Lulus Input (<30 dtk)", WHITE),
        ("Training Jadwal & Rotasi", "12 Kepala Kelompok", "22–23 Des 2025", "Simulasi Grid Jadwal Gilir", "Mandiri Kelola 3 Shift", WHITE),
        ("Dasbor Preskriptif & Pareto", "Kepala Unit & Kasie", "28–29 Des 2025", "Studi Kasus Analisis Cacat", "Mampu Servis Terarah", WHITE),
        ("Pendampingan On-The-Job", "Seluruh Tim Meja Mesin", "Jan 2026 (W1–4)", "Mentoring Meja Mesin Gilir", "Nol Eror Entri Data PO", FILL_LIGHT_GREEN)
    ]
    for r_idx, (p_prog, p_part, p_time, p_meth, p_out, bg_col) in enumerate(rows_l, start=1):
        is_bold = (r_idx == 4)
        txt_col = DARK_GREEN if is_bold else DARK_TEXT
        style_table_cell(t_l.cell(r_idx, 0), p_prog, font_size=6.5, bold=True, text_color=DARK_NAVY, fill_color=bg_col)
        style_table_cell(t_l.cell(r_idx, 1), p_part, font_size=6.5, bold=False, text_color=DARK_TEXT, fill_color=bg_col, align=PP_ALIGN.CENTER)
        style_table_cell(t_l.cell(r_idx, 2), p_time, font_size=6.5, bold=False, text_color=DARK_TEXT, fill_color=bg_col, align=PP_ALIGN.CENTER)
        style_table_cell(t_l.cell(r_idx, 3), p_meth, font_size=6.5, bold=False, text_color=DARK_TEXT, fill_color=bg_col)
        style_table_cell(t_l.cell(r_idx, 4), p_out, font_size=6.5, bold=is_bold, text_color=txt_col, fill_color=bg_col, align=PP_ALIGN.CENTER)

    # In-App Guidance Card below Table
    card_inapp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(5.36), Inches(5.81), Inches(0.86))
    card_inapp.fill.solid()
    card_inapp.fill.fore_color.rgb = WHITE
    card_inapp.line.color.rgb = BORDER_PURPLE
    card_inapp.line.width = Pt(1.0)
    tf = card_inapp.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.02)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)

    p = tf.paragraphs[0]
    p.text = "FITUR IN-APP GUIDANCE & MODUL PANDUAN INTERAKTIF TERINTEGRASI"
    p.font.name = 'Arial'
    p.font.size = Pt(7.4)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    p1 = tf.add_paragraph()
    r1 = p1.add_run()
    r1.text = "• Panduan Interaktif Tersemat: "
    r1.font.bold = True
    r1.font.size = Pt(6.5)
    r1.font.color.rgb = DARK_NAVY
    r2 = p1.add_run()
    r2.text = "Sistem dilengkapi tooltip dan bantuan langkah per langkah langsung di aplikasi web."
    r2.font.size = Pt(6.5)
    r2.font.color.rgb = DARK_TEXT

    p2 = tf.add_paragraph()
    r1 = p2.add_run()
    r1.text = "• Pelatihan Mandiri Cepat: "
    r1.font.bold = True
    r1.font.size = Pt(6.5)
    r1.font.color.rgb = DARK_NAVY
    r2 = p2.add_run()
    r2.text = "Operator baru dapat mahir menginput data dalam < 1 hari tanpa ketergantungan dokumen fisik."
    r2.font.size = Pt(6.5)
    r2.font.color.rgb = DARK_TEXT

    # Label under left panel
    tb_src_l2 = slide.shapes.add_textbox(Inches(0.72), Inches(6.26), Inches(5.81), Inches(0.14))
    tf_src_l2 = tb_src_l2.text_frame
    tf_src_l2.word_wrap = False
    tf_src_l2.margin_top = Inches(0)
    tf_src_l2.margin_left = Inches(0)
    p = tf_src_l2.paragraphs[0]
    p.text = "*(Bukti Pelaksanaan: Daftar Hadir Pelatihan, Berita Acara Sosialisasi & Nilai Post-Test Terarsip)"
    p.font.name = 'Arial'
    p.font.size = Pt(6.3)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY

    # Right Panel: 2. MEKANISME AUDIT KEBERLANJUTAN, KONTROL BERJENJANG & ROADMAP
    header_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.78), Inches(3.08), Inches(5.95), Inches(0.34))
    header_right.fill.solid()
    header_right.fill.fore_color.rgb = PURPLE
    header_right.line.fill.background()
    tf = header_right.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "2. MEKANISME AUDIT KEBERLANJUTAN, KONTROL BERJENJANG & ROADMAP"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.78), Inches(3.42), Inches(5.95), Inches(3.02))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = FILL_LIGHT_PURPLE
    box_right.line.color.rgb = BORDER_PURPLE
    box_right.line.width = Pt(1.0)

    # 3 Cards on Right Panel
    right_cards2 = [
        (
            Inches(3.48), Inches(0.92), NAVY, BORDER_NAVY,
            "1. MEKANISME KONTROL 3 TINGKAT (THREE-TIER AUDIT GOVERNANCE)",
            "• Tingkat 1 (Harian/Shift): ", "Kepala Kelompok periksa checklist kelengkapan input PO sebelum serah terima.",
            "• Tingkat 2 (Mingguan/Bulanan): ", "Supervisor Unit Cetak rekonsiliasi log SIRINE vs SAP ZPPRSIPPC0012 & Verifikasi."
        ),
        (
            Inches(4.46), Inches(0.88), GREEN, BORDER_GREEN,
            "2. ELIMINASI KETERGANTUNGAN INDIVIDU (ZERO SINGLE POINT OF FAILURE)",
            "• Repositori Kode Terpusat: ", "Source code modular (Laravel Service, Inertia, Vue) tersimpan di Git internal Peruri.",
            "• Dokumentasi Teknis Lengkap: ", "Kamus data & dokumentasi arsitektur siap dipelihara tim IT Peruri mana pun."
        ),
        (
            Inches(5.40), Inches(0.86), PURPLE, BORDER_PURPLE,
            "3. RENCANA KEBERLANJUTAN JANGKA PANJANG (DEVELOPMENT ROADMAP)",
            "• Q3–Q4 2026: ", "Implementasi Rapor Scoring Tim (Grade A–E) & Integrasi Maintenance Log Book.",
            "• 2027: ", "Replikasi arsitektur sistem ke lini Meterai, Paspor RI, Khazanah, dan Uang Kertas."
        )
    ]

    for top_pos, h_size, title_col, border_col, title_txt, r1_lbl, r1_val, r2_lbl, r2_val in right_cards2:
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.90), top_pos, Inches(5.71), h_size)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = border_col
        c.line.width = Pt(1.0)
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.02)
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = title_txt
        p.font.name = 'Arial'
        p.font.size = Pt(7.3)
        p.font.bold = True
        p.font.color.rgb = title_col

        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = r1_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p1.add_run()
        r2.text = r1_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = r2_lbl
        r1.font.bold = True
        r1.font.size = Pt(6.5)
        r1.font.color.rgb = DARK_NAVY
        r2 = p2.add_run()
        r2.text = r2_val
        r2.font.size = Pt(6.5)
        r2.font.color.rgb = DARK_TEXT

    # Label under right panel
    tb_src_r2 = slide.shapes.add_textbox(Inches(6.90), Inches(6.26), Inches(5.71), Inches(0.14))
    tf_src_r2 = tb_src_r2.text_frame
    tf_src_r2.word_wrap = False
    tf_src_r2.margin_top = Inches(0)
    tf_src_r2.margin_left = Inches(0)
    p = tf_src_r2.paragraphs[0]
    p.text = "*(Audit Rutin: Sesuai Siklus Audit Internal ISO 9001:2015 & Rencana Kerja Operasional 2026)"
    p.font.name = 'Arial'
    p.font.size = Pt(6.3)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Bottom Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(6.52), Inches(12.13), Inches(0.44))
    banner.fill.solid()
    banner.fill.fore_color.rgb = FILL_LIGHT_GREEN
    banner.line.color.rgb = BORDER_GREEN
    banner.line.width = Pt(1.0)
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "• Jaminan Kemandirian & Keberlanjutan Sistem: "
    r1.font.name = 'Arial'
    r1.font.size = Pt(7.5)
    r1.font.bold = True
    r1.font.color.rgb = DARK_GREEN
    r2 = p.add_run()
    r2.text = "Kombinasi pelatihan 100% operator (42 orang / 3 gilir), modul in-app terintegrasi, audit berjenjang 3 tingkat, dan repositori kode internal memastikan DSS SIRINE 4.0 beroperasi mandiri dan berkelanjutan tanpa bergantung pada satu individu."
    r2.font.name = 'Arial'
    r2.font.size = Pt(7.4)
    r2.font.color.rgb = DARK_TEXT

# -------------------------------------------------------------
# MAIN EXECUTION
# -------------------------------------------------------------
def main():
    prs = pptx.Presentation('Kerangka_Presentasi_Peserta_IAKA_2026 (1).pptx')
    print(f"Total slides initially: {len(prs.slides)}")

    # Slide 16 (index 15) is Point 7 Sustainability Template
    # We will build Slide 07.1 on slide 16 (idx 15)
    slide_7_1 = prs.slides[15]
    build_slide_7_1(slide_7_1)
    print("Slide 07.1 built successfully on index 15.")

    # Check if slide 07.2 already exists or needs to be inserted
    # In the current 18-slide deck:
    # idx 15: Point 7
    # idx 16: Point 8 (Lesson Learned)
    # idx 17: Closing (Terima Kasih)
    # We need to insert a new slide right after index 15 (at index 16) for 07.2!
    if len(prs.slides) == 18:
        new_slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Move new_slide to index 16 (right after index 15)
        sldIdLst = prs.slides._sldIdLst
        slide_elem = sldIdLst[-1]
        sldIdLst.remove(slide_elem)
        sldIdLst.insert(16, slide_elem)
        slide_7_2 = prs.slides[16]
    else:
        slide_7_2 = prs.slides[16]

    build_slide_7_2(slide_7_2)
    print("Slide 07.2 built successfully on index 16.")

    prs.save('Kerangka_Presentasi_Peserta_IAKA_2026 (1).pptx')
    print(f"Saved! Kerangka_Presentasi_Peserta_IAKA_2026 (1).pptx updated. Total slides: {len(prs.slides)}")

if __name__ == '__main__':
    main()
