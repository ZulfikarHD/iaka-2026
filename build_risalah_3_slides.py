import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
import sys
import os

# ==========================================
# PALETTE & DESIGN CONSTANTS (Peruri Theme)
# ==========================================
NAVY = RGBColor(32, 40, 143)         # #20288F - Primary Brand Navy
PURPLE = RGBColor(106, 44, 201)      # #6A2CC9 - Accent Purple
DARK_NAVY = RGBColor(26, 35, 126)    # #1A237E - Deep Navy
GREEN = RGBColor(46, 125, 50)        # #2E7D32 - Success Green
DARK_GREEN = RGBColor(27, 94, 32)    # #1B5E20
RED = RGBColor(198, 40, 40)          # #C62828 - Alert Red
ORANGE = RGBColor(230, 81, 0)        # #E65100 - Warning Amber
DARK_TEXT = RGBColor(42, 42, 61)     # #2A2A3D - Body Text Charcoal
MUTED_TEXT = RGBColor(110, 110, 130) # #6E6E82 - Muted Text Gray
WHITE = RGBColor(255, 255, 255)

# Container Fills
FILL_LIGHT_NAVY = RGBColor(244, 247, 254)   # #F4F7FE
FILL_LIGHT_PURPLE = RGBColor(248, 246, 254) # #F8F6FE
FILL_LIGHT_GREEN = RGBColor(232, 245, 233)  # #E8F5E9
FILL_LIGHT_AMBER = RGBColor(255, 249, 230)  # #FFF9E6
FILL_LIGHT_RED = RGBColor(255, 245, 245)    # #FFF5F5

# Container Borders
BORDER_NAVY = RGBColor(210, 220, 245)       # #D2DCF5
BORDER_PURPLE = RGBColor(217, 206, 247)     # #D9CEF7
BORDER_GREEN = RGBColor(165, 214, 167)      # #A5D6A7
BORDER_AMBER = RGBColor(255, 224, 130)      # #FFE082
BORDER_RED = RGBColor(255, 205, 210)        # #FFCDD2

def create_presentation():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    return prs, blank_layout

def add_header(slide, badge_num, title_text, purpose_text):
    if os.path.exists('wave_ornament.png'):
        slide.shapes.add_picture('wave_ornament.png', Inches(-0.118), Inches(-0.196), Inches(1.584), Inches(0.891))
    
    if os.path.exists('peruri_logo.png'):
        slide.shapes.add_picture('peruri_logo.png', Inches(12.09), Inches(0.013), Inches(1.247), Inches(0.694))

    # Badge Box
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.55), Inches(0.75), Inches(0.68))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PURPLE
    badge.line.color.rgb = PURPLE
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = badge_num
    p.font.name = 'Arial'
    p.font.size = Pt(18 if len(badge_num) <= 2 else 15)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    tb_title = slide.shapes.add_textbox(Inches(1.42), Inches(0.50), Inches(10.50), Inches(0.75))
    tf = tb_title.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Purpose Subtitle
    tb_purpose = slide.shapes.add_textbox(Inches(0.55), Inches(1.30), Inches(12.23), Inches(0.32))
    tf = tb_purpose.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Tujuan bagian ini:  "
    r1.font.name = 'Arial'
    r1.font.size = Pt(10)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = purpose_text
    r2.font.name = 'Arial'
    r2.font.size = Pt(10)
    r2.font.color.rgb = DARK_TEXT

    # Footer
    tb_foot = slide.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(12.23), Inches(0.22))
    tf = tb_foot.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = "IAKA 2026 — Risalah Latar Belakang & Identifikasi Masalah  ·  DSS SIRINE 4.0 Unit Cetak Pita Cukai (Perum Peruri)"
    p.font.name = 'Arial'
    p.font.size = Pt(8)
    p.font.color.rgb = MUTED_TEXT

def add_kpi_card(slide, left, top, width, height, title, value, subtext, color, fill_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(left + 0.08), Inches(top + 0.06), Inches(width - 0.16), Inches(height - 0.10))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)

    p1 = tf.paragraphs[0]
    p1.text = title.upper()
    p1.font.name = 'Arial'
    p1.font.size = Pt(7.5)
    p1.font.bold = True
    p1.font.color.rgb = color

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.name = 'Arial'
    p2.font.size = Pt(14.5)
    p2.font.bold = True
    p2.font.color.rgb = color

    p3 = tf.add_paragraph()
    p3.text = subtext
    p3.font.name = 'Arial'
    p3.font.size = Pt(7.5)
    p3.font.color.rgb = DARK_TEXT

def add_card(slide, left, top, width, height, fill_color, border_color, border_width=1.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(border_width)
    return card

def add_section_header(slide, left, top, width, height, title, bg_color):
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = bg_color
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(8.5)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    return hdr


# ==============================================================================
# SLIDE 1 OF 3: NARASI 1, 2 & 3
# (1. Dinamika Pasar Global, 2. Skala Tender DJBC, 3. Arah Strategis Peruri)
# ==============================================================================
def build_slide_1(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide, "01",
        "Konteks Latar Belakang (1/3): Dinamika Pasar Global, Skala Tender DJBC & Mandat Peruri",
        "Bangun urgensi dari standar presisi percetakan sekuriti internasional (Intergraf/WCO), kontrak tender DJBC Kemenkeu RI (177,6M LK), dan mandat korporasi PP 06/2019."
    )

    # 4 Top Hero KPIs
    add_kpi_card(slide, 0.55, 1.68, 2.95, 0.98, "TARGET VOLUME PPIC", "160.000.000 LK", "Kapasitas Rata-rata Tahunan Unit Cetak", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 3.62, 1.68, 2.95, 0.98, "REALISASI PESANAN 2025", "177.636.930 LK", "Pesanan Aktual Modul SAP ZPPRSIPPC0012", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 6.69, 1.68, 2.95, 0.98, "PERAN FISKAL APBN", "Ratusan Triliun", "Penerimaan Cukai Negara (PCHT & MMEA)", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 9.76, 1.68, 3.02, 0.98, "9 MESIN CETAK OFFSET", "4 KMR · 2 RYB · 3 GTO", "Pola 3 Shift 24/7 (±42 Personel Operator)", DARK_NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)

    # 3 Large Vertical Columns (Sections 1, 2, 3)
    col_w = 3.98
    gap = 0.14
    top_y = 2.76
    card_h = 4.30

    # Column 1: Bab 1. Dinamika Pasar Global: Standar Presisi & Keunggulan Bersaing
    add_card(slide, 0.55, top_y, col_w, card_h, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_section_header(slide, 0.55, top_y, col_w, 0.28, "1. DINAMIKA PASAR GLOBAL: STANDAR PRESISE & BERSAING", NAVY)

    tb_1 = slide.shapes.add_textbox(Inches(0.65), Inches(top_y + 0.32), Inches(col_w - 0.20), Inches(card_h - 0.38))
    tf_1 = tb_1.text_frame
    tf_1.word_wrap = True
    tf_1.margin_top = Inches(0)
    tf_1.margin_left = Inches(0)

    p1_items = [
        ("Standar Internasional (Intergraf & WCO):",
         "Kepatuhan terhadap standar global percetakan sekuriti tinggi (*high-security printing*) wajib presisi tinggi guna mencegah pemalsuan (*anti-counterfeiting*). Menerapkan fitur pengamanan berlapis (*multi-layer security*): kertas serat pengaman tak kasat mata (*security fibers*), tinta UV berpendar, ornamen garis halus *guilloche*, *microtext*, dan pita hologram presisi mikron."),
        ("4 Kriteria Utama Keunggulan Bersaing:",
         "1. Jaminan Autentikasi Tanpa Cacat (Zero-Defect): Cacat fisik (blobor, bintik, misregister) memicu sengketa keabsahan hukum dan gagal verifikasi aparat di lapangan.\n2. Daya Saing Harga (Cost Competitiveness): Inschiet tinggi mendongkrak HPP bahan baku mahal.\n3. Ketepatan Distribusi (Strict Delivery SLA): Menjamin kelancaran industri dan kas APBN.\n4. Akuntabilitas Bahan Baku (Chain of Custody & Zero Leakage): Pengawasan ketat & pemusnahan resmi lembar rusak.")
    ]
    for idx, (h, b) in enumerate(p1_items):
        p = tf_1.paragraphs[0] if idx == 0 else tf_1.add_paragraph()
        p.space_after = Pt(5)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(8.5)
        r0.font.bold = True
        r0.font.color.rgb = NAVY
        r1 = p.add_run()
        r1.text = h + "\n"
        r1.font.name = 'Arial'
        r1.font.size = Pt(8.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(7.8)
        r2.font.color.rgb = DARK_TEXT

    # Column 2: Bab 2. Skala Tender Nasional: Kontrak Pengadaan & Integritas Fiskal DJBC
    add_card(slide, 0.55 + col_w + gap, top_y, col_w, card_h, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 0.55 + col_w + gap, top_y, col_w, 0.28, "2. SKALA TENDER NASIONAL: KONTRAK & FISKAL DJBC", PURPLE)

    tb_2 = slide.shapes.add_textbox(Inches(0.55 + col_w + gap + 0.10), Inches(top_y + 0.32), Inches(col_w - 0.20), Inches(card_h - 0.38))
    tf_2 = tb_2.text_frame
    tf_2.word_wrap = True
    tf_2.margin_top = Inches(0)
    tf_2.margin_left = Inches(0)

    p2_items = [
        ("Peran Fiskal Pita Cukai (PCHT & MMEA):",
         "Pita cukai merupakan instrumen fiskal resmi dan bukti pelunasan penerimaan cukai negara yang menyumbang ratusan triliun rupiah ke kas APBN. Skala tahunan rata-rata 160.000.000 Lembar Cetak (Realisasi 2025: 177.636.930 LK — SAP ZPPRSIPPC0012)."),
        ("3 Klausul Kualitas & Kepatuhan Kontrak Tender:",
         "1. Spesifikasi Mutu Mutlak: Seluruh fitur pengaman (kertas sekuriti, tinta UV, guilloche, hologram) wajib tercetak presisi tanpa deviasi warna, register, maupun densitas. Toleransi cacat ditekan mendekati nol.\n2. Rekonsiliasi Ketat Lembar Rusak (HCTS): Setiap lembar rusak dikategorikan sebagai HCTS dan wajib dipertanggungjawabkan via berita acara pemusnahan resmi bersama pengawas.\n3. Jaminan Ketepatan SLA: Pengiriman pesanan bernilai ratusan juta lembar wajib tepat jadwal; mitigasi risiko sanksi penalti keterlambatan kontrak.")
    ]
    for idx, (h, b) in enumerate(p2_items):
        p = tf_2.paragraphs[0] if idx == 0 else tf_2.add_paragraph()
        p.space_after = Pt(5)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(8.5)
        r0.font.bold = True
        r0.font.color.rgb = PURPLE
        r1 = p.add_run()
        r1.text = h + "\n"
        r1.font.name = 'Arial'
        r1.font.size = Pt(8.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(7.8)
        r2.font.color.rgb = DARK_TEXT

    # Column 3: Bab 3. Arah Strategis & Mandat Korporasi Perum Peruri
    add_card(slide, 0.55 + 2 * (col_w + gap), top_y, col_w + 0.04, card_h, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_section_header(slide, 0.55 + 2 * (col_w + gap), top_y, col_w + 0.04, 0.28, "3. ARAH STRATEGIS & MANDAT KORPORASI PERURI", DARK_GREEN)

    tb_3 = slide.shapes.add_textbox(Inches(0.55 + 2 * (col_w + gap) + 0.10), Inches(top_y + 0.32), Inches(col_w - 0.16), Inches(card_h - 0.38))
    tf_3 = tb_3.text_frame
    tf_3.word_wrap = True
    tf_3.margin_top = Inches(0)
    tf_3.margin_left = Inches(0)

    p3_items = [
        ("Mandat Tunggal PP No. 06 Tahun 2019:",
         "Perum Peruri mengemban mandat tunggal untuk mencetak Uang Rupiah dan dokumen sekuriti negara bernilai tinggi bagi Republik Indonesia dengan jaminan keaslian dan operational excellence berstandar dunia."),
        ("3 Sasaran Strategis Direksi & Manajemen:",
         "1. Pengendalian Biaya Bahan Baku (Cost Leadership & Material Protection): Proteksi bahan baku mahal (kertas sekuriti & tinta khusus) melalui penekanan rasio pemborosan bahan (inschiet) guna menjaga HPP efisien.\n2. Pencapaian Keunggulan Operasional (Operational Excellence): Menyelaraskan proses manufaktur dengan standar ISO 9001:2015 agar kapasitas 9 mesin memenuhi kontrak tanpa pembengkakan afval.\n3. Digitalisasi Area Kerja (Smart Factory & INDI 4.0): Mentransformasikan pencatatan manual di meja mesin menjadi aliran data digital terintegrasi."),
        ("Muara Lini Operasional Terbesar:",
         "Seluruh target strategis korporasi bermuara pada lini produksi dengan volume pekerjaan terbesar: **Unit Cetak Pita Cukai**.")
    ]
    for idx, (h, b) in enumerate(p3_items):
        p = tf_3.paragraphs[0] if idx == 0 else tf_3.add_paragraph()
        p.space_after = Pt(4)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(8.5)
        r0.font.bold = True
        r0.font.color.rgb = DARK_GREEN
        r1 = p.add_run()
        r1.text = h + "\n"
        r1.font.name = 'Arial'
        r1.font.size = Pt(8.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(7.8)
        r2.font.color.rgb = DARK_TEXT


# ==============================================================================
# SLIDE 2 OF 3: NARASI 4
# (4. Realitas Lapangan di Unit Cetak Pita Cukai, 4.1 Fluktuasi 2025, 4.2 Data Silo, 4.3 Missing Link)
# ==============================================================================
def build_slide_2(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide, "02",
        "Konteks Latar Belakang (2/3): Realitas Lapangan 9 Mesin, Fluktuasi Baseline & Data Silo",
        "Paparkan profil operasional 9 mesin (3 shift 24/7), pembuktian kapabilitas Q2 vs lonjakan desain Q4, fenomena data silo meja mesin vs SAP, dan ketiadaan atribusi data."
    )

    # 4 Top Hero KPIs
    add_kpi_card(slide, 0.55, 1.68, 2.95, 0.98, "9 MESIN CETAK OFFSET", "4 KMR · 2 RYB · 3 GTO", "Pola 3 Shift 24/7 (±42 Personel Operator)", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 3.62, 1.68, 2.95, 0.98, "BASELINE INSCHIET 2025", "4,61%", "Puncak Q4: 5,11% (SAP ZPPRSIPPC0012)", RED, FILL_LIGHT_RED, BORDER_RED)
    add_kpi_card(slide, 6.69, 1.68, 2.95, 0.98, "KAPABILITAS Q2 2025", "3,97%", "Terbukti Mampu Stabil di Bawah Target 4,00%", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 9.76, 1.68, 3.02, 0.98, "DOWNTIME INVESTIGASI", "> 8 Jam / Mesin", "Pemeriksaan Spekulatif Bergilir (> 1 Shift)", ORANGE, FILL_LIGHT_AMBER, BORDER_AMBER)

    top_y = 2.76
    card_h = 4.30

    # Panel 1 (Left): Bab 4. Realitas Lapangan & Tabel 1.1 Parameter Kapasitas 2025
    col_w_1 = 3.85
    add_card(slide, 0.55, top_y, col_w_1, card_h, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_section_header(slide, 0.55, top_y, col_w_1, 0.28, "4. PARAMETER OPERASIONAL & TABEL 1.1 (2025)", NAVY)

    # Embedded Table 1.1 with 11 rows
    table_shape = slide.shapes.add_table(11, 3, Inches(0.62), Inches(top_y + 0.32), Inches(col_w_1 - 0.14), Inches(3.90))
    table = table_shape.table
    table.columns[0].width = Inches(1.80)
    table.columns[1].width = Inches(1.05)
    table.columns[2].width = Inches(0.86)

    t1_data = [
        ("Parameter Operasional", "Nilai / Data", "Sumber Valid"),
        ("Jumlah Mesin Cetak Aktif", "9 Mesin (4 KMR, 2 RYB, 3 GTO)", "Khazanah & Verif"),
        ("Pola Gilir Kerja (Shift)", "3 Shift (24 Jam)", "SOP Unit Cetak"),
        ("Personel Operator", "±42 Personel", "Seksi Cetak"),
        ("Target Volume PPIC", "160.000.000 LK", "PPIC Peruri"),
        ("Volume Aktual 2025", "177.636.930 LK", "SAP ZPPRSIPPC0012"),
        ("Inschiet Q1 2025", "4,72%", "SAP / QC Verif"),
        ("Inschiet Q2 2025", "3,97% (Stabil)", "SAP / QC Verif"),
        ("Inschiet Q3 2025", "4,64%", "SAP / QC Verif"),
        ("Inschiet Q4 2025", "5,11% (Puncak)", "SAP / QC Verif"),
        ("BASELINE INSCHIET 2025", "4,61% (Rata-rata)", "SAP ZPPRSIPPC0012")
    ]

    for r_idx, row in enumerate(t1_data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = 'Arial'
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_NAVY
                p.font.size = Pt(7.5)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            elif r_idx == 10:
                cell.fill.solid()
                cell.fill.fore_color.rgb = FILL_LIGHT_RED
                p.font.size = Pt(7.5)
                p.font.bold = True
                p.font.color.rgb = RED
                p.alignment = PP_ALIGN.CENTER if c_idx == 1 else PP_ALIGN.LEFT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else FILL_LIGHT_NAVY
                p.font.size = Pt(7)
                p.font.color.rgb = DARK_TEXT
                if c_idx == 1:
                    p.alignment = PP_ALIGN.CENTER
                    p.font.bold = True
                else:
                    p.alignment = PP_ALIGN.LEFT

    # Panel 2 (Center): Bab 4.1 Analisis Fluktuasi Baseline 2025: Kapabilitas Q2 vs Lonjakan Q4
    col_w_2 = 4.10
    add_card(slide, 4.50, top_y, col_w_2, card_h, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 4.50, top_y, col_w_2, 0.28, "4.1 ANALISIS FLUKTUASI BASELINE 2025 (%)", PURPLE)

    # Native Clustered Column Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1 (Jan-Mar)', 'Q2 (Apr-Jun)', 'Q3 (Jul-Sep)', 'Q4 (Okt-Des)', 'Avg 2025']
    chart_data.add_series('Inschiet (%)', (4.72, 3.97, 4.64, 5.11, 4.61))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(4.60), Inches(top_y + 0.32), Inches(3.90), Inches(2.15),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    plots = chart.plots[0]
    plots.has_data_labels = True
    for series in plots.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = PURPLE
        dl = series.data_labels
        dl.font.name = 'Arial'
        dl.font.size = Pt(8.5)
        dl.font.bold = True
        dl.font.color.rgb = DARK_NAVY
        dl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    # Sub-box below chart: Analisis Q2 vs Q4
    tb_q2q4 = slide.shapes.add_textbox(Inches(4.60), Inches(top_y + 2.50), Inches(3.90), Inches(1.70))
    tf_qq = tb_q2q4.text_frame
    tf_qq.word_wrap = True
    tf_qq.margin_top = Inches(0)
    tf_qq.margin_left = Inches(0)

    p_q2_h = tf_qq.paragraphs[0]
    p_q2_h.text = "• Pembuktian Kapabilitas Q2 2025 (3,97%):"
    p_q2_h.font.name = 'Arial'
    p_q2_h.font.size = Pt(8)
    p_q2_h.font.bold = True
    p_q2_h.font.color.rgb = GREEN
    p_q2_t = tf_qq.add_paragraph()
    p_q2_t.text = "Membuktikan 9 mesin & operator secara teknis mampu beroperasi di bawah toleransi 4,00% saat parameter mesin dan order stabil."
    p_q2_t.font.name = 'Arial'
    p_q2_t.font.size = Pt(7.5)
    p_q2_t.font.color.rgb = DARK_TEXT
    p_q2_t.space_after = Pt(3)

    p_q4_h = tf_qq.add_paragraph()
    p_q4_h.text = "• Lonjakan Pesanan Desain Baru Q4 (5,11%):"
    p_q4_h.font.name = 'Arial'
    p_q4_h.font.size = Pt(8)
    p_q4_h.font.bold = True
    p_q4_h.font.color.rgb = RED
    p_q4_t = tf_qq.add_paragraph()
    p_q4_t.text = "Lonjakan tajam (+1,14 pp vs Q2) terjadi saat unit menerima pesanan desain baru dalam jumlah besar. Tanpa diagnostik data harian di meja mesin, make-ready memanjang dan deviasi mutu terlambat dideteksi."
    p_q4_t.font.name = 'Arial'
    p_q4_t.font.size = Pt(7.5)
    p_q4_t.font.color.rgb = DARK_TEXT

    # Panel 3 (Right): Bab 4.2 Kesenjangan Data Silo & 4.3 Kebutaan Atribusi (The Missing Link)
    col_w_3 = 4.14
    add_card(slide, 8.70, top_y, col_w_3, card_h, FILL_LIGHT_AMBER, BORDER_AMBER)
    add_section_header(slide, 8.70, top_y, col_w_3, 0.28, "4.2 & 4.3 DATA SILO & KEBUTAAN ATRIBUSI", ORANGE)

    tb_silo = slide.shapes.add_textbox(Inches(8.80), Inches(top_y + 0.32), Inches(col_w_3 - 0.20), Inches(card_h - 0.38))
    tf_s = tb_silo.text_frame
    tf_s.word_wrap = True
    tf_s.margin_top = Inches(0)
    tf_s.margin_left = Inches(0)

    silo_items = [
        ("Pemisahan Aliran Data (Data Silo):",
         "Pencatatan di meja 9 mesin manual pada BUKU FOLIO FISIK (data jam jalan & operator terisolasi, baru direkap berbulan-bulan kemudian) ≠ Hasil verifikasi disortir 1–2 hari lalu diinput ke SAP ZPPRSIPPC0012 sebagai RINGKASAN KERUSAKAN GLOBAL tanpa nomor mesin & shift."),
        ("Kebutaan Atribusi (The Missing Link):",
         "Mengetahui jenis cacat global (blobor, bintik, misregister) terbukti tidak cukup tanpa tahu di mesin mana dan kondisi operasional shift mana yang memicunya."),
        ("3 Implikasi Titik Buta Lapangan:",
         "1. Investigasi Spekulatif (> 8 Jam): Teknisi memeriksa 9 mesin coba-coba (trial-and-error).\n2. Dilema Komponen Mesin vs Shift: Penurunan performa fisik (rol karet mengeras/licin, blanket turun elastisitas, penjepit silinder kendur) vs variasi make-ready / kelelahan Shift Malam (23.00–07.00 WIB).\n3. Evaluasi Operator Tertunda: Umpan balik harian ±42 operator terputus.")
    ]
    for idx, (h, b) in enumerate(silo_items):
        p = tf_s.paragraphs[0] if idx == 0 else tf_s.add_paragraph()
        p.space_after = Pt(4)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(8)
        r0.font.bold = True
        r0.font.color.rgb = RED if idx > 0 else ORANGE
        r1 = p.add_run()
        r1.text = h + "\n"
        r1.font.name = 'Arial'
        r1.font.size = Pt(8)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = DARK_TEXT


# ==============================================================================
# SLIDE 3 OF 3: NARASI 5 & 6
# (5. Skala Dampak Finansial & Matriks Inaction, 6. Urgensi Solusi DSS SIRINE 4.0 & Realisasi S1)
# ==============================================================================
def build_slide_3(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide, "03",
        "Konteks Latar Belakang (3/3): Simulasi Finansial, Risiko Pembiaran & Urgensi Solusi",
        "Sajikan simulasi beban finansial baseline (Rp 24,56 M), valuasi 1% inschiet, matriks 5 pilar Cost of Inaction, arsitektur solusi DSS SIRINE 4.0, dan kertas kerja realisasi S1 2026."
    )

    # 4 Top Hero KPIs
    add_kpi_card(slide, 0.55, 1.68, 2.95, 0.98, "BEBAN BIAYA BASELINE*", "Rp 24,56 Miliar*", "8.189.062 LK Rusak/Thn (@ Rp 3.000 / LK)", RED, FILL_LIGHT_RED, BORDER_RED)
    add_kpi_card(slide, 3.62, 1.68, 2.95, 0.98, "VALUASI TIAP 1% INSCHIET", "Rp 4,80 M – 5,33 M", "Penyelamatan 1,60M – 1,78M Lembar/Tahun", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 6.69, 1.68, 2.95, 0.98, "REALISASI INSCHIET S1", "3,89% (Q2: 3,33%)", "Turun -1,28 pp (-27,77%) vs Baseline 4,61%", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 9.76, 1.68, 3.02, 0.98, "PENGHEMATAN BIAYA S1", "Rp 2,23 Miliar", "743.234 LK Diselamatkan | Proyeksi Rp 6,82 M/Thn", DARK_GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)

    # Top Flow Banner: Bab 6 Urgensi Solusi & Arsitektur Integrasi 3 Titik Aliran Data
    top_flow_w = 12.23
    add_card(slide, 0.55, 2.76, top_flow_w, 0.88, FILL_LIGHT_PURPLE, BORDER_PURPLE)

    tb_flow = slide.shapes.add_textbox(Inches(0.65), Inches(2.78), Inches(top_flow_w - 0.20), Inches(0.82))
    tf_f = tb_flow.text_frame
    tf_f.word_wrap = True
    tf_f.margin_top = Inches(0)
    tf_f.margin_left = Inches(0)

    p_f1 = tf_f.paragraphs[0]
    r_f0 = p_f1.add_run()
    r_f0.text = "6. KESIMPULAN & URGENSI SOLUSI INOVASI (DSS SIRINE 4.0) — INTEGRASI 3 TITIK ALIRAN DATA:  "
    r_f0.font.name = 'Arial'
    r_f0.font.size = Pt(8)
    r_f0.font.bold = True
    r_f0.font.color.rgb = PURPLE
    r_f1 = p_f1.add_run()
    r_f1.text = "Data Transaksi Meja Mesin (< 30 dtk)  ⇄  Modul SAP Production Order (ZPPRSIPPC0012)  ⇄  Hasil Sortir Verifikasi Mutu (HCTS)\n"
    r_f1.font.name = 'Arial'
    r_f1.font.size = Pt(8)
    r_f1.font.bold = True
    r_f1.font.color.rgb = NAVY

    p_f2 = tf_f.add_paragraph()
    r_f2 = p_f2.add_run()
    r_f2.text = "RUNTUTAN DATA GRANULAR & TINDAKAN PRESKRIPTIF:  "
    r_f2.font.name = 'Arial'
    r_f2.font.size = Pt(8)
    r_f2.font.bold = True
    r_f2.font.color.rgb = DARK_GREEN
    r_f3 = p_f2.add_run()
    r_f3.text = "Nomor PO  ➔  Nomor Mesin (9 Mesin)  ➔  Pola Gilir (Shift 1/2/3)  ➔  Tim Operator  ➔  Kategori Cacat Cetak  ➔  Tindakan Korektif Cepat"
    r_f3.font.name = 'Arial'
    r_f3.font.size = Pt(8)
    r_f3.font.color.rgb = DARK_TEXT

    # Bottom 3 Panels (Sections 5.1/5.2, 5.3, 6 Table 1.3)
    bot_y = 3.72
    bot_h = 3.34
    col_w_b1 = 3.85
    col_w_b2 = 3.85
    col_w_b3 = 4.25
    gap_b = 0.14

    # Bottom-Left Panel: Bab 5.1 & 5.2 Simulasi Skala Dampak Finansial & Valuasi 1%
    add_card(slide, 0.55, bot_y, col_w_b1, bot_h, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_section_header(slide, 0.55, bot_y, col_w_b1, 0.25, "5.1 & 5.2 SIMULASI FINANSIAL & VALUASI 1%", NAVY)

    tb_fin = slide.shapes.add_textbox(Inches(0.65), Inches(bot_y + 0.28), Inches(col_w_b1 - 0.20), Inches(bot_h - 0.35))
    tf_fin = tb_fin.text_frame
    tf_fin.word_wrap = True
    tf_fin.margin_top = Inches(0)
    tf_fin.margin_left = Inches(0)

    fin_items = [
        ("Skenario A (Standar 160 Juta LK):",
         "160.000.000 LK × 4,61% = 7.376.000 LK Rusak → 7.376.000 × Rp 3.000* = Rp 22,13 Miliar / Tahun (Rp 1,84 Miliar / Bulan)."),
        ("Skenario B (Aktual 177,6 Juta LK 2025):",
         "177.636.930 LK × 4,61% = 8.189.062 LK Rusak → 8.189.062 × Rp 3.000* = Rp 24,56 Miliar / Tahun (Rp 2,05 Miliar / Bulan)."),
        ("Valuasi Penurunan Tiap 1,00% Inschiet:",
         "• Standar 160M LK: 1.600.000 LK diselamatkan = Efisiensi Rp 4,80 Miliar / Tahun.\n• Aktual 177,6M LK: 1.776.369 LK diselamatkan = Efisiensi Rp 5,33 Miliar / Tahun.\n*Catatan: Nilai Rp 3.000/LK adalah estimasi biaya cetak untuk simulasi cost avoidance (bukan HPP rahasia Peruri).")
    ]
    for idx, (h, b) in enumerate(fin_items):
        p = tf_fin.paragraphs[0] if idx == 0 else tf_fin.add_paragraph()
        p.space_after = Pt(3)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(7.5)
        r0.font.bold = True
        r0.font.color.rgb = NAVY if idx < 2 else GREEN
        r1 = p.add_run()
        r1.text = h + "\n"
        r1.font.name = 'Arial'
        r1.font.size = Pt(7.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(7)
        r2.font.color.rgb = DARK_TEXT

    # Bottom-Center Panel: Bab 5.3 Matriks Evaluasi 5 Pilar Cost of Inaction (Tabel 1.2)
    add_card(slide, 0.55 + col_w_b1 + gap_b, bot_y, col_w_b2, bot_h, FILL_LIGHT_RED, BORDER_RED)
    add_section_header(slide, 0.55 + col_w_b1 + gap_b, bot_y, col_w_b2, 0.25, "5.3 MATRIKS 5 PILAR RISIKO PEMBIARAN (INACTION)", RED)

    tb_in = slide.shapes.add_textbox(Inches(0.55 + col_w_b1 + gap_b + 0.10), Inches(bot_y + 0.28), Inches(col_w_b2 - 0.20), Inches(bot_h - 0.35))
    tf_in = tb_in.text_frame
    tf_in.word_wrap = True
    tf_in.margin_top = Inches(0)
    tf_in.margin_left = Inches(0)

    in_items = [
        ("1. Biaya (Cost) [KRITIS]:", "Akumulasi pemborosan bahan baku kertas sekuriti & tinta mencapai Rp 22,13 M s.d. Rp 24,56 Miliar/tahun & pembengkakan tambah cetak."),
        ("2. Mutu (Quality) [TINGGI]:", "Inschiet berfluktuasi hingga 5,11% akibat penanganan suku cadang mesin yang reaktif & coba-coba."),
        ("3. Kepatuhan [TINGGI]:", "Pencatatan manual menyulitkan audit riwayat lot produksi ISO 9001:2015."),
        ("4. K3L / ESG [SEDANG]:", "Timbulan afval 7,37M–8,18M LK/tahun (±60–65 Ton kertas) & kelelahan shift malam."),
        ("5. Layanan (SLA) [TINGGI]:", "Siklus cetak ulang lambat menunda serah terima pita cukai ke DJBC & ancaman denda penalti.")
    ]
    for idx, (h, b) in enumerate(in_items):
        p = tf_in.paragraphs[0] if idx == 0 else tf_in.add_paragraph()
        p.space_after = Pt(2)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(7.5)
        r0.font.bold = True
        r0.font.color.rgb = RED
        r1 = p.add_run()
        r1.text = h + " "
        r1.font.name = 'Arial'
        r1.font.size = Pt(7.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(7)
        r2.font.color.rgb = DARK_TEXT

    # Bottom-Right Panel: Bab 6. Tabel 1.3 Kertas Kerja Realisasi S1 2026 & Capaian Nyata
    add_card(slide, 0.55 + col_w_b1 + gap_b + col_w_b2 + gap_b, bot_y, col_w_b3, bot_h, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_section_header(slide, 0.55 + col_w_b1 + gap_b + col_w_b2 + gap_b, bot_y, col_w_b3, 0.25, "6. TABEL 1.3 REALISASI S1 2026 & VALUE CREATION", DARK_GREEN)

    # Embedded Table 1.3
    table_shape_3 = slide.shapes.add_table(4, 8, Inches(0.55 + col_w_b1 + gap_b + col_w_b2 + gap_b + 0.06), Inches(bot_y + 0.28), Inches(col_w_b3 - 0.12), Inches(1.50))
    table_3 = table_shape_3.table
    table_3.columns[0].width = Inches(1.10)
    table_3.columns[1].width = Inches(0.48)
    table_3.columns[2].width = Inches(0.42)
    table_3.columns[3].width = Inches(0.46)
    table_3.columns[4].width = Inches(0.46)
    table_3.columns[5].width = Inches(0.46)
    table_3.columns[6].width = Inches(0.46)
    table_3.columns[7].width = Inches(0.45)

    t3_headers = ["Periode", "Vol (n)", "Inschiet", "Deviasi", "Tgt 4,61%", "Realisasi", "Saved", "Saving*"]
    t3_rows = [
        ("Q1 2026", "57,38M", "4,34%", "-0,27 pp", "2,64M lb", "2,49M lb", "154,9K LK", "Rp 464,8Jt"),
        ("Q2 2026", "45,96M", "3,33%", "-1,28 pp", "2,11M lb", "1,53M lb", "588,3K LK", "Rp 1,76 M"),
        ("TOTAL S1", "103,3M", "3,89%", "-0,72 pp", "4,76M lb", "4,02M lb", "743,2K LK", "Rp 2,23 M")
    ]

    for c_idx, h_text in enumerate(t3_headers):
        cell = table_3.cell(0, c_idx)
        cell.margin_left = Inches(0.02)
        cell.margin_right = Inches(0.02)
        cell.margin_top = Inches(0.01)
        cell.margin_bottom = Inches(0.01)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.name = 'Arial'
        p.font.size = Pt(6.5)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for r_idx, r_data in enumerate(t3_rows):
        for c_idx, val in enumerate(r_data):
            cell = table_3.cell(r_idx + 1, c_idx)
            cell.margin_left = Inches(0.02)
            cell.margin_right = Inches(0.02)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            cell.fill.solid()
            cell.fill.fore_color.rgb = FILL_LIGHT_GREEN if r_idx == 2 else (WHITE if r_idx % 2 == 0 else FILL_LIGHT_NAVY)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = 'Arial'
            p.font.size = Pt(6.5)
            p.font.color.rgb = DARK_GREEN if r_idx == 2 else DARK_TEXT
            p.font.bold = (r_idx == 2 or c_idx in [0, 6, 7])
            p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT

    # Capaian Nyata & Value Creation below Table 1.3
    tb_val3 = slide.shapes.add_textbox(Inches(0.55 + col_w_b1 + gap_b + col_w_b2 + gap_b + 0.08), Inches(bot_y + 1.82), Inches(col_w_b3 - 0.16), Inches(1.45))
    tf_v3 = tb_val3.text_frame
    tf_v3.word_wrap = True
    tf_v3.margin_top = Inches(0)
    tf_v3.margin_left = Inches(0)

    val_bullets = [
        ("1. Inschiet Tembus 3,33% (Q2):", "Penurunan -1,28 pp (-27,77%) menyelamatkan 743.234 lembar kertas sekuriti."),
        ("2. Efisiensi Finansial S1:", "Saving riil Rp 2,23 Miliar (proyeksi Rp 6,82 M/thn) dengan biaya dev Rp 0,-."),
        ("3. Reduksi Waktu Henti:", "Pemeriksaan mesin turun dari > 8 jam menjadi < 2–4 jam (efisiensi 50%–75%)."),
        ("4. Jaminan Kepatuhan SLA:", "Kepatuhan 100% pengiriman DJBC Kemenkeu RI & penguatan tender nasional.")
    ]
    for idx, (h, b) in enumerate(val_bullets):
        p = tf_v3.paragraphs[0] if idx == 0 else tf_v3.add_paragraph()
        p.space_after = Pt(2)
        r0 = p.add_run()
        r0.text = "✔ "
        r0.font.name = 'Arial'
        r0.font.size = Pt(7)
        r0.font.bold = True
        r0.font.color.rgb = GREEN
        r1 = p.add_run()
        r1.text = h + " "
        r1.font.name = 'Arial'
        r1.font.size = Pt(7)
        r1.font.bold = True
        r1.font.color.rgb = DARK_NAVY
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(6.8)
        r2.font.color.rgb = DARK_TEXT


# ==============================================================================
# MAIN GENERATOR
# ==============================================================================
def main():
    prs, blank_layout = create_presentation()

    print("Building Slide 1 of 3: Narasi 1, 2, 3 (Pasar Global, Tender DJBC, Mandat Peruri)...")
    build_slide_1(prs, blank_layout)

    print("Building Slide 2 of 3: Narasi 4 (Realitas Lapangan 9 Mesin, Fluktuasi 2025, Data Silo & Kebutaan Atribusi)...")
    build_slide_2(prs, blank_layout)

    print("Building Slide 3 of 3: Narasi 5 & 6 (Simulasi Finansial, Matriks Inaction, Urgensi DSS SIRINE 4.0 & Realisasi S1)...")
    build_slide_3(prs, blank_layout)

    output_path = "Presentasi_Risalah_Latar_Belakang_IAKA_2026.pptx"
    prs.save(output_path)
    print(f"\n[SUCCESS] Master 3-Slide Presentation for Risalah successfully generated: {output_path}")

    prs.save("Presentasi_Risalah_Latar_Belakang_IAKA_2026_3_Slides.pptx")

if __name__ == "__main__":
    main()
