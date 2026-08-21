import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
import sys
import os

# ==========================================
# PALETTE & DESIGN CONSTANTS (Peruri Theme)
# ==========================================
NAVY = RGBColor(32, 40, 143)         # #20288F - Primary Brand Navy
PURPLE = RGBColor(106, 44, 201)      # #6A2CC9 - Accent Purple
DARK_NAVY = RGBColor(26, 35, 126)    # #1A237E - Deep Navy
GREEN = RGBColor(46, 125, 50)        # #2E7D32 - Success Green
DARK_GREEN = RGBColor(27, 94, 32)    # #1B5E20
RED = RGBColor(198, 40, 40)          # #C62828 - Alert Red
ORANGE = RGBColor(230, 81, 0)        # #E65100 - Warning Amber
DARK_TEXT = RGBColor(42, 42, 61)     # #2A2A3D - Body Text Charcoal
MUTED_TEXT = RGBColor(110, 110, 130) # #6E6E82 - Muted Text Gray
WHITE = RGBColor(255, 255, 255)

# Container Fills
FILL_LIGHT_NAVY = RGBColor(244, 247, 254)   # #F4F7FE
FILL_LIGHT_PURPLE = RGBColor(248, 246, 254) # #F8F6FE
FILL_LIGHT_GREEN = RGBColor(232, 245, 233)  # #E8F5E9
FILL_LIGHT_AMBER = RGBColor(255, 249, 230)  # #FFF9E6
FILL_LIGHT_RED = RGBColor(255, 245, 245)    # #FFF5F5

# Container Borders
BORDER_NAVY = RGBColor(210, 220, 245)       # #D2DCF5
BORDER_PURPLE = RGBColor(217, 206, 247)     # #D9CEF7
BORDER_GREEN = RGBColor(165, 214, 167)      # #A5D6A7
BORDER_AMBER = RGBColor(255, 224, 130)      # #FFE082
BORDER_RED = RGBColor(255, 205, 210)        # #FFCDD2

def create_presentation():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    return prs, blank_layout

def add_header(slide, badge_num, title_text, purpose_text):
    if os.path.exists('wave_ornament.png'):
        slide.shapes.add_picture('wave_ornament.png', Inches(-0.118), Inches(-0.196), Inches(1.584), Inches(0.891))
    
    if os.path.exists('peruri_logo.png'):
        slide.shapes.add_picture('peruri_logo.png', Inches(12.09), Inches(0.013), Inches(1.247), Inches(0.694))

    # Badge Box
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.55), Inches(0.75), Inches(0.68))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PURPLE
    badge.line.color.rgb = PURPLE
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = badge_num
    p.font.name = 'Arial'
    p.font.size = Pt(18 if len(badge_num) <= 2 else 15)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    tb_title = slide.shapes.add_textbox(Inches(1.42), Inches(0.50), Inches(10.50), Inches(0.75))
    tf = tb_title.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Purpose Subtitle
    tb_purpose = slide.shapes.add_textbox(Inches(0.55), Inches(1.30), Inches(12.23), Inches(0.32))
    tf = tb_purpose.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Tujuan bagian ini:  "
    r1.font.name = 'Arial'
    r1.font.size = Pt(10)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = purpose_text
    r2.font.name = 'Arial'
    r2.font.size = Pt(10)
    r2.font.color.rgb = DARK_TEXT

    # Footer
    tb_foot = slide.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(12.23), Inches(0.22))
    tf = tb_foot.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = "IAKA 2026 — Kerangka Presentasi Peserta  ·  DSS SIRINE 4.0 Unit Cetak Pita Cukai (Perum Peruri)"
    p.font.name = 'Arial'
    p.font.size = Pt(8)
    p.font.color.rgb = MUTED_TEXT

def add_kpi_card(slide, left, top, width, height, title, value, subtext, color, fill_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(left + 0.08), Inches(top + 0.06), Inches(width - 0.16), Inches(height - 0.10))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)

    p1 = tf.paragraphs[0]
    p1.text = title.upper()
    p1.font.name = 'Arial'
    p1.font.size = Pt(7.5)
    p1.font.bold = True
    p1.font.color.rgb = color

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.name = 'Arial'
    p2.font.size = Pt(14.5)
    p2.font.bold = True
    p2.font.color.rgb = color

    p3 = tf.add_paragraph()
    p3.text = subtext
    p3.font.name = 'Arial'
    p3.font.size = Pt(7.5)
    p3.font.color.rgb = DARK_TEXT

def add_card(slide, left, top, width, height, fill_color, border_color, border_width=1.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(border_width)
    return card

def add_section_header(slide, left, top, width, height, title, bg_color):
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = bg_color
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(8.5)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    return hdr


# ==============================================================================
# SLIDE 1 OF 3: KONTEKS STRATEGIS, SKALA TENDER & PARAMETER 9 MESIN
# (Sections 1, 2, 3, 4 of Risalah)
# ==============================================================================
def build_slide_1(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide, "01",
        "Konteks Strategis Percetakan Sekuriti, Mandat Tender DJBC & Realitas Operasional 9 Mesin",
        "Bangun konteks pasar sekuriti global (Intergraf/WCO), komitmen fiskal tender DJBC (177,6M LK), arah korporasi PP 06/2019, dan parameter kapasitas 9 mesin offset."
    )

    # 4 Top Hero KPIs
    add_kpi_card(slide, 0.55, 1.68, 2.95, 0.98, "REALISASI PESANAN 2025", "177.636.930 LK", "Target PPIC: 160M LK | SAP ZPPRSIPPC0012", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 3.62, 1.68, 2.95, 0.98, "9 MESIN CETAK OFFSET", "4 KMR · 2 RYB · 3 GTO", "Pola 3 Shift 24/7 (Pagi, Sore, Malam) | ±42 Personel", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 6.69, 1.68, 2.95, 0.98, "BASELINE INSCHIET 2025", "4,61%", "Q1: 4,72% · Q2: 3,97% · Q3: 4,64% · Q4: 5,11%", RED, FILL_LIGHT_RED, BORDER_RED)
    add_kpi_card(slide, 9.76, 1.68, 3.02, 0.98, "PERAN FISKAL APBN", "Ratusan Triliun", "Penerimaan Cukai Negara (PCHT & MMEA)", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)

    # 3 Large Vertical Panels (Covering Sections 1, 2, 3, 4)
    col_w = 3.98
    gap = 0.14
    top_y = 2.76
    card_h = 4.30

    # Panel 1: Dinamika Pasar Global & Mandat Korporasi Peruri (Sections 1 & 3)
    add_card(slide, 0.55, top_y, col_w, card_h, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_section_header(slide, 0.55, top_y, col_w, 0.28, "1. PASAR GLOBAL & ARAH STRATEGIS PERURI", NAVY)

    tb_1 = slide.shapes.add_textbox(Inches(0.65), Inches(top_y + 0.32), Inches(col_w - 0.20), Inches(card_h - 0.38))
    tf_1 = tb_1.text_frame
    tf_1.word_wrap = True
    tf_1.margin_top = Inches(0)
    tf_1.margin_left = Inches(0)

    p1_items = [
        ("Standar Internasional (Intergraf & WCO):",
         "Industri percetakan sekuriti tinggi menuntut fitur pengamanan berlapis (multi-layer security): kertas serat pengaman tak kasat mata, tinta sekuriti berpendar UV, ornamen guilloche, microtext rapat, dan pita hologram presisi mikron guna mencegah pemalsuan (anti-counterfeiting)."),
        ("4 Kriteria Keunggulan Bersaing:",
         "1. Zero-Defect Assurance: Cacat blobor, bintik, atau misregister memicu sengketa keabsahan hukum dokumen.\n2. Cost Competitiveness: Inschiet tinggi mendongkrak biaya pokok produksi bahan baku mahal.\n3. Strict Delivery SLA: Kepastian jadwal pasokan nasional ketat.\n4. Chain of Custody & Zero Leakage: Wajib pemusnahan resmi lembar rusak."),
        ("Mandat PP No. 06/2019 & 3 Sasaran Korporasi:",
         "Peruri mengemban mandat mencetak dokumen sekuriti negara. Direksi menetapkan: (1) Cost Leadership & Material Protection; (2) Operational Excellence (ISO 9001:2015); (3) Digitalisasi Area Kerja (INDI 4.0). Seluruh sasaran bermuara di Unit Cetak Pita Cukai.")
    ]
    for idx, (h, b) in enumerate(p1_items):
        p = tf_1.paragraphs[0] if idx == 0 else tf_1.add_paragraph()
        p.space_after = Pt(4)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(8)
        r0.font.bold = True
        r0.font.color.rgb = NAVY
        r1 = p.add_run()
        r1.text = h + " "
        r1.font.name = 'Arial'
        r1.font.size = Pt(8)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = DARK_TEXT

    # Panel 2: Skala Tender Nasional & Klausul DJBC Kemenkeu RI (Section 2)
    add_card(slide, 0.55 + col_w + gap, top_y, col_w, card_h, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 0.55 + col_w + gap, top_y, col_w, 0.28, "2. KOMITMEN TENDER NASIONAL DJBC KEMENKEU RI", PURPLE)

    tb_2 = slide.shapes.add_textbox(Inches(0.55 + col_w + gap + 0.10), Inches(top_y + 0.32), Inches(col_w - 0.20), Inches(card_h - 0.38))
    tf_2 = tb_2.text_frame
    tf_2.word_wrap = True
    tf_2.margin_top = Inches(0)
    tf_2.margin_left = Inches(0)

    p2_items = [
        ("Peran Fiskal Pita Cukai (PCHT & MMEA):",
         "Pita cukai merupakan instrumen fiskal resmi dan bukti pelunasan cukai negara bernilai ratusan triliun rupiah ke kas APBN. Peran vital ini melandasi klausul kontrak pengadaan yang sangat ketat."),
        ("1. Spesifikasi Mutu Mutlak (Nol Toleransi):",
         "Seluruh fitur pengamanan fisik (kertas sekuriti, tinta UV, guilloche, hologram) wajib tercetak presisi tanpa deviasi warna, register, maupun densitas. Cacat fisik berisiko menggagalkan verifikasi keaslian oleh aparat DJBC di lapangan."),
        ("2. Rekonsiliasi Ketat Lembar Rusak (HCTS):",
         "Setiap lembar rusak dikategorikan sebagai Hasil Cetak Tidak Sempurna (HCTS) dan wajib dipertanggungjawabkan melalui berita acara pemusnahan resmi bersama pengawas. Tingginya inschiet menambah beban audit fisik."),
        ("3. Jaminan Service Level Agreement (SLA):",
         "Pengiriman pesanan bernilai ratusan juta lembar wajib tepat jadwal guna menjamin kelancaran pabrik rokok/MMEA dan arus kas APBN. Proses cetak ulang yang panjang berisiko terkena sanksi penalti kontrak.")
    ]
    for idx, (h, b) in enumerate(p2_items):
        p = tf_2.paragraphs[0] if idx == 0 else tf_2.add_paragraph()
        p.space_after = Pt(4)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(8)
        r0.font.bold = True
        r0.font.color.rgb = PURPLE
        r1 = p.add_run()
        r1.text = h + " "
        r1.font.name = 'Arial'
        r1.font.size = Pt(8)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = DARK_TEXT

    # Panel 3: Realitas Lapangan & Tabel Parameter Kapasitas 2025 (Section 4, Tabel 1.1)
    add_card(slide, 0.55 + 2 * (col_w + gap), top_y, col_w + 0.04, card_h, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_section_header(slide, 0.55 + 2 * (col_w + gap), top_y, col_w + 0.04, 0.28, "3. PARAMETER OPERASIONAL 9 MESIN & BASELINE 2025", DARK_GREEN)

    # Embedded Table 1.1 inside Panel 3
    table_shape = slide.shapes.add_table(11, 3, Inches(0.55 + 2 * (col_w + gap) + 0.06), Inches(top_y + 0.34), Inches(col_w - 0.08), Inches(3.88))
    table = table_shape.table
    table.columns[0].width = Inches(1.85)
    table.columns[1].width = Inches(1.10)
    table.columns[2].width = Inches(0.95)

    t1_data = [
        ("Parameter Operasional", "Nilai / Data", "Sumber Valid"),
        ("Jumlah Mesin Cetak Aktif", "9 Mesin (4 KMR, 2 RYB, 3 GTO)", "Khazanah & Verif"),
        ("Pola Gilir Kerja (Shift)", "3 Shift (24 Jam)", "SOP Unit Cetak"),
        ("Personel Operator", "±42 Personel", "Seksi Cetak"),
        ("Target Volume PPIC", "160.000.000 LK", "PPIC Peruri"),
        ("Volume Aktual 2025", "177.636.930 LK", "SAP ZPPRSIPPC0012"),
        ("Inschiet Q1 2025", "4,72%", "SAP / QC Verif"),
        ("Inschiet Q2 2025", "3,97% (Terkendali)", "SAP / QC Verif"),
        ("Inschiet Q3 2025", "4,64%", "SAP / QC Verif"),
        ("Inschiet Q4 2025", "5,11% (Puncak)", "SAP / QC Verif"),
        ("BASELINE INSCHIET 2025", "4,61% (Rata-rata)", "SAP ZPPRSIPPC0012")
    ]

    for r_idx, row in enumerate(t1_data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = 'Arial'
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_NAVY
                p.font.size = Pt(7.5)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            elif r_idx == 10:
                cell.fill.solid()
                cell.fill.fore_color.rgb = FILL_LIGHT_RED
                p.font.size = Pt(7.5)
                p.font.bold = True
                p.font.color.rgb = RED
                p.alignment = PP_ALIGN.CENTER if c_idx == 1 else PP_ALIGN.LEFT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else FILL_LIGHT_NAVY
                p.font.size = Pt(7)
                p.font.color.rgb = DARK_TEXT
                if c_idx == 1:
                    p.alignment = PP_ALIGN.CENTER
                    p.font.bold = True
                else:
                    p.alignment = PP_ALIGN.LEFT


# ==============================================================================
# SLIDE 2 OF 3: IDENTIFIKASI MASALAH, DATA SILO, SIMULASI FINANSIAL & MATRIKS RISIKO
# (Sections 4.1, 4.2, 4.3, 5 of Risalah)
# ==============================================================================
def build_slide_2(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide, "02",
        "Identifikasi Masalah: Fluktuasi Inschiet, Fenomena Data Silo, Simulasi Finansial & Risiko",
        "Analisis tren fluktuasi baseline 2025, dekonstruksi pemisahan aliran data (data silo), simulasi beban finansial Rp 24,56 Miliar, dan matriks 5 pilar Cost of Inaction."
    )

    # 4 Top Hero KPIs
    add_kpi_card(slide, 0.55, 1.68, 2.95, 0.98, "BEBAN BASELINE AKTUAL", "Rp 24,56 Miliar*", "8.189.062 LK Rusak/Thn (@ Rp 3.000 / LK)", RED, FILL_LIGHT_RED, BORDER_RED)
    add_kpi_card(slide, 3.62, 1.68, 2.95, 0.98, "VALUASI TIAP 1% INSCHIET", "Rp 4,80 M – 5,33 M", "Penyelamatan 1,60M – 1,78M Lembar/Thn", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 6.69, 1.68, 2.95, 0.98, "DOWNTIME PERBAIKAN", "> 8 Jam / Mesin", "Investigasi spekulatif bergilir (> 1 shift)", ORANGE, FILL_LIGHT_AMBER, BORDER_AMBER)
    add_kpi_card(slide, 9.76, 1.68, 3.02, 0.98, "LONJAKAN DESAIN Q4", "5,11% (+1,14 pp)", "Anomali akhir tahun tanpa diagnostik meja mesin", RED, FILL_LIGHT_RED, BORDER_RED)

    # Grid of 4 Dense Quadrants
    quad_w = 6.04
    quad_h = 2.12
    top_y1 = 2.76
    top_y2 = 4.96

    # Quadrant 1 (Top-Left): Native Chart Baseline Inschiet 2025 & Analisis Q2 vs Q4 (Section 4.1)
    add_card(slide, 0.55, top_y1, quad_w, quad_h, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 0.55, top_y1, quad_w, 0.25, "1. TREN INSCHIET BASELINE 2025 (%) & PEMBUKTIAN KAPABILITAS (Q2 VS Q4)", PURPLE)

    # Native Clustered Column Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1 (Jan-Mar)', 'Q2 (Apr-Jun)', 'Q3 (Jul-Sep)', 'Q4 (Okt-Des)', 'Avg 2025']
    chart_data.add_series('Inschiet (%)', (4.72, 3.97, 4.64, 5.11, 4.61))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.65), Inches(top_y1 + 0.28), Inches(3.30), Inches(1.78),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    plots = chart.plots[0]
    plots.has_data_labels = True
    for series in plots.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = PURPLE
        dl = series.data_labels
        dl.font.name = 'Arial'
        dl.font.size = Pt(8)
        dl.font.bold = True
        dl.font.color.rgb = DARK_NAVY
        dl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    # Side callout Q2 vs Q4
    tb_q = slide.shapes.add_textbox(Inches(4.02), Inches(top_y1 + 0.28), Inches(2.50), Inches(1.78))
    tf_q = tb_q.text_frame
    tf_q.word_wrap = True
    tf_q.margin_top = Inches(0)
    tf_q.margin_left = Inches(0)
    p_q1 = tf_q.paragraphs[0]
    p_q1.text = "• Q2 (3,97%): Bukti 9 mesin mampu stabil < 4,00% saat order repetitif."
    p_q1.font.name = 'Arial'
    p_q1.font.size = Pt(7.5)
    p_q1.font.bold = True
    p_q1.font.color.rgb = GREEN
    p_q1.space_after = Pt(2)
    p_q2 = tf_q.add_paragraph()
    p_q2.text = "• Q4 (5,11%): Lonjakan saat pesanan desain baru masuk tanpa diagnostik harian di meja mesin sehingga make-ready lama dan cacat telat dideteksi."
    p_q2.font.name = 'Arial'
    p_q2.font.size = Pt(7)
    p_q2.font.color.rgb = DARK_TEXT

    # Quadrant 2 (Top-Right): Anatomi Pemisahan Data (Data Silo) & Kebutaan Atribusi (Sections 4.2 & 4.3)
    add_card(slide, 6.73, top_y1, quad_w + 0.05, quad_h, FILL_LIGHT_AMBER, BORDER_AMBER)
    add_section_header(slide, 6.73, top_y1, quad_w + 0.05, 0.25, "2. ANATOMI DATA SILO & KEBUTAAN ATRIBUSI (THE MISSING LINK)", ORANGE)

    tb_ds = slide.shapes.add_textbox(Inches(6.83), Inches(top_y1 + 0.28), Inches(quad_w - 0.15), Inches(quad_h - 0.32))
    tf_ds = tb_ds.text_frame
    tf_ds.word_wrap = True
    tf_ds.margin_top = Inches(0)
    tf_ds.margin_left = Inches(0)

    ds_items = [
        ("Meja Mesin vs SAP Kantor:",
         "Pencatatan di 9 meja mesin manual pada BUKU FOLIO FISIK (data terisolasi, baru direkap berbulan-bulan kemudian) ≠ Hasil verifikasi QC diinput ke SAP ZPPRSIPPC0012 sebagai RINGKASAN GLOBAL tingkat unit tanpa nomor mesin, PO, dan shift."),
        ("Implikasi Kebutaan Atribusi (Attribution Blindness):",
         "1. Investigasi Spekulatif (> 8 Jam): Saat cacat blobor/bintik naik, teknisi memeriksa 9 mesin coba-coba (trial-and-error).\n2. Dilema Komponen vs Operator: Sulit memisahkan penurunan fisik mesin (rol karet mengeras/licin, blanket kempes, gripper kendur) vs variasi make-ready / kelelahan Shift Malam (23.00–07.00 WIB).\n3. Evaluasi Operator Tertunda: Umpan balik harian bagi ±42 operator terputus.")
    ]
    for idx, (h, b) in enumerate(ds_items):
        p = tf_ds.paragraphs[0] if idx == 0 else tf_ds.add_paragraph()
        p.space_after = Pt(2)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(7.5)
        r0.font.bold = True
        r0.font.color.rgb = RED
        r1 = p.add_run()
        r1.text = h + " "
        r1.font.name = 'Arial'
        r1.font.size = Pt(7.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(7)
        r2.font.color.rgb = DARK_TEXT

    # Quadrant 3 (Bottom-Left): Simulasi Skala Dampak Finansial (Section 5.1 & 5.2)
    add_card(slide, 0.55, top_y2, quad_w, quad_h - 0.05, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_section_header(slide, 0.55, top_y2, quad_w, 0.25, "3. SIMULASI SKALA DAMPAK FINANSIAL BASELINE 2025 & VALUASI 1% INSCHIET", NAVY)

    tb_fin = slide.shapes.add_textbox(Inches(0.65), Inches(top_y2 + 0.28), Inches(quad_w - 0.20), Inches(quad_h - 0.35))
    tf_fin = tb_fin.text_frame
    tf_fin.word_wrap = True
    tf_fin.margin_top = Inches(0)
    tf_fin.margin_left = Inches(0)

    fin_items = [
        ("Skenario A (Standar 160 Juta LK):",
         "160.000.000 LK × 4,61% = 7.376.000 LK Rusak → 7.376.000 × Rp 3.000* = Rp 22,13 Miliar / Tahun (Rp 1,84 M/Bulan)."),
        ("Skenario B (Realisasi 177,6 Juta LK):",
         "177.636.930 LK × 4,61% = 8.189.062 LK Rusak → 8.189.062 × Rp 3.000* = Rp 24,56 Miliar / Tahun (Rp 2,05 M/Bulan)."),
        ("Valuasi Penghematan Tiap 1,00% Penurunan Inschiet:",
         "• Standar 160M LK: 1.600.000 LK diselamatkan = Efisiensi Rp 4,80 Miliar / Tahun.\n• Aktual 177,6M LK: 1.776.369 LK diselamatkan = Efisiensi Rp 5,33 Miliar / Tahun.\n*Catatan: Nilai Rp 3.000/LK adalah estimasi internal untuk simulasi cost avoidance (bukan HPP rahasia Peruri).")
    ]
    for idx, (h, b) in enumerate(fin_items):
        p = tf_fin.paragraphs[0] if idx == 0 else tf_fin.add_paragraph()
        p.space_after = Pt(2)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(7.5)
        r0.font.bold = True
        r0.font.color.rgb = NAVY if idx < 2 else GREEN
        r1 = p.add_run()
        r1.text = h + " "
        r1.font.name = 'Arial'
        r1.font.size = Pt(7.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(7)
        r2.font.color.rgb = DARK_TEXT

    # Quadrant 4 (Bottom-Right): Matriks Evaluasi 5 Pilar Cost of Inaction (Section 5.3, Tabel 1.2)
    add_card(slide, 6.73, top_y2, quad_w + 0.05, quad_h - 0.05, FILL_LIGHT_RED, BORDER_RED)
    add_section_header(slide, 6.73, top_y2, quad_w + 0.05, 0.25, "4. MATRIKS EVALUASI 5 PILAR RISIKO PEMBIARAN (COST OF INACTION)", RED)

    tb_inaction = slide.shapes.add_textbox(Inches(6.83), Inches(top_y2 + 0.28), Inches(quad_w - 0.15), Inches(quad_h - 0.35))
    tf_in = tb_inaction.text_frame
    tf_in.word_wrap = True
    tf_in.margin_top = Inches(0)
    tf_in.margin_left = Inches(0)

    inaction_items = [
        ("1. Biaya (Cost) [KRITIS]:", "Akumulasi pemborosan bahan baku Rp 22,13 M s.d. Rp 24,56 Miliar/tahun & pembengkakan biaya tambah cetak."),
        ("2. Mutu (Quality) [TINGGI]:", "Inschiet berfluktuasi hingga 5,11% akibat penanganan suku cadang mesin yang reaktif & coba-coba."),
        ("3. Kepatuhan [TINGGI]:", "Buku manual menyulitkan penelusuran riwayat lot produksi saat audit ISO 9001:2015."),
        ("4. K3L / ESG [SEDANG]:", "Timbulan limbah afval 7,37M–8,18M LK/tahun (±60–65 Ton kertas) & kelelahan fisik shift malam."),
        ("5. Layanan (SLA) [TINGGI]:", "Siklus cetak ulang lambat memperlambat serah terima ke DJBC & ancaman denda penalti kontrak.")
    ]
    for idx, (h, b) in enumerate(inaction_items):
        p = tf_in.paragraphs[0] if idx == 0 else tf_in.add_paragraph()
        p.space_after = Pt(2)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(7.5)
        r0.font.bold = True
        r0.font.color.rgb = RED
        r1 = p.add_run()
        r1.text = h + " "
        r1.font.name = 'Arial'
        r1.font.size = Pt(7.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(7)
        r2.font.color.rgb = DARK_TEXT


# ==============================================================================
# SLIDE 3 OF 3: SOLUSI DSS SIRINE 4.0, VALIDASI S1 2026 & VALUE CREATION KORPORASI
# (Section 6 & Document Synthesis of Risalah)
# ==============================================================================
def build_slide_3(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide, "03",
        "Solusi Inovasi DSS SIRINE 4.0, Validasi Empiris Semester 1 2026 & Dampak Strategis",
        "Tunjukkan arsitektur integrasi 3 titik aliran data, kertas kerja realisasi penurunan inschiet ke 3,33%, penghematan Rp 2,23 Miliar, dan reduksi downtime 50%–75%."
    )

    # 4 Top Hero KPIs
    add_kpi_card(slide, 0.55, 1.68, 2.95, 0.98, "INSCHIET AKHIR Q2 2026", "3,33%", "Turun -1,28 pp (-27,77%) vs Baseline 4,61%", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 3.62, 1.68, 2.95, 0.98, "RATA-RATA SEMESTER 1", "3,89%", "Volume: 103.345.688 LK (< Target 4,00%)", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 6.69, 1.68, 2.95, 0.98, "LEMBAR DISELAMATKAN", "743.234 LK", "Kertas Sekuriti Fisik Diselamatkan selama S1", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 9.76, 1.68, 3.02, 0.98, "EFISIENSI BIAYA S1", "Rp 2,23 Miliar", "Proyeksi Tahunan: Rp 6,82 M | Dev: Rp 0,-", RED, FILL_LIGHT_AMBER, BORDER_AMBER)

    # Top Flow Banner: Arsitektur Integrasi 3 Titik Aliran Data & Runtutan Granular
    top_flow_w = 12.23
    add_card(slide, 0.55, 2.76, top_flow_w, 0.92, FILL_LIGHT_PURPLE, BORDER_PURPLE)

    tb_flow = slide.shapes.add_textbox(Inches(0.65), Inches(2.78), Inches(top_flow_w - 0.20), Inches(0.86))
    tf_f = tb_flow.text_frame
    tf_f.word_wrap = True
    tf_f.margin_top = Inches(0)
    tf_f.margin_left = Inches(0)

    p_f1 = tf_f.paragraphs[0]
    r_f0 = p_f1.add_run()
    r_f0.text = "INTEGRASI 3 TITIK ALIRAN DATA OPERASIONAL SEKETIKA:  "
    r_f0.font.name = 'Arial'
    r_f0.font.size = Pt(8)
    r_f0.font.bold = True
    r_f0.font.color.rgb = PURPLE
    r_f1 = p_f1.add_run()
    r_f1.text = "Data Transaksi Meja Mesin (< 30 dtk)  ⇄  Modul SAP Production Order (ZPPRSIPPC0012)  ⇄  Hasil Sortir Verifikasi Mutu (HCTS)\n"
    r_f1.font.name = 'Arial'
    r_f1.font.size = Pt(8)
    r_f1.font.bold = True
    r_f1.font.color.rgb = NAVY

    p_f2 = tf_f.add_paragraph()
    r_f2 = p_f2.add_run()
    r_f2.text = "RUNTUTAN DATA GRANULAR & TINDAKAN PRESKRIPTIF:  "
    r_f2.font.name = 'Arial'
    r_f2.font.size = Pt(8)
    r_f2.font.bold = True
    r_f2.font.color.rgb = DARK_GREEN
    r_f3 = p_f2.add_run()
    r_f3.text = "Nomor PO  ➔  Nomor Mesin (9 Mesin)  ➔  Pola Gilir (Shift 1/2/3)  ➔  Tim Operator  ➔  Kategori Cacat Cetak  ➔  Tindakan Korektif Cepat"
    r_f3.font.name = 'Arial'
    r_f3.font.size = Pt(8)
    r_f3.font.color.rgb = DARK_TEXT

    # Bottom Split Panels
    bot_y = 3.76
    bot_h = 3.30
    left_w = 4.80
    right_w = 7.28
    gap = 0.15

    # Left Bottom Panel: Native Chart Penurunan Inschiet + Reduksi Downtime
    add_card(slide, 0.55, bot_y, left_w, bot_h, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 0.55, bot_y, left_w, 0.25, "TREN PENURUNAN INSCHIET BASELINE VS REALISASI S1 2026 (%)", PURPLE)

    # Native Clustered Column Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Baseline 2025', 'Q1 2026 (Adaptasi)', 'Q2 2026 (Presisi)']
    chart_data.add_series('Inschiet (%)', (4.61, 4.34, 3.33))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.65), Inches(bot_y + 0.28), Inches(4.60), Inches(2.05),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    plots = chart.plots[0]
    plots.has_data_labels = True
    for series in plots.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = GREEN
        dl = series.data_labels
        dl.font.name = 'Arial'
        dl.font.size = Pt(9)
        dl.font.bold = True
        dl.font.color.rgb = DARK_GREEN
        dl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    # Callout Downtime Box
    dt_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(bot_y + 2.40), Inches(4.60), Inches(0.80))
    dt_box.fill.solid()
    dt_box.fill.fore_color.rgb = WHITE
    dt_box.line.color.rgb = GREEN
    dt_box.line.width = Pt(1.2)
    tf_dt = dt_box.text_frame
    tf_dt.word_wrap = True
    tf_dt.margin_top = Inches(0.04)
    tf_dt.margin_left = Inches(0.06)
    tf_dt.margin_right = Inches(0.06)
    p_dt = tf_dt.paragraphs[0]
    r_d0 = p_dt.add_run()
    r_d0.text = "EFISIENSI WAKTU PEMERIKSAAN MESIN (DOWNTIME):  "
    r_d0.font.name = 'Arial'
    r_d0.font.size = Pt(7.5)
    r_d0.font.bold = True
    r_d0.font.color.rgb = GREEN
    r_d1 = p_dt.add_run()
    r_d1.text = "Durasi penanganan mesin bermasalah terpangkas dari > 1 shift (> 8 jam) menjadi < 2–4 jam (reduksi jam henti produktif sebesar 50% s.d. 75%)."
    r_d1.font.name = 'Arial'
    r_d1.font.size = Pt(7)
    r_d1.font.color.rgb = DARK_TEXT

    # Right Bottom Panel: Tabel 1.3 Kertas Kerja Realisasi Finansial + 4 Pilar Nilai Tambah Korporasi
    add_card(slide, 0.55 + left_w + gap, bot_y, right_w, bot_h, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_section_header(slide, 0.55 + left_w + gap, bot_y, right_w, 0.25, "KERTAS KERJA REALISASI PENGHEMATAN BIAYA & TRANSFORMASI KORPORASI", DARK_GREEN)

    # Full Table 1.3
    table_shape = slide.shapes.add_table(4, 8, Inches(0.55 + left_w + gap + 0.08), Inches(bot_y + 0.30), Inches(right_w - 0.16), Inches(1.50))
    table = table_shape.table
    table.columns[0].width = Inches(1.30)
    table.columns[1].width = Inches(0.85)
    table.columns[2].width = Inches(0.68)
    table.columns[3].width = Inches(0.78)
    table.columns[4].width = Inches(0.85)
    table.columns[5].width = Inches(0.85)
    table.columns[6].width = Inches(0.92)
    table.columns[7].width = Inches(0.89)

    t3_headers = ["Periode", "Vol Produksi", "Inschiet", "Deviasi", "Target 4,61%", "Realisasi", "Diselamatkan", "Nilai Saving*"]
    t3_rows = [
        ("Q1 2026 (Adaptasi)", "57.385.254", "4,34%", "-0,27 pp", "2.645.460 lb", "2.490.520 lb", "154.940 LK", "Rp 464,82 Jt"),
        ("Q2 2026 (Presisi)", "45.960.434", "3,33%", "-1,28 pp", "2.118.776 lb", "1.530.482 lb", "588.294 LK", "Rp 1,76 M"),
        ("TOTAL S1 2026", "103.345.688", "3,89%", "-0,72 pp", "4.764.236 lb", "4.021.002 lb", "743.234 LK", "Rp 2,23 Miliar")
    ]

    for c_idx, h_text in enumerate(t3_headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.name = 'Arial'
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for r_idx, r_data in enumerate(t3_rows):
        for c_idx, val in enumerate(r_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = FILL_LIGHT_GREEN if r_idx == 2 else (WHITE if r_idx % 2 == 0 else FILL_LIGHT_NAVY)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = 'Arial'
            p.font.size = Pt(7)
            p.font.color.rgb = DARK_GREEN if r_idx == 2 else DARK_TEXT
            p.font.bold = (r_idx == 2 or c_idx in [0, 6, 7])
            p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT

    # Multi-Dimensional Value Creation Box (Bottom-Right)
    val_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55 + left_w + gap + 0.08), Inches(bot_y + 1.86), Inches(right_w - 0.16), Inches(1.34))
    val_box.fill.solid()
    val_box.fill.fore_color.rgb = FILL_LIGHT_NAVY
    val_box.line.color.rgb = BORDER_NAVY
    val_box.line.width = Pt(1.2)
    tf_vb = val_box.text_frame
    tf_vb.word_wrap = True
    tf_vb.margin_top = Inches(0.04)
    tf_vb.margin_left = Inches(0.08)
    tf_vb.margin_right = Inches(0.08)

    val_pillars = [
        ("1. Mutu & Operasional:", "Inschiet tembus 3,33% di Q2 2026, menyelamatkan 743.234 lembar kertas sekuriti fisik, dan memotong jam henti perbaikan mesin 50%–75%."),
        ("2. Finansial & Investasi:", "Cost avoidance nyata Rp 2,23 Miliar pada S1 2026 (proyeksi Rp 6,82 Miliar/tahun) dengan biaya in-house development Rp 0,- (payback period seketika)."),
        ("3. Budaya Kerja & Data:", "Eliminasi buku folio fisik (< 30 detik input PO), pembinaan objektif harian bagi ±42 operator tanpa prasangka antar-shift, kesiapan INDI 4.0."),
        ("4. Kepatuhan & Tender DJBC:", "Jaminan 100% kepatuhan SLA pengiriman ke Bea Cukai, eliminasi penalti kontrak, dan penguatan keunggulan bersaing Peruri dalam tender nasional.")
    ]
    for idx, (h, b) in enumerate(val_pillars):
        p = tf_vb.paragraphs[0] if idx == 0 else tf_vb.add_paragraph()
        p.space_after = Pt(2)
        r0 = p.add_run()
        r0.text = "✔ "
        r0.font.name = 'Arial'
        r0.font.size = Pt(7.5)
        r0.font.bold = True
        r0.font.color.rgb = GREEN
        r1 = p.add_run()
        r1.text = h + " "
        r1.font.name = 'Arial'
        r1.font.size = Pt(7.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_NAVY
        r2 = p.add_run()
        r2.text = b
        r2.font.name = 'Arial'
        r2.font.size = Pt(7)
        r2.font.color.rgb = DARK_TEXT


# ==============================================================================
# MAIN EXECUTION
# ==============================================================================
def main():
    prs, blank_layout = create_presentation()

    print("Building Ultra-Dense Slide 1 of 3 (Sections 1, 2, 3, 4: Pasar Global, Tender DJBC, Mandat & 9 Mesin Offset)...")
    build_slide_1(prs, blank_layout)

    print("Building Ultra-Dense Slide 2 of 3 (Sections 4.1, 4.2, 4.3, 5: Fluktuasi 2025, Data Silo, Finansial & Risiko Inaction)...")
    build_slide_2(prs, blank_layout)

    print("Building Ultra-Dense Slide 3 of 3 (Section 6 & Synthesis: Solusi DSS SIRINE 4.0, Validasi S1 2026 & Value Creation)...")
    build_slide_3(prs, blank_layout)

    output_path = "Presentasi_Risalah_Latar_Belakang_IAKA_2026.pptx"
    prs.save(output_path)
    print(f"\n[SUCCESS] Master 3-Slide Ultra-Dense Presentation successfully saved to: {output_path}")

if __name__ == "__main__":
    main()
