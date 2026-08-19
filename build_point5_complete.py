import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
import sys
import os

# --- PALETTE CONSTANTS ---
NAVY = RGBColor(32, 40, 143)         # #20288F - Primary Brand
PURPLE = RGBColor(106, 44, 201)      # #6A2CC9 - Accent
DARK_NAVY = RGBColor(26, 35, 126)    # #1A237E
GREEN = RGBColor(46, 125, 50)        # #2E7D32 - Success
DARK_GREEN = RGBColor(27, 94, 32)    # #1B5E20
RED = RGBColor(198, 40, 40)          # #C62828 - Alert
ORANGE = RGBColor(230, 81, 0)        # #E65100 - Warning
DARK_TEXT = RGBColor(42, 42, 61)     # #2A2A3D - Body Text
MUTED_TEXT = RGBColor(110, 110, 130) # #6E6E82 - Muted Text
WHITE = RGBColor(255, 255, 255)

# Container Fills
FILL_LIGHT_NAVY = RGBColor(244, 247, 254)   # #F4F7FE
FILL_LIGHT_PURPLE = RGBColor(248, 246, 254) # #F8F6FE
FILL_LIGHT_GREEN = RGBColor(232, 245, 233)  # #E8F5E9
FILL_LIGHT_AMBER = RGBColor(255, 249, 230)  # #FFF9E6
FILL_LIGHT_RED = RGBColor(255, 245, 245)    # #FFF5F5

# Borders
BORDER_NAVY = RGBColor(210, 220, 245)       # #D2DCF5
BORDER_PURPLE = RGBColor(217, 206, 247)     # #D9CEF7
BORDER_GREEN = RGBColor(165, 214, 167)      # #A5D6A7
BORDER_AMBER = RGBColor(255, 224, 130)      # #FFE082
BORDER_RED = RGBColor(255, 205, 210)        # #FFCDD2

def add_header(slide, badge_num, title_text, purpose_text):
    # Wave ornament
    if os.path.exists('wave_ornament.png'):
        slide.shapes.add_picture('wave_ornament.png', Inches(-0.118), Inches(-0.196), Inches(1.584), Inches(0.891))
    # Peruri Logo
    if os.path.exists('peruri_logo.png'):
        slide.shapes.add_picture('peruri_logo.png', Inches(12.09), Inches(0.013), Inches(1.247), Inches(0.694))

    # Badge Box
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(0.65), Inches(0.80), Inches(0.72))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PURPLE
    badge.line.color.rgb = PURPLE
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = badge_num
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    tb_title = slide.shapes.add_textbox(Inches(1.55), Inches(0.60), Inches(10.20), Inches(0.80))
    tf = tb_title.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Purpose Subtitle
    tb_purpose = slide.shapes.add_textbox(Inches(0.60), Inches(1.48), Inches(11.80), Inches(0.38))
    tf = tb_purpose.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Tujuan bagian ini:  "
    r1.font.name = 'Arial'
    r1.font.size = Pt(11.5)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = purpose_text
    r2.font.name = 'Arial'
    r2.font.size = Pt(11.5)
    r2.font.color.rgb = DARK_TEXT

    # Footer
    tb_foot = slide.shapes.add_textbox(Inches(0.60), Inches(7.02), Inches(12.13), Inches(0.24))
    tf = tb_foot.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = "IAKA 2026 — Kerangka Presentasi Peserta  ·  DSS SIRINE 4.0 Unit Cetak Pita Cukai"
    p.font.name = 'Arial'
    p.font.size = Pt(8.5)
    p.font.color.rgb = MUTED_TEXT

def add_kpi_card(slide, left, top, width, height, title, value, subtext, color, fill_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.04)
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)

    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = 'Arial'
    p0.font.size = Pt(8.5)
    p0.font.bold = True
    p0.font.color.rgb = color
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = value
    p1.font.name = 'Arial'
    p1.font.size = Pt(15.5)
    p1.font.bold = True
    p1.font.color.rgb = color
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = subtext
    p2.font.name = 'Arial'
    p2.font.size = Pt(7.5)
    p2.font.color.rgb = DARK_TEXT
    p2.alignment = PP_ALIGN.CENTER

def style_table_cell(cell, text, font_size=7.2, bold=False, text_color=DARK_TEXT, fill_color=None, align=PP_ALIGN.LEFT):
    cell.text = ""
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    else:
        cell.fill.background()
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = align

# -------------------------------------------------------------
# BUILD SLIDE 05.1
# -------------------------------------------------------------
def build_slide_5_1(slide):
    # Clear shapes on template slide
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    add_header(
        slide,
        badge_num="05.1",
        title_text="Implementasi: Lingkup Uji Coba Lapangan, Peran Tim & Penanganan Kendala",
        purpose_text="Buktikan solusi benar-benar dijalankan di lini produksi riil dengan pembagian peran terstruktur dan penyelesaian kendala yang adaptif."
    )

    # Top 4 Hero KPI Cards
    add_kpi_card(slide, 0.60, 1.90, 2.86, 1.08, "LINGKUP LINI CETAK (MVP)", "6 Mesin Utama", "4 Komori (KMR 1–4) & 2 Ryobi (RYB 1–2) + 3 GTO", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 3.69, 1.90, 2.86, 1.08, "POLA KERJA & SDM TERLIBAT", "3 Gilir (24 Jam)", "±42 Operator Cetak & Kepala Kelompok Non-stop", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 6.78, 1.90, 2.86, 1.08, "SAMPEL PENGUJIAN S1 2026", "103,3 Juta LK", "Q1: 57,38 Jt Lembar | Q2: 45,96 Jt Lembar Cetak", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 9.87, 1.90, 2.86, 1.08, "KEPATUHAN INPUT DIGITAL", "100% PO Terekam", "Penarikan Buku Folio per 1 Jan 2026 (BA-002)", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)

    # Left Panel: 1. LINGKUP PENGUJIAN LAPANGAN & STRUKTUR PERAN (PIC)
    header_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.60), Inches(3.08), Inches(5.90), Inches(0.34))
    header_left.fill.solid()
    header_left.fill.fore_color.rgb = NAVY
    header_left.line.fill.background()
    tf = header_left.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "1. LINGKUP PENGUJIAN LAPANGAN & STRUKTUR PERAN (PIC)"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(3.42), Inches(5.90), Inches(3.02))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = FILL_LIGHT_NAVY
    box_left.line.color.rgb = BORDER_NAVY
    box_left.line.width = Pt(1.0)

    # Left Card 1: Parameter Lokasi & Periode
    card_l1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(3.50), Inches(5.66), Inches(1.30))
    card_l1.fill.solid()
    card_l1.fill.fore_color.rgb = WHITE
    card_l1.line.color.rgb = BORDER_NAVY
    card_l1.line.width = Pt(1.0)
    tf = card_l1.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    
    p = tf.paragraphs[0]
    p.text = "PARAMETER LOKASI, SARANA & PERIODE UJI COBA"
    p.font.name = 'Arial'
    p.font.size = Pt(8.2)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p1 = tf.add_paragraph()
    r1 = p1.add_run()
    r1.text = "• Lokasi & Sarana: "
    r1.font.bold = True
    r1.font.size = Pt(7.2)
    r1.font.color.rgb = NAVY
    r2 = p1.add_run()
    r2.text = "Gedung Cetak Sekuriti Karawang (Unit Cetak Pita Cukai). Memanfaatkan PC meja kontrol mesin & intranet internal Peruri (Rp 0 biaya lisensi/alat)."
    r2.font.size = Pt(7.2)
    r2.font.color.rgb = DARK_TEXT

    p2 = tf.add_paragraph()
    r1 = p2.add_run()
    r1.text = "• Linimasa Pengujian: "
    r1.font.bold = True
    r1.font.size = Pt(7.2)
    r1.font.color.rgb = NAVY
    r2 = p2.add_run()
    r2.text = "Nov–Des 2025 (Sosialisasi & IK) ➔ 1–2 Jan 2026 (Go-Live & Tarik Buku Folio) ➔ Jan–Mar 2026 (Fase 1 MVP / Adaptasi) ➔ Apr–Jun 2026 (Fase 2 Preskriptif Penuh)."
    r2.font.size = Pt(7.2)
    r2.font.color.rgb = DARK_TEXT

    # Left Card 2: Matriks Peran & Tanggung Jawab
    card_l2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(4.88), Inches(5.66), Inches(1.48))
    card_l2.fill.solid()
    card_l2.fill.fore_color.rgb = WHITE
    card_l2.line.color.rgb = BORDER_PURPLE
    card_l2.line.width = Pt(1.0)
    tf = card_l2.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)

    p = tf.paragraphs[0]
    p.text = "MATRIKS TANGGUNG JAWAB & STRUKTUR PERAN (PIC)"
    p.font.name = 'Arial'
    p.font.size = Pt(8.2)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    roles = [
        ("• Pembina & Fasilitator (Kadep Khazanah & Verifikasi SBU HSS & Kasie Cetak): ", "Menjamin wewenang lintas seksi (Cetak, Verifikasi, Perawatan) & legalitas SOP."),
        ("• Operator Meja Mesin (PIC Kelompok - ±42 Personel): ", "Input data PO via form digital < 30 detik seketika proses cetak selesai di mesin."),
        ("• Kepala Kelompok (Pengawas Shift): ", "Kelola jadwal gilir mingguan & verifikasi kelengkapan data di akhir gilir kerja."),
        ("• Kepala Unit & Teknisi Mesin: ", "Gunakan dasbor inschiet mesin untuk servis presisi & bimbingan teknis shift.")
    ]
    for r_title, r_desc in roles:
        p_r = tf.add_paragraph()
        r1 = p_r.add_run()
        r1.text = r_title
        r1.font.bold = True
        r1.font.size = Pt(7.1)
        r1.font.color.rgb = PURPLE
        r2 = p_r.add_run()
        r2.text = r_desc
        r2.font.size = Pt(7.1)
        r2.font.color.rgb = DARK_TEXT

    # Right Panel: 2. KENDALA YANG DIHADAPI & CARA PENANGGULANGANNYA
    header_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.78), Inches(3.08), Inches(5.95), Inches(0.34))
    header_right.fill.solid()
    header_right.fill.fore_color.rgb = PURPLE
    header_right.line.fill.background()
    tf = header_right.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "2. KENDALA YANG DIHADAPI & CARA PENANGGULANGANNYA"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.78), Inches(3.42), Inches(5.95), Inches(3.02))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = FILL_LIGHT_PURPLE
    box_right.line.color.rgb = BORDER_PURPLE
    box_right.line.width = Pt(1.0)

    # Problem Cards (3 Cards)
    problems = [
        (
            Inches(3.50), Inches(0.90), RED, BORDER_RED,
            "KENDALA 1: RESISTENSI KEBIASAAN INPUT DI MEJA MESIN",
            "• Masalah di Lapangan: ", "Operator merasa terbebani input digital di PC saat mengawasi jalannya lembar cetak cepat.",
            "• Cara Penanggulangan: ", "Rancang Lean UX dengan Autofill otomatis data SAP & regu (< 30 detik) + shortcut Ctrl+S."
        ),
        (
            Inches(4.46), Inches(0.90), ORANGE, BORDER_AMBER,
            "KENDALA 2: MASA ADAPTASI & DISIPLIN INPUT PADA Q1 2026",
            "• Masalah di Lapangan: ", "Disiplin entri belum seragam di awal tahun; inschiet Q1 baru turun tipis ke 4,34% (-5,9%).",
            "• Cara Penanggulangan: ", "Daily coaching tiap gilir, verifikasi penutup shift oleh Kepala Kelompok, & penegakan IK-001."
        ),
        (
            Inches(5.42), Inches(0.94), NAVY, BORDER_NAVY,
            "KENDALA 3: KEBIASAAN SERVIS MESIN SECARA ACAK/BERGILIR",
            "• Masalah di Lapangan: ", "Teknisi awalnya masih memeriksa 6 mesin bergilir saat ada laporan cacat umum.",
            "• Cara Penanggulangan: ", "Terapkan SOP-004: wajib cek Pareto modul cacat sebelum bongkar (downtime turun ke < 2–4 jam)."
        )
    ]

    for top_pos, h_size, title_col, border_col, title_txt, p_label, p_txt, s_label, s_txt in problems:
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.90), top_pos, Inches(5.71), h_size)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = border_col
        c.line.width = Pt(1.0)
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.04)
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = title_txt
        p.font.name = 'Arial'
        p.font.size = Pt(7.8)
        p.font.bold = True
        p.font.color.rgb = title_col

        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = p_label
        r1.font.bold = True
        r1.font.size = Pt(7.0)
        r1.font.color.rgb = title_col
        r2 = p1.add_run()
        r2.text = p_txt
        r2.font.size = Pt(7.0)
        r2.font.color.rgb = DARK_TEXT

        p2 = tf.add_paragraph()
        r1 = p2.add_run()
        r1.text = s_label
        r1.font.bold = True
        r1.font.size = Pt(7.0)
        r1.font.color.rgb = GREEN
        r2 = p2.add_run()
        r2.text = s_txt
        r2.font.size = Pt(7.0)
        r2.font.color.rgb = DARK_TEXT

    # Bottom Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(6.52), Inches(12.13), Inches(0.44))
    banner.fill.solid()
    banner.fill.fore_color.rgb = FILL_LIGHT_AMBER
    banner.line.color.rgb = BORDER_AMBER
    banner.line.width = Pt(1.0)
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "• Prinsip Keberhasilan Lapangan: "
    r1.font.name = 'Arial'
    r1.font.size = Pt(7.5)
    r1.font.bold = True
    r1.font.color.rgb = ORANGE
    r2 = p.add_run()
    r2.text = "Kendala diselesaikan bukan dengan instruksi birokratis, melainkan melalui penyederhanaan antarmuka (Lean UX < 30 detik), kepastian standar kerja resmi (IK & SOP), serta pendampingan langsung di setiap gilir kerja."
    r2.font.name = 'Arial'
    r2.font.size = Pt(7.5)
    r2.font.color.rgb = DARK_TEXT

# -------------------------------------------------------------
# BUILD SLIDE 05.2
# -------------------------------------------------------------
def build_slide_5_2(slide):
    # Clear shapes on slide
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    add_header(
        slide,
        badge_num="05.2",
        title_text="Implementasi: Validasi Data Before vs After & Realisasi vs Rencana",
        purpose_text="Sajikan bukti data Before vs After terverifikasi dengan jumlah sampel (n) dan perbandingan realisasi vs rencana target."
    )

    # Top 4 Hero KPI Cards
    add_kpi_card(slide, 0.60, 1.90, 2.86, 1.08, "PENURUNAN INSCHIET (Q2 2026)", "-1,28 pp (-27,8%)", "Baseline 4,61% ➔ 3,33% (Sumber: SAP ZPPRSIPPC0012)", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 3.69, 1.90, 2.86, 1.08, "WAKTU SERVIS PER MESIN", "< 2–4 Jam", "Pemeriksaan terpangkas ≥ 50%–75% (dari > 8 Jam)", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 6.78, 1.90, 2.86, 1.08, "WAKTU REKAP EVALUASI HARIAN", "0 Menit / Hari", "Otomatis seketika (hemat ±45 menit/hari)", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 9.87, 1.90, 2.86, 1.08, "CAPAIAN TARGET MUTU FASE 1", "210% Tercapai", "Target < 4,00% (-0,61 pp) terlampaui ke 3,33%", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)

    # Left Panel: 1. VALIDASI DATA BEFORE VS AFTER (JUMLAH SAMPEL N LENGKAP)
    header_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.60), Inches(3.08), Inches(6.15), Inches(0.34))
    header_left.fill.solid()
    header_left.fill.fore_color.rgb = NAVY
    header_left.line.fill.background()
    tf = header_left.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "1. VALIDASI DATA BEFORE VS AFTER (JUMLAH SAMPEL N LENGKAP)"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(3.42), Inches(6.15), Inches(3.02))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = FILL_LIGHT_NAVY
    box_left.line.color.rgb = BORDER_NAVY
    box_left.line.width = Pt(1.0)

    # Mini title for chart
    tb_c_title = slide.shapes.add_textbox(Inches(0.72), Inches(3.46), Inches(5.91), Inches(0.24))
    tf = tb_c_title.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Grafik Tren Inschiet: Baseline vs Target MVP vs Realisasi S1 2026 (%) "
    r1.font.name = 'Arial'
    r1.font.size = Pt(7.5)
    r1.font.bold = True
    r1.font.color.rgb = NAVY
    r2 = p.add_run()
    r2.text = "— Sumber: SAP ZPPRSIPPC0012"
    r2.font.name = 'Arial'
    r2.font.size = Pt(6.8)
    r2.font.italic = True
    r2.font.color.rgb = MUTED_TEXT

    # Native PPTX Clustered Column Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Baseline 2025', 'Target MVP', 'Realisasi Q1', 'Realisasi Q2']
    chart_data.add_series('Inschiet (%)', (4.61, 4.00, 4.34, 3.33))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.72), Inches(3.68), Inches(5.91), Inches(1.30),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.name = 'Arial'
    data_labels.font.size = Pt(8.0)
    data_labels.font.bold = True
    data_labels.font.color.rgb = DARK_NAVY
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    # Series styling
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = PURPLE

    # Table of 5 parameters Before vs After
    table_shape_l = slide.shapes.add_table(5, 4, Inches(0.72), Inches(5.02), Inches(5.91), Inches(1.18))
    t_l = table_shape_l.table
    t_l.columns[0].width = Inches(1.35)
    t_l.columns[1].width = Inches(1.30)
    t_l.columns[2].width = Inches(1.30)
    t_l.columns[3].width = Inches(1.96)

    headers_l = ["Parameter Kinerja", "Before (Pra-Inovasi)", "After (DSS SIRINE 4.0)", "Perubahan (Δ) & Jumlah Data (n)"]
    for c_idx, h in enumerate(headers_l):
        style_table_cell(t_l.cell(0, c_idx), h, font_size=7.0, bold=True, text_color=WHITE, fill_color=NAVY, align=PP_ALIGN.CENTER if c_idx>0 else PP_ALIGN.LEFT)

    rows_l = [
        ("Tingkat Inschiet (%)", "4,61% (Puncak: 5,11%)", "3,33% (Q2 2026)", "-1,28 pp (-27,8%) | n=177,6M ➔ 45,9M LK", FILL_LIGHT_GREEN),
        ("Downtime Servis Mesin", "> 8 Jam / Mesin", "< 2–4 Jam / Mesin", "Turun ≥ 50%–75% | 6 Mesin Utama", WHITE),
        ("Waktu Rekap Evaluasi", "± 45 Menit / Hari", "0 Menit (Otomatis)", "Hemat 100% | Rekap Manual Hilang", FILL_LIGHT_GREEN),
        ("Input Data per PO", "± 3–5 Menit / PO", "< 30 Detik / PO", "Hemat ≥ 85% | n = 100% Order PO", WHITE)
    ]
    for r_idx, (p_name, b_val, a_val, diff_val, bg_col) in enumerate(rows_l, start=1):
        style_table_cell(t_l.cell(r_idx, 0), p_name, font_size=6.8, bold=True, text_color=DARK_NAVY, fill_color=bg_col)
        style_table_cell(t_l.cell(r_idx, 1), b_val, font_size=6.8, bold=False, text_color=RED if "5,11%" in b_val or "> 8" in b_val else DARK_TEXT, fill_color=bg_col, align=PP_ALIGN.CENTER)
        style_table_cell(t_l.cell(r_idx, 2), a_val, font_size=6.8, bold=True, text_color=GREEN if "3,33%" in a_val or "<" in a_val or "0" in a_val else DARK_TEXT, fill_color=bg_col, align=PP_ALIGN.CENTER)
        style_table_cell(t_l.cell(r_idx, 3), diff_val, font_size=6.8, bold=True, text_color=GREEN if "-" in diff_val or "Hemat" in diff_val or "Turun" in diff_val else DARK_NAVY, fill_color=bg_col)

    # Source label under left table
    tb_src_5_2 = slide.shapes.add_textbox(Inches(0.72), Inches(6.22), Inches(5.91), Inches(0.18))
    tf_src_5_2 = tb_src_5_2.text_frame
    tf_src_5_2.word_wrap = False
    tf_src_5_2.margin_top = Inches(0)
    tf_src_5_2.margin_left = Inches(0)
    p = tf_src_5_2.paragraphs[0]
    p.text = "*(Sumber Data Terverifikasi: Modul SAP ZPPRSIPPC0012 & Rekapitulasi Unit Verifikasi Mutu S1 2026)"
    p.font.name = 'Arial'
    p.font.size = Pt(6.5)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY

    # Right Panel: 2. PERBANDINGAN REALISASI VS RENCANA TARGET & LINIMASA
    header_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.90), Inches(3.08), Inches(5.83), Inches(0.34))
    header_right.fill.solid()
    header_right.fill.fore_color.rgb = PURPLE
    header_right.line.fill.background()
    tf = header_right.text_frame
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = "2. PERBANDINGAN REALISASI VS RENCANA TARGET & LINIMASA"
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE

    box_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.90), Inches(3.42), Inches(5.83), Inches(3.02))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = FILL_LIGHT_PURPLE
    box_right.line.color.rgb = BORDER_PURPLE
    box_right.line.width = Pt(1.0)

    # Card R1: Evaluation of Realization vs Target
    card_r1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.00), Inches(3.48), Inches(5.63), Inches(1.58))
    card_r1.fill.solid()
    card_r1.fill.fore_color.rgb = WHITE
    card_r1.line.color.rgb = BORDER_PURPLE
    card_r1.line.width = Pt(1.0)
    tf = card_r1.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.04)
    tf.margin_left = Inches(0.06)
    p = tf.paragraphs[0]
    p.text = "EVALUASI REALISASI VS TARGET SASARAN KINERJA (KPI)"
    p.font.name = 'Arial'
    p.font.size = Pt(7.8)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # Table for Realization vs Plan KPI
    table_shape_r = slide.shapes.add_table(5, 4, Inches(7.06), Inches(3.72), Inches(5.51), Inches(1.28))
    t_r = table_shape_r.table
    t_r.columns[0].width = Inches(1.35)
    t_r.columns[1].width = Inches(1.25)
    t_r.columns[2].width = Inches(1.35)
    t_r.columns[3].width = Inches(1.56)

    headers_r = ["Indikator Kinerja", "Rencana Target", "Realisasi Aktual", "Status Capaian"]
    for c_idx, h in enumerate(headers_r):
        style_table_cell(t_r.cell(0, c_idx), h, font_size=7.0, bold=True, text_color=WHITE, fill_color=PURPLE, align=PP_ALIGN.CENTER if c_idx>0 else PP_ALIGN.LEFT)

    rows_r = [
        ("Tingkat Inschiet", "< 4,00% (-0,61 pp)", "3,33% (-1,28 pp)", "Melampaui Target (210%)", FILL_LIGHT_GREEN),
        ("Penghematan Tahunan", "Rp 3,25 Miliar/Thn", "Rp 6,82 Miliar/Thn", "Melampaui Target (210%)", FILL_LIGHT_GREEN),
        ("Downtime Servis Mesin", "< 4 Jam / Tindakan", "< 2–4 Jam / Mesin", "100% Tercapai", WHITE),
        ("Waktu Rekap Harian", "< 5 Menit / Hari", "0 Menit (Otomatis)", "100% Tercapai", WHITE)
    ]
    for r_idx, (kpi_name, plan_val, act_val, status_val, bg_col) in enumerate(rows_r, start=1):
        style_table_cell(t_r.cell(r_idx, 0), kpi_name, font_size=6.8, bold=True, text_color=DARK_NAVY, fill_color=bg_col)
        style_table_cell(t_r.cell(r_idx, 1), plan_val, font_size=6.8, bold=False, text_color=DARK_TEXT, fill_color=bg_col, align=PP_ALIGN.CENTER)
        style_table_cell(t_r.cell(r_idx, 2), act_val, font_size=6.8, bold=True, text_color=GREEN, fill_color=bg_col, align=PP_ALIGN.CENTER)
        style_table_cell(t_r.cell(r_idx, 3), status_val, font_size=6.8, bold=True, text_color=GREEN, fill_color=bg_col, align=PP_ALIGN.CENTER)

    # Card R2: Evaluation of Roadmap Timeline
    card_r2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.00), Inches(5.12), Inches(5.63), Inches(1.24))
    card_r2.fill.solid()
    card_r2.fill.fore_color.rgb = WHITE
    card_r2.line.color.rgb = BORDER_GREEN
    card_r2.line.width = Pt(1.0)
    tf = card_r2.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.04)
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)

    p = tf.paragraphs[0]
    p.text = "REALISASI JADWAL LINIMASA (ROADMAP) OKT 2025 – JUN 2026"
    p.font.name = 'Arial'
    p.font.size = Pt(7.8)
    p.font.bold = True
    p.font.color.rgb = GREEN

    timeline_items = [
        ("• Tahap 1–2 (Okt 2025 – Jan 2026): ", "Form digital, IK-001/SOP-004, penarikan buku folio (BA-002) ➔ Tepat Waktu (100%)."),
        ("• Tahap 3 (Jan – Mar 2026 / Q1): ", "Uji coba lini 6 mesin (Masa Adaptasi Q1) ➔ Tepat Waktu (Inschiet 4,34%)."),
        ("• Tahap 4 (Apr – Jun 2026 / Q2): ", "Implementasi preskriptif penuh ➔ Tepat Waktu (Inschiet 3,33% / Hemat Rp 6,82 M)."),
        ("• Evaluasi Deviasi Proyek: ", "0 Hari Deviasi (Seluruh 6 tahapan terlaksana 100% on-schedule).")
    ]
    for t_label, t_desc in timeline_items:
        p_t = tf.add_paragraph()
        r1 = p_t.add_run()
        r1.text = t_label
        r1.font.bold = True
        r1.font.size = Pt(6.8)
        r1.font.color.rgb = DARK_GREEN if "Deviasi" in t_label else DARK_NAVY
        r2 = p_t.add_run()
        r2.text = t_desc
        r2.font.size = Pt(6.8)
        r2.font.color.rgb = DARK_TEXT

    # Bottom Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(6.52), Inches(12.13), Inches(0.44))
    banner.fill.solid()
    banner.fill.fore_color.rgb = FILL_LIGHT_GREEN
    banner.line.color.rgb = BORDER_GREEN
    banner.line.width = Pt(1.0)
    tf = banner.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.06)
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "• Kesimpulan Realisasi: "
    r1.font.name = 'Arial'
    r1.font.size = Pt(7.5)
    r1.font.bold = True
    r1.font.color.rgb = DARK_GREEN
    r2 = p.add_run()
    r2.text = "Pada periode uji coba Semester 1 2026 (103,3 Juta LK), efisiensi riil tercatat Rp 2,23 Miliar (743.234 LK diselamatkan). Seluruh 5 indikator kinerja melampaui target awal tanpa ada deviasi keterlambatan jadwal proyek."
    r2.font.name = 'Arial'
    r2.font.size = Pt(7.5)
    r2.font.color.rgb = DARK_TEXT

# -------------------------------------------------------------
# MAIN EXECUTION
# -------------------------------------------------------------
def main():
    prs = pptx.Presentation('Kerangka_Presentasi_Peserta_IAKA_2026 (1).pptx')
    print(f"Total slides initially: {len(prs.slides)}")

    # Slide 12 (index 11) is currently Point 5 Template
    slide_5_1 = prs.slides[11]
    build_slide_5_1(slide_5_1)
    print("Slide 05.1 built successfully.")

    # Check if slide 13 is already 05.2 or if we need to insert it
    # If total slides is 16, slide 12 is 06 Dampak, so we insert 05.2 right after 05.1
    # If slide 13 was already inserted from previous run, let's check
    if len(prs.slides) == 16:
        new_slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Move new_slide to index 12 (right after index 11)
        sldIdLst = prs.slides._sldIdLst
        slide_elem = sldIdLst[-1]
        sldIdLst.remove(slide_elem)
        sldIdLst.insert(12, slide_elem)
        slide_5_2 = prs.slides[12]
    else:
        # Check text of slide 13
        slide_5_2 = prs.slides[12]

    build_slide_5_2(slide_5_2)
    print("Slide 05.2 built successfully.")

    prs.save('Kerangka_Presentasi_Peserta_IAKA_2026 (1).pptx')
    print(f"Saved! Total slides now: {len(prs.slides)}")

if __name__ == '__main__':
    main()
