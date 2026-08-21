import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
import sys
import os

# ==========================================
# PALETTE & DESIGN CONSTANTS (Peruri Theme)
# ==========================================
NAVY = RGBColor(32, 40, 143)         # #20288F - Primary Brand Navy
PURPLE = RGBColor(106, 44, 201)      # #6A2CC9 - Accent Purple
DARK_NAVY = RGBColor(26, 35, 126)    # #1A237E - Deep Navy
GREEN = RGBColor(46, 125, 50)        # #2E7D32 - Success Green
DARK_GREEN = RGBColor(27, 94, 32)    # #1B5E20
RED = RGBColor(198, 40, 40)          # #C62828 - Alert Red
ORANGE = RGBColor(230, 81, 0)        # #E65100 - Warning Amber
DARK_TEXT = RGBColor(42, 42, 61)     # #2A2A3D - Body Text Charcoal
MUTED_TEXT = RGBColor(110, 110, 130) # #6E6E82 - Muted Text Gray
WHITE = RGBColor(255, 255, 255)

# Container Fills
FILL_LIGHT_NAVY = RGBColor(244, 247, 254)   # #F4F7FE
FILL_LIGHT_PURPLE = RGBColor(248, 246, 254) # #F8F6FE
FILL_LIGHT_GREEN = RGBColor(232, 245, 233)  # #E8F5E9
FILL_LIGHT_AMBER = RGBColor(255, 249, 230)  # #FFF9E6
FILL_LIGHT_RED = RGBColor(255, 245, 245)    # #FFF5F5

# Container Borders
BORDER_NAVY = RGBColor(210, 220, 245)       # #D2DCF5
BORDER_PURPLE = RGBColor(217, 206, 247)     # #D9CEF7
BORDER_GREEN = RGBColor(165, 214, 167)      # #A5D6A7
BORDER_AMBER = RGBColor(255, 224, 130)      # #FFE082
BORDER_RED = RGBColor(255, 205, 210)        # #FFCDD2

def create_presentation():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide layout
    return prs, blank_layout

def add_header(slide, badge_num, title_text, purpose_text):
    # Wave ornament top left
    if os.path.exists('wave_ornament.png'):
        slide.shapes.add_picture('wave_ornament.png', Inches(-0.118), Inches(-0.196), Inches(1.584), Inches(0.891))
    
    # Peruri Logo top right
    if os.path.exists('peruri_logo.png'):
        slide.shapes.add_picture('peruri_logo.png', Inches(12.09), Inches(0.013), Inches(1.247), Inches(0.694))

    # Badge Box
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(0.65), Inches(0.80), Inches(0.72))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PURPLE
    badge.line.color.rgb = PURPLE
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = badge_num
    p.font.name = 'Arial'
    p.font.size = Pt(20 if len(badge_num) <= 3 else 16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    tb_title = slide.shapes.add_textbox(Inches(1.55), Inches(0.60), Inches(10.20), Inches(0.80))
    tf = tb_title.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(23)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Purpose Subtitle
    tb_purpose = slide.shapes.add_textbox(Inches(0.60), Inches(1.48), Inches(11.80), Inches(0.38))
    tf = tb_purpose.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Tujuan bagian ini:  "
    r1.font.name = 'Arial'
    r1.font.size = Pt(11.5)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = purpose_text
    r2.font.name = 'Arial'
    r2.font.size = Pt(11.5)
    r2.font.color.rgb = DARK_TEXT

    # Footer
    tb_foot = slide.shapes.add_textbox(Inches(0.60), Inches(7.02), Inches(12.13), Inches(0.24))
    tf = tb_foot.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = "IAKA 2026 — Kerangka Presentasi Peserta  ·  DSS SIRINE 4.0 Unit Cetak Pita Cukai"
    p.font.name = 'Arial'
    p.font.size = Pt(8.5)
    p.font.color.rgb = MUTED_TEXT

def add_kpi_card(slide, left, top, width, height, title, value, subtext, color, fill_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(left + 0.10), Inches(top + 0.08), Inches(width - 0.20), Inches(height - 0.14))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)

    p1 = tf.paragraphs[0]
    p1.text = title.upper()
    p1.font.name = 'Arial'
    p1.font.size = Pt(8.5)
    p1.font.bold = True
    p1.font.color.rgb = color

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.name = 'Arial'
    p2.font.size = Pt(16.5)
    p2.font.bold = True
    p2.font.color.rgb = color

    p3 = tf.add_paragraph()
    p3.text = subtext
    p3.font.name = 'Arial'
    p3.font.size = Pt(8)
    p3.font.color.rgb = DARK_TEXT

def add_card(slide, left, top, width, height, fill_color, border_color, border_width=1.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(border_width)
    return card

def add_section_header(slide, left, top, width, height, title, bg_color):
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = bg_color
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    return hdr


# ==============================================================================
# SLIDE 1 OF 3 (TEMPLATE SLIDE 01): LATAR BELAKANG
# ==============================================================================
def build_slide_01_template(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide, "01",
        "Latar Belakang: Urgensi Stabilitas Mutu & Pencatatan Unit Cetak",
        "Bangun konteks: mengapa perbaikan sistem pencatatan & stabilitas mutu penting bagi operasional dan perusahaan."
    )

    col_w = 3.82
    gap = 0.33
    top_y = 1.95
    card_h = 4.95

    # 3 Column Panels matching the exact template:
    # 1. KONDISI SAAT INI DI UNIT CETAK
    # 2. MENGAPA INI PENTING BAGI PERUSAHAAN
    # 3. DATA AWAL PEMICU INOVASI
    columns_data = [
        ("1. KONDISI SAAT INI DI UNIT CETAK", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY, [
            ("Operasional 24 Jam Non-stop (9 Mesin):",
             " Unit Cetak Pita Cukai beroperasi 3 gilir kerja (Pagi 07–15, Sore 15–23, Malam 23–07 WIB) dengan 9 mesin cetak offset (4 Komori: KMR 1–4, 2 Ryobi: RYB 1–2, 3 GTO: GTO 1–3) dan didukung ±42 operator."),
            ("Pencatatan Fisik Manual di Meja Mesin:",
             " Hasil produksi tiap mesin dan operator dicatat manual pada buku folio fisik di meja kontrol. Data transaksi terisolasi dan hanya menumpuk di area mesin tanpa sinkronisasi jaringan."),
            ("Rekapitulasi Manual Terpisah & Lambat:",
             " Buku folio baru direkapitulasi secara manual oleh Kepala Kelompok saat evaluasi triwulan pegawai, sementara modul SAP di kantor hanya menyajikan tabel angka mentah."),
            ("Laporan Mutu Terpusat Global:",
             " Laporan dari Unit Verifikasi QC dan SAP ZPPRSIPPC0012 hanya menampilkan ringkasan kerusakan global di level unit tanpa atribusi nomor mesin, nomor PO, dan regu gilir kerja.")
        ]),
        ("2. MENGAPA INI PENTING BAGI PERUSAHAAN", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE, [
            ("Mandat Dokumen Sekuriti Negara (PP 06/2019):",
             " Mencetak Pita Cukai (PCHT & MMEA) di bawah kontrak tender DJBC Kemenkeu RI sebagai instrumen pelunasan fiskal ratusan triliun rupiah ke kas APBN."),
            ("Skala Volume Masif Terverifikasi:",
             " Target kapasitas rata-rata 160 Juta Lembar/tahun, dengan realisasi pesanan aktual 2025 mencapai 177.636.930 Lembar Cetak (Sumber: Modul SAP ZPPRSIPPC0012)."),
            ("Standar Industri Sekuriti (Intergraf & WCO):",
             " Menuntut jaminan mutu tanpa cacat (zero-defect) pada kertas serat pengaman, tinta UV, guilloche, dan hologram; rekonsiliasi ketat pemusnahan resmi lembar rusak HCTS (zero leakage); serta kepatuhan mutlak SLA distribusi."),
            ("Pengendalian Biaya Bahan Baku (Cost Leadership):",
             " Tingginya harga kertas sekuriti dan tinta khusus mewajibkan penekanan rasio pemborosan bahan (inschiet) guna menjaga efisiensi HPP dan daya saing tender.")
        ]),
        ("3. DATA AWAL PEMICU INOVASI", DARK_NAVY, FILL_LIGHT_NAVY, BORDER_NAVY, [
            ("Capaian SIRINE 3.5 (2024):",
             " Berhasil menurunkan inschiet tahunan ke 4,06% dengan menyajikan data jenis kerusakan umum tingkat unit."),
            ("Titik Jenuh & Fluktuasi Baseline 2025:",
             " Sepanjang 2025 inschiet kembali berfluktuasi dengan rata-rata 4,61% (Q1: 4,72%, Q2: 3,97%, Q3: 4,64%, Q4: 5,11% — Konsolidasi SAP ZPPRSIPPC0012 & Verifikasi Mutu)."),
            ("Tantangan Pesanan Desain Baru Q4:",
             " Saat order desain baru melonjak di Q4, inschiet naik tajam ke puncaknya 5,11% (+1,14 pp vs Q2) karena di area mesin tidak tersedia diagnostik data harian sehingga make-ready memanjang."),
            ("Pemicu Kebutuhan DSS SIRINE 4.0:",
             " Diperlukan sistem terintegrasi yang menghubungkan transaksi meja mesin, SAP, dan sortir verifikasi secara seketika untuk mengeliminasi titik buta.")
        ])
    ]

    for idx, (col_title, col_col, col_fill, col_bord, bullets) in enumerate(columns_data):
        cx = 0.60 + idx * (col_w + gap)
        add_card(slide, cx, top_y, col_w, card_h, col_fill, col_bord)
        add_section_header(slide, cx, top_y, col_w, 0.40, col_title, col_col)

        tb = slide.shapes.add_textbox(Inches(cx + 0.15), Inches(top_y + 0.48), Inches(col_w - 0.30), Inches(card_h - 0.58))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)

        for b_idx, (b_head, b_desc) in enumerate(bullets):
            p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
            p.space_after = Pt(7)
            r0 = p.add_run()
            r0.text = "• "
            r0.font.name = 'Arial'
            r0.font.size = Pt(9.5)
            r0.font.bold = True
            r0.font.color.rgb = col_col
            r1 = p.add_run()
            r1.text = b_head
            r1.font.name = 'Arial'
            r1.font.size = Pt(9.5)
            r1.font.bold = True
            r1.font.color.rgb = DARK_TEXT
            r2 = p.add_run()
            r2.text = " " + b_desc
            r2.font.name = 'Arial'
            r2.font.size = Pt(8.5)
            r2.font.color.rgb = DARK_TEXT


# ==============================================================================
# SLIDE 2 OF 3 (TEMPLATE SLIDE 02.1): MASALAH - TITIK BUTA OPERASIONAL
# ==============================================================================
def build_slide_02_1_template(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide, "02.1",
        "Masalah: Keterbatasan Data Granular & Titik Buta Operasional",
        "Nyatakan masalah operasional nyata dan terukur — apa yang terjadi, di unit mana, sejak kapan, dan siapa terdampak."
    )

    col_w = 3.82
    gap = 0.33
    top_y = 1.95
    card_h = 4.95

    # 3 Column Panels matching the exact template:
    # 1. TITIK BUTA PASCA-SIRINE 2024
    # 2. DUA KASUS NYATA DI LAPANGAN
    # 3. PIHAK TERDAMPAK LANGSUNG
    columns_data = [
        ("1. TITIK BUTA PASCA-SIRINE 2024", RED, FILL_LIGHT_RED, BORDER_RED, [
            ("Capaian SIRINE 3.5 (2024):",
             " Berhasil memetakan jenis kerusakan apa yang mendominasi dan menekan inschiet ke 4,06% di akhir 2024."),
            ("Kesenjangan Baru di 2025 (The Missing Link):",
             " Data jenis kerusakan sudah ada di SAP kantor, tetapi dengan 9 mesin cetak (4 Komori, 2 Ryobi, 3 GTO) dan 3 shift, belum ada data untuk tahu mesin mana dan shift mana yang menjadi sumber masalah utama."),
            ("Ketiadaan Atribusi Kondisi Kerja:",
             " Belum dapat dibedakan secara objektif apakah tingginya inschiet dipicu oleh degradasi mekanis komponen mesin atau variasi metode kerja/kelelahan operator shift malam."),
            ("Lokasi & Periode Masalah:",
             " Berpusat pada lini cetak sheet-fed offset Unit Cetak Pita Cukai, Dept. Khazanah & Verifikasi sejak Januari 2025.")
        ]),
        ("2. DUA KASUS NYATA DI LAPANGAN", ORANGE, FILL_LIGHT_AMBER, BORDER_AMBER, [
            ("Kasus 1 – Penelusuran Spekulatif (> 8 Jam):",
             " Bulan Juni 2025 terjadi lonjakan cacat blobor. Tanpa data per-mesin, teknisi memeriksa 9 mesin satu per satu secara coba-coba (trial-and-error)."),
            ("Dampak Downtime Kasus 1:",
             " Sumber masalah ternyata di Komori 3 (KMR3). Karena diperiksa terakhir, timbul downtime > 8 jam (> 1 shift) dan mesin terus mencetak lembar rusak."),
            ("Kasus 2 – Dilema Mesin vs Operator Shift:",
             " Mesin KMR3 pada Shift Malam mencatat inschiet 8,5%, sedangkan Shift Pagi di mesin yang sama hanya 2,5%. Tanpa data shift, teknisi selalu menyetel ulang mesin padahal masalahnya pada variasi make-ready / kelelahan sirkadian Shift Malam (23.00–07.00 WIB)."),
            ("Akar Mekanis Komponen Fisik Mesin:",
             " Rol karet tinta/air mengeras/licin (blobor/bintik), blanket aus/turun elastisitas (ghosting), penjepit silinder kendur (misregister & plooi).")
        ]),
        ("3. PIHAK TERDAMPAK LANGSUNG", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE, [
            ("Operator Cetak (±42 Orang):",
             " Sulit memantau hasil mutu hariannya secara mandiri dan kehilangan kesempatan memperoleh pembinaan (coaching) teknis objektif."),
            ("Kepala Kelompok & Kepala Unit:",
             " Menanggung beban rekapitulasi manual buku folio dan kesulitan memberi feedback terarah kepada tim kerja harian."),
            ("Teknisi Pemeliharaan (Maintenance):",
             " Kehilangan jam produktif (> 8 jam per mesin) akibat mencari sumber masalah secara acak tanpa prioritas data."),
            ("PPIC & Manajemen Unit:",
             " Sulit merencanakan alokasi 9 mesin cetak secara presisi karena profil keandalan mesin aktual tidak tersedia."),
            ("Pelanggan Utama (DJBC Kemenkeu RI):",
             " Menghadapi risiko keterlambatan pengiriman pita cukai akibat siklus cetak ulang (tambah cetak) lembar rusak.")
        ])
    ]

    for idx, (col_title, col_col, col_fill, col_bord, bullets) in enumerate(columns_data):
        cx = 0.60 + idx * (col_w + gap)
        add_card(slide, cx, top_y, col_w, card_h, col_fill, col_bord)
        add_section_header(slide, cx, top_y, col_w, 0.40, col_title, col_col)

        tb = slide.shapes.add_textbox(Inches(cx + 0.15), Inches(top_y + 0.48), Inches(col_w - 0.30), Inches(card_h - 0.58))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)

        for b_idx, (b_head, b_desc) in enumerate(bullets):
            p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
            p.space_after = Pt(7)
            r0 = p.add_run()
            r0.text = "• "
            r0.font.name = 'Arial'
            r0.font.size = Pt(9.5)
            r0.font.bold = True
            r0.font.color.rgb = col_col
            r1 = p.add_run()
            r1.text = b_head
            r1.font.name = 'Arial'
            r1.font.size = Pt(9.5)
            r1.font.bold = True
            r1.font.color.rgb = DARK_TEXT
            r2 = p.add_run()
            r2.text = " " + b_desc
            r2.font.name = 'Arial'
            r2.font.size = Pt(8.5)
            r2.font.color.rgb = DARK_TEXT


# ==============================================================================
# SLIDE 3 OF 3 (TEMPLATE SLIDE 02.2): MASALAH - BASELINE & RISIKO PEMBIARAN
# ==============================================================================
def build_slide_02_2_template(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide, "02.2",
        "Masalah: Baseline Permasalahan & Risiko Bila Dibiarkan",
        "Tunjukkan baseline permasalahan operasional terukur dan dampak multidimensi jika kondisi dibiarkan tanpa intervensi."
    )

    # 4 Top Hero KPIs (Exact coordinates and layout from template)
    add_kpi_card(slide, 0.60, 1.90, 2.86, 1.08, "BASELINE INSCHIET 2025", "4,61%", "Puncak Q4: 5,11% (SAP ZPPRSIPPC0012)", RED, FILL_LIGHT_RED, BORDER_RED)
    add_kpi_card(slide, 3.69, 1.90, 2.86, 1.08, "JUMLAH LEMBAR RUSAK", "8.189.062 LK", "Total order: 177.636.930 Lembar Cetak", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 6.78, 1.90, 2.86, 1.08, "POTENSI BIAYA TERBUANG*", "Rp 24,56 Miliar*", "*Simulasi estimasi @ Rp 3.000 / LK", RED, FILL_LIGHT_AMBER, BORDER_AMBER)
    add_kpi_card(slide, 9.87, 1.90, 2.86, 1.08, "DOWNTIME PERBAIKAN", "> 8 Jam / Mesin", "Pemeriksaan spekulatif bergilir (> 1 shift)", ORANGE, FILL_LIGHT_RED, BORDER_RED)

    # Left Container: TREN INSCHIET BASELINE PER KUARTAL 2025
    left_w = 5.95
    add_card(slide, 0.60, 3.44, left_w, 3.02, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 0.60, 3.08, left_w, 0.36, "TREN INSCHIET BASELINE PER KUARTAL 2025", PURPLE)

    # Native Clustered Column Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1 (Jan-Mar)', 'Q2 (Apr-Jun)', 'Q3 (Jul-Sep)', 'Q4 (Okt-Des)', 'Rata-rata 2025']
    chart_data.add_series('Inschiet (%)', (4.72, 3.97, 4.64, 5.11, 4.61))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.75), Inches(3.50), Inches(5.65), Inches(1.85),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    plots = chart.plots[0]
    plots.has_data_labels = True
    for series in plots.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = PURPLE
        dl = series.data_labels
        dl.font.name = 'Arial'
        dl.font.size = Pt(9.5)
        dl.font.bold = True
        dl.font.color.rgb = DARK_NAVY
        dl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    # Subtext underneath chart
    tb_chart_sub = slide.shapes.add_textbox(Inches(0.78), Inches(5.40), Inches(5.59), Inches(1.00))
    tf_cs = tb_chart_sub.text_frame
    tf_cs.word_wrap = True
    tf_cs.margin_top = Inches(0)
    tf_cs.margin_left = Inches(0)
    tf_cs.margin_right = Inches(0)

    p1 = tf_cs.paragraphs[0]
    p1.text = "Rata-rata Tahunan: 4,61%  |  Sumber: SAP ZPPRSIPPC0012 & QC Verifikasi"
    p1.font.name = 'Arial'
    p1.font.size = Pt(8.5)
    p1.font.bold = True
    p1.font.color.rgb = DARK_NAVY
    p1.space_after = Pt(2)

    p2 = tf_cs.add_paragraph()
    r_q2_h = p2.add_run()
    r_q2_h.text = "• Q2 (3,97%): "
    r_q2_h.font.name = 'Arial'
    r_q2_h.font.size = Pt(8)
    r_q2_h.font.bold = True
    r_q2_h.font.color.rgb = GREEN
    r_q2_t = p2.add_run()
    r_q2_t.text = "Bukti lini cetak mampu mencapai target < 4,00% saat kondisi terkontrol.\n"
    r_q2_t.font.name = 'Arial'
    r_q2_t.font.size = Pt(8)
    r_q2_t.font.color.rgb = DARK_TEXT

    r_q4_h = p2.add_run()
    r_q4_h.text = "• Q4 (5,11%): "
    r_q4_h.font.name = 'Arial'
    r_q4_h.font.size = Pt(8)
    r_q4_h.font.bold = True
    r_q4_h.font.color.rgb = RED
    r_q4_t = p2.add_run()
    r_q4_t.text = "Lonjakan saat pesanan desain baru masuk tanpa diagnostik harian di meja mesin.\n"
    r_q4_t.font.name = 'Arial'
    r_q4_t.font.size = Pt(8)
    r_q4_t.font.color.rgb = DARK_TEXT

    r_val_h = p2.add_run()
    r_val_h.text = "• Valuasi Tiap 1% Inschiet: "
    r_val_h.font.name = 'Arial'
    r_val_h.font.size = Pt(8)
    r_val_h.font.bold = True
    r_val_h.font.color.rgb = DARK_GREEN
    r_val_t = p2.add_run()
    r_val_t.text = "1.600.000 – 1.776.369 LK diselamatkan ≈ Rp 4,80 M – Rp 5,33 Miliar / Tahun."
    r_val_t.font.name = 'Arial'
    r_val_t.font.size = Pt(8)
    r_val_t.font.color.rgb = DARK_TEXT

    # Right Container: DAMPAK RISIKO BILA DIBIARKAN (5 PILAR)
    right_w = 5.95
    add_card(slide, 6.78, 3.44, right_w, 3.02, FILL_LIGHT_RED, BORDER_RED)
    add_section_header(slide, 6.78, 3.08, right_w, 0.36, "DAMPAK RISIKO BILA DIBIARKAN (5 PILAR)", RED)

    tb_right = slide.shapes.add_textbox(Inches(6.92), Inches(3.52), Inches(right_w - 0.28), Inches(2.88))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True
    tf_r.margin_top = Inches(0)
    tf_r.margin_left = Inches(0)
    tf_r.margin_right = Inches(0)

    inaction_bullets = [
        ("1. Biaya (Cost) [KRITIS]:",
         "Akumulasi pemborosan biaya cetak mencapai Rp 22,13 M s.d. Rp 24,56 Miliar/tahun (8.189.062 LK × asumsi Rp 3.000*) dan pembengkakan biaya tambah cetak."),
        ("2. Mutu (Quality) [TINGGI]:",
         "Inschiet berfluktuasi tanpa kendali (puncak 5,11%). Penanganan mesin bersifat sementara karena tidak menyentuh akar masalah per-mesin."),
        ("3. Kepatuhan (Compliance) [TINGGI]:",
         "Akuntabilitas pelacakan (traceability) dokumen negara terhambat karena pencatatan manual di buku folio tidak dapat diaudit digital (ISO 9001:2015)."),
        ("4. K3L & Lingkungan (ESG) [SEDANG]:",
         "Timbulan limbah padat kertas sekuriti mencapai 7,37 – 8,18 Juta LK/tahun (±60–65 ton kertas terbuang) dan peningkatan kelelahan operator di shift malam."),
        ("5. Layanan (Service SLA) [TINGGI]:",
         "Risiko keterlambatan serah terima pesanan ke DJBC akibat antrean cetak ulang lembar rusak yang mengancam reputasi SLA dan memicu penalti kontrak.")
    ]

    for idx, (b_head, b_desc) in enumerate(inaction_bullets):
        p = tf_r.paragraphs[0] if idx == 0 else tf_r.add_paragraph()
        p.space_after = Pt(4)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(8.5)
        r0.font.bold = True
        r0.font.color.rgb = RED
        r1 = p.add_run()
        r1.text = b_head + " "
        r1.font.name = 'Arial'
        r1.font.size = Pt(8.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.name = 'Arial'
        r2.font.size = Pt(8)
        r2.font.color.rgb = DARK_TEXT

    # Bottom Confidentiality Note Box
    note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(6.52), Inches(12.13), Inches(0.44))
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = FILL_LIGHT_NAVY
    note_box.line.color.rgb = BORDER_NAVY
    note_box.line.width = Pt(1.0)
    tf_n = note_box.text_frame
    tf_n.word_wrap = True
    tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_n.margin_left = Inches(0.12)
    tf_n.margin_right = Inches(0.12)
    tf_n.margin_top = Inches(0)
    tf_n.margin_bottom = Inches(0)
    p_n = tf_n.paragraphs[0]
    p_n.text = "*Catatan Asumsi Finansial (Confidentiality Protection): Nilai Rp 3.000/lembar adalah estimasi internal (kertas sekuriti, tinta khusus, depresiasi & operasional mesin, tenaga kerja) semata-mata untuk simulasi dampak inovasi (cost avoidance), BUKAN harga jual resmi atau rincian biaya pokok produksi resmi dari Perum Peruri yang bersifat rahasia perusahaan."
    p_n.font.name = 'Arial'
    p_n.font.size = Pt(7.5)
    p_n.font.italic = True
    p_n.font.color.rgb = MUTED_TEXT


# ==============================================================================
# MAIN EXECUTION
# ==============================================================================
def main():
    prs, blank_layout = create_presentation()

    print("Building Slide 1 of 3: Badge 01 (Latar Belakang: Urgensi Stabilitas Mutu & Pencatatan Unit Cetak)...")
    build_slide_01_template(prs, blank_layout)

    print("Building Slide 2 of 3: Badge 02.1 (Masalah: Keterbatasan Data Granular & Titik Buta Operasional)...")
    build_slide_02_1_template(prs, blank_layout)

    print("Building Slide 3 of 3: Badge 02.2 (Masalah: Baseline Permasalahan & Risiko Bila Dibiarkan)...")
    build_slide_02_2_template(prs, blank_layout)

    output_path = "Presentasi_Risalah_Latar_Belakang_IAKA_2026.pptx"
    prs.save(output_path)
    print(f"\n[SUCCESS] Exact Template 3-Slide Presentation successfully generated: {output_path}")

    # Also update the 3-slides backup file
    prs.save("Presentasi_Risalah_Latar_Belakang_IAKA_2026_3_Slides.pptx")

if __name__ == "__main__":
    main()
