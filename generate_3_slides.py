import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
import sys
import os

# ==========================================
# PALETTE & DESIGN CONSTANTS (Peruri Theme)
# ==========================================
NAVY = RGBColor(32, 40, 143)         # #20288F - Primary Brand Navy
PURPLE = RGBColor(106, 44, 201)      # #6A2CC9 - Accent Purple
DARK_NAVY = RGBColor(26, 35, 126)    # #1A237E - Deep Navy
GREEN = RGBColor(46, 125, 50)        # #2E7D32 - Success Green
DARK_GREEN = RGBColor(27, 94, 32)    # #1B5E20
RED = RGBColor(198, 40, 40)          # #C62828 - Alert Red
ORANGE = RGBColor(230, 81, 0)        # #E65100 - Warning Amber
DARK_TEXT = RGBColor(42, 42, 61)     # #2A2A3D - Body Text Charcoal
MUTED_TEXT = RGBColor(110, 110, 130) # #6E6E82 - Muted Text Gray
WHITE = RGBColor(255, 255, 255)

# Container Fills
FILL_LIGHT_NAVY = RGBColor(244, 247, 254)   # #F4F7FE
FILL_LIGHT_PURPLE = RGBColor(248, 246, 254) # #F8F6FE
FILL_LIGHT_GREEN = RGBColor(232, 245, 233)  # #E8F5E9
FILL_LIGHT_AMBER = RGBColor(255, 249, 230)  # #FFF9E6
FILL_LIGHT_RED = RGBColor(255, 245, 245)    # #FFF5F5

# Container Borders
BORDER_NAVY = RGBColor(210, 220, 245)       # #D2DCF5
BORDER_PURPLE = RGBColor(217, 206, 247)     # #D9CEF7
BORDER_GREEN = RGBColor(165, 214, 167)      # #A5D6A7
BORDER_AMBER = RGBColor(255, 224, 130)      # #FFE082
BORDER_RED = RGBColor(255, 205, 210)        # #FFCDD2

def create_presentation():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide layout
    return prs, blank_layout

def add_header(slide, badge_num, title_text, purpose_text):
    # Wave ornament top left
    if os.path.exists('wave_ornament.png'):
        slide.shapes.add_picture('wave_ornament.png', Inches(-0.118), Inches(-0.196), Inches(1.584), Inches(0.891))
    
    # Peruri Logo top right
    if os.path.exists('peruri_logo.png'):
        slide.shapes.add_picture('peruri_logo.png', Inches(12.09), Inches(0.013), Inches(1.247), Inches(0.694))

    # Badge Box
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(0.60), Inches(0.75), Inches(0.70))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PURPLE
    badge.line.color.rgb = PURPLE
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = badge_num
    p.font.name = 'Arial'
    p.font.size = Pt(20 if len(badge_num) <= 2 else 16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    tb_title = slide.shapes.add_textbox(Inches(1.50), Inches(0.55), Inches(10.30), Inches(0.80))
    tf = tb_title.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Purpose Subtitle
    tb_purpose = slide.shapes.add_textbox(Inches(0.60), Inches(1.40), Inches(12.13), Inches(0.35))
    tf = tb_purpose.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Tujuan bagian ini:  "
    r1.font.name = 'Arial'
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE
    r2 = p.add_run()
    r2.text = purpose_text
    r2.font.name = 'Arial'
    r2.font.size = Pt(11)
    r2.font.color.rgb = DARK_TEXT

    # Footer
    tb_foot = slide.shapes.add_textbox(Inches(0.60), Inches(7.06), Inches(12.13), Inches(0.24))
    tf = tb_foot.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = "IAKA 2026 — Kerangka Presentasi Peserta  ·  DSS SIRINE 4.0 Unit Cetak Pita Cukai (Perum Peruri)"
    p.font.name = 'Arial'
    p.font.size = Pt(8.5)
    p.font.color.rgb = MUTED_TEXT

def add_kpi_card(slide, left, top, width, height, title, value, subtext, color, fill_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.05))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(left + 0.10), Inches(top + 0.08), Inches(width - 0.20), Inches(height - 0.14))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)

    p1 = tf.paragraphs[0]
    p1.text = title.upper()
    p1.font.name = 'Arial'
    p1.font.size = Pt(8.5)
    p1.font.bold = True
    p1.font.color.rgb = color

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.name = 'Arial'
    p2.font.size = Pt(16.5)
    p2.font.bold = True
    p2.font.color.rgb = color

    p3 = tf.add_paragraph()
    p3.text = subtext
    p3.font.name = 'Arial'
    p3.font.size = Pt(8)
    p3.font.color.rgb = DARK_TEXT

def add_card(slide, left, top, width, height, fill_color, border_color, border_width=1.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(border_width)
    return card

def add_section_header(slide, left, top, width, height, title, bg_color):
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = bg_color
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0)
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    return hdr


# ==============================================================================
# SLIDE 1 OF 3: KONTEKS STRATEGIS, SKALA 9 MESIN & BASELINE MASALAH 2025 (BADGE 01)
# ==============================================================================
def build_slide_1_compressed(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide, "01",
        "Konteks Strategis, Skala Operasional 9 Mesin & Dinamika Baseline 2025",
        "Bangun konteks pasar sekuriti global, komitmen tender DJBC Kemenkeu RI (177,6M LK), profil 9 mesin, dan analisis fluktuasi baseline."
    )

    # 4 Top Hero KPIs
    add_kpi_card(slide, 0.60, 1.82, 2.90, 1.10, "SKALA PESANAN AKTUAL 2025", "177.636.930 LK", "Rata-rata: 160M LK | Modul SAP ZPPRSIPPC0012", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 3.65, 1.82, 2.90, 1.10, "9 MESIN CETAK OFFSET", "4 KMR · 2 RYB · 3 GTO", "Pola 3 Shift 24/7 (Pagi, Sore, Malam) | ±42 Operator", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 6.70, 1.82, 2.90, 1.10, "BASELINE INSCHIET 2025", "4,61%", "Q1: 4,72% · Q2: 3,97% · Q3: 4,64% · Q4: 5,11%", RED, FILL_LIGHT_RED, BORDER_RED)
    add_kpi_card(slide, 9.75, 1.82, 2.98, 1.10, "BEBAN BIAYA BASELINE*", "Rp 24,56 Miliar*", "8.189.062 Lembar Rusak/Thn (@ Rp 3.000 / LK)", ORANGE, FILL_LIGHT_AMBER, BORDER_AMBER)

    # Left Container: Konteks Pasar Global, Tender DJBC & Mandat Peruri
    col_w_left = 6.20
    add_card(slide, 0.60, 3.00, col_w_left, 3.95, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_section_header(slide, 0.60, 3.00, col_w_left, 0.32, "1. TUNTUTAN PASAR SEKURITI, TENDER DJBC & MANDAT PERURI", NAVY)

    tb_left = slide.shapes.add_textbox(Inches(0.72), Inches(3.36), Inches(col_w_left - 0.24), Inches(3.50))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_top = Inches(0)
    tf_l.margin_left = Inches(0)

    left_sections = [
        ("Standar Industri Sekuriti Internasional (Intergraf & WCO):",
         " Dokumen negara wajib presisi tinggi guna mencegah pemalsuan (*anti-counterfeiting*). Menerapkan *multi-layer security*: kertas berserat pengaman, tinta UV berpendar, garis halus *guilloche*, *microtext*, dan pita hologram."),
        ("Klausul Kepatuhan Tender DJBC Kemenkeu RI (PCHT & MMEA):",
         " Pita cukai adalah instrumen penerimaan APBN (ratusan triliun rupiah). Kontrak menuntut: (1) Jaminan mutu mutlak tanpa cacat (*zero-defect*); (2) Rekonsiliasi & pemusnahan resmi lembar rusak HCTS (*zero leakage*); (3) Ketepatan jadwal pasokan (*strict SLA*) tanpa keterlambatan cetak ulang."),
        ("Mandat PP No. 06/2019 & 3 Sasaran Strategis Peruri:",
         " (1) *Cost Leadership & Material Protection* (proteksi bahan baku mahal dari pemborosan *inschiet*); (2) *Operational Excellence* (standar ISO 9001:2015); (3) *Digitalisasi Area Kerja* (INDI 4.0). Seluruh sasaran bermuara pada lini operasional terbesar: **Unit Cetak Pita Cukai**.")
    ]

    for idx, (head, body) in enumerate(left_sections):
        p = tf_l.paragraphs[0] if idx == 0 else tf_l.add_paragraph()
        p.space_after = Pt(4)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(8.5)
        r0.font.bold = True
        r0.font.color.rgb = NAVY
        r1 = p.add_run()
        r1.text = head
        r1.font.name = 'Arial'
        r1.font.size = Pt(8.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = body
        r2.font.name = 'Arial'
        r2.font.size = Pt(8)
        r2.font.color.rgb = DARK_TEXT

    # Right Container: Visualisasi Tren Baseline 2025 & Analisis Fluktuasi
    col_w_right = 5.78
    add_card(slide, 6.95, 3.00, col_w_right, 3.95, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 6.95, 3.00, col_w_right, 0.32, "2. TREN FLUKTUASI BASELINE INSCHIET 2025 (%) & PEMBUKTIAN KAPABILITAS", PURPLE)

    # Native Clustered Column Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1 (Jan-Mar)', 'Q2 (Apr-Jun)', 'Q3 (Jul-Sep)', 'Q4 (Okt-Des)', 'Avg 2025']
    chart_data.add_series('Inschiet (%)', (4.72, 3.97, 4.64, 5.11, 4.61))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(7.05), Inches(3.38), Inches(5.55), Inches(2.05),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    plots = chart.plots[0]
    plots.has_data_labels = True
    for series in plots.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = PURPLE
        dl = series.data_labels
        dl.font.name = 'Arial'
        dl.font.size = Pt(9.5)
        dl.font.bold = True
        dl.font.color.rgb = DARK_NAVY
        dl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    # Split Box underneath Chart (Q2 vs Q4 Analysis)
    box_q2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(5.48), Inches(2.70), Inches(1.35))
    box_q2.fill.solid()
    box_q2.fill.fore_color.rgb = FILL_LIGHT_GREEN
    box_q2.line.color.rgb = BORDER_GREEN
    box_q2.line.width = Pt(1.2)
    tf_q2 = box_q2.text_frame
    tf_q2.word_wrap = True
    tf_q2.margin_top = Inches(0.04)
    tf_q2.margin_left = Inches(0.06)
    tf_q2.margin_right = Inches(0.06)
    p_q2 = tf_q2.paragraphs[0]
    p_q2.text = "PEMBUKTIAN KAPABILITAS Q2 (3,97%):"
    p_q2.font.name = 'Arial'
    p_q2.font.size = Pt(8)
    p_q2.font.bold = True
    p_q2.font.color.rgb = GREEN
    p_q2_b = tf_q2.add_paragraph()
    p_q2_b.text = "Membuktikan 9 mesin & operator mampu beroperasi di bawah toleransi 4,00% saat parameter mesin dan order stabil."
    p_q2_b.font.name = 'Arial'
    p_q2_b.font.size = Pt(7.5)
    p_q2_b.font.color.rgb = DARK_TEXT

    box_q4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.85 + 1.95), Inches(5.48), Inches(2.70), Inches(1.35))
    box_q4.fill.solid()
    box_q4.fill.fore_color.rgb = FILL_LIGHT_RED
    box_q4.line.color.rgb = BORDER_RED
    box_q4.line.width = Pt(1.2)
    tf_q4 = box_q4.text_frame
    tf_q4.word_wrap = True
    tf_q4.margin_top = Inches(0.04)
    tf_q4.margin_left = Inches(0.06)
    tf_q4.margin_right = Inches(0.06)
    p_q4 = tf_q4.paragraphs[0]
    p_q4.text = "LONJAKAN DESAIN BARU Q4 (5,11%):"
    p_q4.font.name = 'Arial'
    p_q4.font.size = Pt(8)
    p_q4.font.bold = True
    p_q4.font.color.rgb = RED
    p_q4_b = tf_q4.add_paragraph()
    p_q4_b.text = "Lonjakan (+1,14 pp vs Q2) terjadi saat order baru masuk tanpa sistem diagnostik harian di meja mesin sehingga make-ready lama."
    p_q4_b.font.name = 'Arial'
    p_q4_b.font.size = Pt(7.5)
    p_q4_b.font.color.rgb = DARK_TEXT


# ==============================================================================
# SLIDE 2 OF 3: KESENJANGAN ALIRAN DATA, TITIK BUTA & MATRIKS RISIKO (BADGE 02)
# ==============================================================================
def build_slide_2_compressed(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide, "02",
        "Kesenjangan Aliran Data (Data Silo), Titik Buta Lapangan & Risiko Pembiaran",
        "Ungkap akar masalah data silo meja mesin vs SAP, kendala diagnostik spekulatif > 8 jam, dan evaluasi 5 pilar Cost of Inaction."
    )

    # Top Flow: Data Silo Architecture & Attribution Blindness
    add_card(slide, 0.60, 1.82, 12.13, 1.45, FILL_LIGHT_AMBER, BORDER_AMBER)
    add_section_header(slide, 0.60, 1.82, 12.13, 0.28, "ANATOMI PEMISAHAN ALIRAN DATA (DATA SILO) & FENOMENA KEBUTAAN ATRIBUSI (ATTRIBUTION BLINDNESS)", ORANGE)

    # Sub-block 1: Meja Mesin
    box_s1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(2.16), Inches(3.60), Inches(1.02))
    box_s1.fill.solid()
    box_s1.fill.fore_color.rgb = WHITE
    box_s1.line.color.rgb = ORANGE
    box_s1.line.width = Pt(1.2)
    tf_s1 = box_s1.text_frame
    tf_s1.word_wrap = True
    tf_s1.margin_top = Inches(0.04)
    tf_s1.margin_left = Inches(0.08)
    tf_s1.margin_right = Inches(0.08)
    p = tf_s1.paragraphs[0]
    p.text = "PENCATATAN DI MEJA MESIN:"
    p.font.name = 'Arial'
    p.font.size = Pt(8.5)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p_b = tf_s1.add_paragraph()
    p_b.text = "• Ditulis manual pada BUKU FOLIO FISIK di meja kontrol 9 mesin cetak.\n• Data jam kerja & operator terisolasi, baru direkap berbulan-bulan kemudian."
    p_b.font.name = 'Arial'
    p_b.font.size = Pt(7.5)
    p_b.font.color.rgb = DARK_TEXT

    # Arrow 1
    arr1 = slide.shapes.add_textbox(Inches(4.35), Inches(2.35), Inches(0.45), Inches(0.45))
    p_a1 = arr1.text_frame.paragraphs[0]
    p_a1.text = "≠"
    p_a1.font.name = 'Arial'
    p_a1.font.size = Pt(22)
    p_a1.font.bold = True
    p_a1.font.color.rgb = RED
    p_a1.alignment = PP_ALIGN.CENTER

    # Sub-block 2: SAP & Verifikasi
    box_s2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.82), Inches(2.16), Inches(3.60), Inches(1.02))
    box_s2.fill.solid()
    box_s2.fill.fore_color.rgb = WHITE
    box_s2.line.color.rgb = PURPLE
    box_s2.line.width = Pt(1.2)
    tf_s2 = box_s2.text_frame
    tf_s2.word_wrap = True
    tf_s2.margin_top = Inches(0.04)
    tf_s2.margin_left = Inches(0.08)
    tf_s2.margin_right = Inches(0.08)
    p = tf_s2.paragraphs[0]
    p.text = "HASIL SORTIR DI VERIFIKASI & SAP:"
    p.font.name = 'Arial'
    p.font.size = Pt(8.5)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p_b = tf_s2.add_paragraph()
    p_b.text = "• Lembar disortir di Unit Verifikasi (waktu tunggu 1–2 hari pasca-cetak).\n• Diinput ke SAP ZPPRSIPPC0012 sebagai RINGKASAN GLOBAL tanpa nomor mesin & shift."
    p_b.font.name = 'Arial'
    p_b.font.size = Pt(7.5)
    p_b.font.color.rgb = DARK_TEXT

    # Arrow 2
    arr2 = slide.shapes.add_textbox(Inches(8.45), Inches(2.35), Inches(0.45), Inches(0.45))
    p_a2 = arr2.text_frame.paragraphs[0]
    p_a2.text = "➔"
    p_a2.font.name = 'Arial'
    p_a2.font.size = Pt(20)
    p_a2.font.bold = True
    p_a2.font.color.rgb = RED
    p_a2.alignment = PP_ALIGN.CENTER

    # Sub-block 3: Blind Spot Consequence
    box_s3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.92), Inches(2.16), Inches(3.68), Inches(1.02))
    box_s3.fill.solid()
    box_s3.fill.fore_color.rgb = FILL_LIGHT_RED
    box_s3.line.color.rgb = RED
    box_s3.line.width = Pt(1.2)
    tf_s3 = box_s3.text_frame
    tf_s3.word_wrap = True
    tf_s3.margin_top = Inches(0.04)
    tf_s3.margin_left = Inches(0.08)
    tf_s3.margin_right = Inches(0.08)
    p = tf_s3.paragraphs[0]
    p.text = "KEBUTAAN ATRIBUSI (ATTRIBUTION BLINDNESS):"
    p.font.name = 'Arial'
    p.font.size = Pt(8.5)
    p.font.bold = True
    p.font.color.rgb = RED
    p_b = tf_s3.add_paragraph()
    p_b.text = "Mengetahui *apa jenis cacatnya* (blobor, bintik, misregister) di tingkat unit, tetapi *BUTA di mesin mana, pada shift berapa, dan faktor mekanis vs operator apa* yang memicunya."
    p_b.font.name = 'Arial'
    p_b.font.size = Pt(7.5)
    p_b.font.color.rgb = DARK_TEXT

    # Bottom-Left Panel: Implikasi Titik Buta Lapangan & 3 Kendala Nyata
    col_w_bl = 5.20
    add_card(slide, 0.60, 3.35, col_w_bl, 3.60, FILL_LIGHT_RED, BORDER_RED)
    add_section_header(slide, 0.60, 3.35, col_w_bl, 0.30, "IMPLIKASI TITIK BUTA DATA TERHADAP OPERASIONAL LAPANGAN", RED)

    tb_bl = slide.shapes.add_textbox(Inches(0.72), Inches(3.70), Inches(col_w_bl - 0.24), Inches(3.15))
    tf_bl = tb_bl.text_frame
    tf_bl.word_wrap = True
    tf_bl.margin_top = Inches(0)
    tf_bl.margin_left = Inches(0)

    blind_points = [
        ("1. Pemeriksaan Spekulatif (> 8 Jam / Mesin):",
         " Saat verifikasi melaporkan kenaikan cacat blobor atau bintik, teknisi tidak tahu mesin mana yang bermasalah dan harus memeriksa 9 mesin bergilir secara *trial-and-error*, sehingga jam henti produktif (*downtime*) melampaui > 1 shift."),
        ("2. Dilema Komponen Mesin vs Cara Kerja Operator:",
         " Sulit memisahkan antara penurunan performa komponen fisik mesin (rol karet mengeras/licin, blanket kempes/turun elastisitas, penjepit silinder kendur) atau variasi *make-ready* dan kelelahan sirkadian pada *Shift* Malam (23.00–07.00 WIB)."),
        ("3. Evaluasi Kerja Operator Tertunda:",
         " Kepala Unit & Kepala Kelompok tidak dapat memberikan bimbingan teknis (*coaching*) harian karena rekam jejak kerja baru direkap berbulan-bulan kemudian.")
    ]
    for idx, (head, body) in enumerate(blind_points):
        p = tf_bl.paragraphs[0] if idx == 0 else tf_bl.add_paragraph()
        p.space_after = Pt(4)
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = 'Arial'
        r0.font.size = Pt(8.5)
        r0.font.bold = True
        r0.font.color.rgb = RED
        r1 = p.add_run()
        r1.text = head
        r1.font.name = 'Arial'
        r1.font.size = Pt(8.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        r2 = p.add_run()
        r2.text = body
        r2.font.name = 'Arial'
        r2.font.size = Pt(8)
        r2.font.color.rgb = DARK_TEXT

    # Bottom-Right Panel: Matriks Risiko Pembiaran (5 Pilar Cost of Inaction & Valuasi Finansial)
    col_w_br = 6.78
    add_card(slide, 5.95, 3.35, col_w_br, 3.60, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 5.95, 3.35, col_w_br, 0.30, "MATRIKS EVALUASI 5 PILAR RISIKO PEMBIARAN (COST OF INACTION)", PURPLE)

    # 5 Structured Mini-Cards Grid
    inaction_rows = [
        ("1. BIAYA (COST)", "KRITIS", RED, FILL_LIGHT_RED, BORDER_RED,
         "Akumulasi pemborosan biaya bahan baku kertas sekuriti & tinta khusus mencapai Rp 22,13 M s.d. Rp 24,56 Miliar/tahun (8,18M LK × Rp 3.000*). Valuasi tiap 1% penurunan inschiet setara Rp 4,80 M – Rp 5,33 Miliar/tahun."),
        ("2. MUTU (QUALITY)", "TINGGI", RED, FILL_LIGHT_RED, BORDER_RED,
         "Tingkat inschiet berfluktuasi hingga 5,11% akibat perbaikan mesin reaktif dan tidak menyentuh akar masalah per-mesin."),
        ("3. KEPATUHAN", "TINGGI", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE,
         "Pencatatan manual di buku folio menyulitkan penelusuran riwayat lot produksi saat audit mutu ISO 9001:2015."),
        ("4. K3L / ESG", "SEDANG", ORANGE, FILL_LIGHT_AMBER, BORDER_AMBER,
         "Timbulan limbah afval 7,37 – 8,18 Juta LK/tahun (±60–65 Ton kertas) & beban fisik kelelahan operator shift malam."),
        ("5. LAYANAN (SLA)", "TINGGI", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY,
         "Siklus tambah cetak lambat menunda serah terima dokumen ke DJBC dan memicu ancaman penalti keterlambatan kontrak.")
    ]

    for idx, (p_title, p_sev, p_col, p_fill, p_bord, p_desc) in enumerate(inaction_rows):
        cy = 3.70 + idx * 0.62
        sub_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.05), Inches(cy), Inches(col_w_br - 0.20), Inches(0.58))
        sub_card.fill.solid()
        sub_card.fill.fore_color.rgb = p_fill
        sub_card.line.color.rgb = p_bord
        sub_card.line.width = Pt(1.0)

        tb_c = slide.shapes.add_textbox(Inches(6.12), Inches(cy + 0.02), Inches(col_w_br - 0.34), Inches(0.54))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        tf_c.margin_top = Inches(0)
        tf_c.margin_left = Inches(0)

        p1 = tf_c.paragraphs[0]
        r0 = p1.add_run()
        r0.text = f"{p_title} [{p_sev}]: "
        r0.font.name = 'Arial'
        r0.font.size = Pt(8)
        r0.font.bold = True
        r0.font.color.rgb = p_col
        r1 = p1.add_run()
        r1.text = p_desc
        r1.font.name = 'Arial'
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = DARK_TEXT


# ==============================================================================
# SLIDE 3 OF 3: SOLUSI DSS SIRINE 4.0, VALIDASI S1 2026 & DAMPAK (BADGE 03)
# ==============================================================================
def build_slide_3_compressed(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    add_header(
        slide, "03",
        "Solusi Terintegrasi DSS SIRINE 4.0, Validasi S1 2026 & Dampak Nyata",
        "Tunjukkan integrasi data meja mesin, modul SAP, dan verifikasi mutu yang mereduksi inschiet ke 3,33% dan menghemat Rp 2,23 Miliar."
    )

    # 4 Top Hero KPIs
    add_kpi_card(slide, 0.60, 1.82, 2.90, 1.10, "INSCHIET AKHIR Q2 2026", "3,33%", "Turun -1,28 pp (-27,8%) vs Baseline 4,61%", GREEN, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_kpi_card(slide, 3.65, 1.82, 2.90, 1.10, "RATA-RATA SEMESTER 1", "3,89%", "Volume: 103.345.688 LK (< Target 4,00%)", PURPLE, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_kpi_card(slide, 6.70, 1.82, 2.90, 1.10, "LEMBAR DISELAMATKAN", "743.234 LK", "Bahan Baku Sekuriti Fisik Bernilai Tinggi", NAVY, FILL_LIGHT_NAVY, BORDER_NAVY)
    add_kpi_card(slide, 9.75, 1.82, 2.98, 1.10, "EFISIENSI FINANSIAL S1", "Rp 2,23 Miliar", "Proyeksi: Rp 6,82 M/Thn | Biaya Dev: Rp 0,-", RED, FILL_LIGHT_AMBER, BORDER_AMBER)

    # Middle Architecture Flow Banner (3 Integrated Points + Granular Chain)
    mid_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.60), Inches(2.98), Inches(12.13), Inches(0.95))
    mid_card.fill.solid()
    mid_card.fill.fore_color.rgb = FILL_LIGHT_PURPLE
    mid_card.line.color.rgb = BORDER_PURPLE
    mid_card.line.width = Pt(1.5)

    tb_m = slide.shapes.add_textbox(Inches(0.72), Inches(3.02), Inches(11.89), Inches(0.88))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    tf_m.margin_top = Inches(0)
    tf_m.margin_left = Inches(0)

    p_m1 = tf_m.paragraphs[0]
    r_m0 = p_m1.add_run()
    r_m0.text = "ARSITEKTUR INTEGRASI 3 TITIK ALIRAN DATA:  "
    r_m0.font.name = 'Arial'
    r_m0.font.size = Pt(8.5)
    r_m0.font.bold = True
    r_m0.font.color.rgb = PURPLE
    r_m1 = p_m1.add_run()
    r_m1.text = "Data Transaksi Meja Mesin (< 30 dtk)  ⇄  Modul SAP Production Order (ZPPRSIPPC0012)  ⇄  Hasil Sortir Verifikasi Mutu (HCTS)\n"
    r_m1.font.name = 'Arial'
    r_m1.font.size = Pt(8.5)
    r_m1.font.bold = True
    r_m1.font.color.rgb = NAVY

    p_m2 = tf_m.add_paragraph()
    r_m2 = p_m2.add_run()
    r_m2.text = "RUNTUTAN DATA GRANULAR:  "
    r_m2.font.name = 'Arial'
    r_m2.font.size = Pt(8)
    r_m2.font.bold = True
    r_m2.font.color.rgb = DARK_GREEN
    r_m3 = p_m2.add_run()
    r_m3.text = "Nomor PO  ➔  Nomor Mesin (9 Mesin)  ➔  Pola Gilir (Shift 1/2/3)  ➔  Tim Operator  ➔  Kategori Cacat Cetak Spesifik  ➔  Tindakan Preskriptif Cepat"
    r_m3.font.name = 'Arial'
    r_m3.font.size = Pt(8)
    r_m3.font.color.rgb = DARK_TEXT

    # Bottom-Left: Native Chart Trend Inschiet + Reduksi Downtime Callout
    col_w_chart = 4.75
    add_card(slide, 0.60, 4.00, col_w_chart, 2.95, FILL_LIGHT_PURPLE, BORDER_PURPLE)
    add_section_header(slide, 0.60, 4.00, col_w_chart, 0.28, "TREN PENURUNAN INSCHIET BASELINE VS REALISASI S1 2026 (%)", PURPLE)

    chart_data = CategoryChartData()
    chart_data.categories = ['Baseline 2025', 'Q1 2026 (Adaptasi)', 'Q2 2026 (Presisi)']
    chart_data.add_series('Inschiet (%)', (4.61, 4.34, 3.33))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.70), Inches(4.30), Inches(4.55), Inches(1.85),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    plots = chart.plots[0]
    plots.has_data_labels = True
    for series in plots.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = GREEN
        dl = series.data_labels
        dl.font.name = 'Arial'
        dl.font.size = Pt(9.5)
        dl.font.bold = True
        dl.font.color.rgb = DARK_GREEN
        dl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

    # Callout Downtime
    dt_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.70), Inches(6.20), Inches(4.55), Inches(0.68))
    dt_box.fill.solid()
    dt_box.fill.fore_color.rgb = WHITE
    dt_box.line.color.rgb = GREEN
    dt_box.line.width = Pt(1.2)
    tf_dt = dt_box.text_frame
    tf_dt.word_wrap = True
    tf_dt.margin_top = Inches(0.04)
    tf_dt.margin_left = Inches(0.06)
    tf_dt.margin_right = Inches(0.06)
    p_dt = tf_dt.paragraphs[0]
    r_d0 = p_dt.add_run()
    r_d0.text = "REDUKSI WAKTU HENTI (DOWNTIME): "
    r_d0.font.name = 'Arial'
    r_d0.font.size = Pt(8)
    r_d0.font.bold = True
    r_d0.font.color.rgb = GREEN
    r_d1 = p_dt.add_run()
    r_d1.text = "Durasi perbaikan terpangkas dari > 1 shift (> 8 jam) menjadi < 2–4 jam (turun 50%–75%)."
    r_d1.font.name = 'Arial'
    r_d1.font.size = Pt(7.5)
    r_d1.font.color.rgb = DARK_TEXT

    # Bottom-Right: Kertas Kerja Realisasi Finansial (Tabel 1.3) + Multi-Impact
    col_w_tbl = 7.23
    add_card(slide, 5.50, 4.00, col_w_tbl, 2.95, FILL_LIGHT_GREEN, BORDER_GREEN)
    add_section_header(slide, 5.50, 4.00, col_w_tbl, 0.28, "KERTAS KERJA REALISASI PENGHEMATAN BIAYA & TRANSFORMASI STRATEGIS", GREEN)

    table_shape = slide.shapes.add_table(4, 6, Inches(5.60), Inches(4.34), Inches(7.03), Inches(1.50))
    table = table_shape.table
    table.columns[0].width = Inches(1.45)
    table.columns[1].width = Inches(1.15)
    table.columns[2].width = Inches(0.90)
    table.columns[3].width = Inches(1.05)
    table.columns[4].width = Inches(1.20)
    table.columns[5].width = Inches(1.28)

    headers = ["Periode", "Volume (n)", "Inschiet", "Deviasi", "Diselamatkan", "Nilai Saving*"]
    rows = [
        ("Q1 2026 (Adaptasi)", "57.385.254", "4,34%", "-0,27 pp", "154.940 LK", "Rp 464,82 Juta"),
        ("Q2 2026 (Presisi)", "45.960.434", "3,33%", "-1,28 pp", "588.294 LK", "Rp 1,76 Miliar"),
        ("TOTAL S1 2026", "103.345.688", "3,89%", "-0,72 pp", "743.234 LK", "Rp 2,23 Miliar")
    ]

    for col_idx, h_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.name = 'Arial'
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for row_idx, r_data in enumerate(rows):
        for col_idx, val in enumerate(r_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = FILL_LIGHT_GREEN if row_idx == 2 else (WHITE if row_idx % 2 == 0 else FILL_LIGHT_NAVY)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = 'Arial'
            p.font.size = Pt(8)
            p.font.color.rgb = DARK_GREEN if row_idx == 2 else DARK_TEXT
            p.font.bold = (row_idx == 2 or col_idx in [0, 4, 5])
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

    # Transformational Value Strip Bottom
    val_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.60), Inches(5.92), Inches(7.03), Inches(0.95))
    val_box.fill.solid()
    val_box.fill.fore_color.rgb = FILL_LIGHT_NAVY
    val_box.line.color.rgb = BORDER_NAVY
    val_box.line.width = Pt(1.2)
    tf_vb = val_box.text_frame
    tf_vb.word_wrap = True
    tf_vb.margin_top = Inches(0.04)
    tf_vb.margin_left = Inches(0.08)
    tf_vb.margin_right = Inches(0.08)

    p_v1 = tf_vb.paragraphs[0]
    r_vb0 = p_v1.add_run()
    r_vb0.text = "KESIMPULAN NILAI TAMBAH KORPORASI:  "
    r_vb0.font.name = 'Arial'
    r_vb0.font.size = Pt(8)
    r_vb0.font.bold = True
    r_vb0.font.color.rgb = NAVY
    r_vb1 = p_v1.add_run()
    r_vb1.text = "DSS SIRINE 4.0 mengamankan efisiensi nyata Rp 2,23 Miliar (proyeksi Rp 6,82 M/tahun), mengeliminasi buku folio fisik (<30 dtk), mewujudkan pembinaan objektif harian bagi ±42 operator, serta memperkuat kepatuhan 100% SLA tender DJBC Kemenkeu RI."
    r_vb1.font.name = 'Arial'
    r_vb1.font.size = Pt(7.5)
    r_vb1.font.color.rgb = DARK_TEXT


# ==============================================================================
# MAIN GENERATOR
# ==============================================================================
def main():
    prs, blank_layout = create_presentation()

    print("Building Slide 1 of 3 (Badge 01: Konteks Strategis, Skala 9 Mesin & Baseline 2025)...")
    build_slide_1_compressed(prs, blank_layout)

    print("Building Slide 2 of 3 (Badge 02: Kesenjangan Data Silo, Titik Buta & Matriks Risiko)...")
    build_slide_2_compressed(prs, blank_layout)

    print("Building Slide 3 of 3 (Badge 03: Solusi DSS SIRINE 4.0, Validasi S1 2026 & Dampak)...")
    build_slide_3_compressed(prs, blank_layout)

    output_path = "Presentasi_Risalah_Latar_Belakang_IAKA_2026_3_Slides.pptx"
    prs.save(output_path)
    print(f"\n[SUCCESS] 3-Slide Master Deck successfully generated: {output_path}")

    # Also overwrite the primary presentation file so user sees the 3-slide version immediately
    primary_output = "Presentasi_Risalah_Latar_Belakang_IAKA_2026.pptx"
    prs.save(primary_output)
    print(f"[SUCCESS] Updated primary file: {primary_output}")

if __name__ == "__main__":
    main()
